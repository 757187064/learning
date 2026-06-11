#!/usr/bin/env python3
import hashlib
import html
import json
import re
import shutil
import subprocess
import zipfile
from collections import Counter, defaultdict
from pathlib import Path
from textwrap import dedent

from curated_quiz_bank import build_curated_quiz_banks, validate_quiz_bank


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "深度学习复习网站"
DATABASE = OUT / "database"
SITE = OUT / "site"
ASSETS = SITE / "assets"
IMAGE_OUT = SITE / "image"
MINDMAP_OUT = SITE / "mindmaps"
GUIDE_OUT = SITE / "guides"
TOPIC_GUIDE_OUT = SITE / "topic_guides"
TOC_IMAGE_OUT = SITE / "guide_toc_images"
CACHE_VERSION = 8
EXCLUDED_SOURCE_MD = {
    "深度学习期末复习_仅基于原始课件版.md",
    "深度学习期末复习_仅基于原始课件版.before_restore.md",
}
FINAL_REVIEW_MD = ROOT / "深度学习期末复习_仅基于原始课件版.md"
FINAL_REVIEW_SITE_MD = TOPIC_GUIDE_OUT / "深度学习期末复习-总复习资料.md"


KEY_TERMS = [
    "TLU", "阈值逻辑单元", "逻辑门", "AND", "OR", "NOT", "XOR", "Constructive Synthesis",
    "感知机", "ADALINE", "学习规则", "梯度更新",
    "梯度下降", "损失函数", "学习率", "Epoch", "Batch", "反向传播", "计算图", "自动微分",
    "Logistic", "Sigmoid", "Softmax", "交叉熵", "BCEWithLogitsLoss", "CrossEntropyLoss",
    "MLP", "多层感知机", "激活函数", "仿射变换", "卷积", "互相关", "卷积核", "步长", "填充",
    "Valid", "Same", "Full", "感受野", "池化", "NCHW", "NHWC", "归一化", "标准化",
    "LeNet", "AlexNet", "ImageFolder", "WeightedRandomSampler", "类别不平衡",
    "梯度消失", "梯度爆炸", "梯度裁剪", "Xavier", "Kaiming", "BatchNorm", "Dropout",
    "数据增强", "权重衰减", "早停", "优化器", "Momentum", "Nesterov", "Adagrad", "RMSProp",
    "Adam", "学习率调度器", "ReduceLROnPlateau", "超参数", "消融实验",
    "VGG", "GoogLeNet", "Inception", "ResNet", "残差连接", "迁移学习",
    "RNN", "隐藏状态", "双向RNN", "堆叠RNN", "GRU", "LSTM", "Padding", "Packing",
    "PackedSequence", "1D卷积", "TCN", "Seq2Seq", "Encoder-Decoder", "Teacher Forcing",
    "注意力机制", "Attention", "Query", "Key", "Value", "Q/K/V", "上下文向量",
    "Scaled Dot-Product Attention", "sqrt(d_k)", "Source Mask", "Multi-Head Attention",
    "Self-Attention", "Cross-Attention", "Target Mask", "Subsequent Mask",
    "位置编码", "Positional Encoding", "Transformer", "RoPE",
    "图学习", "图神经网络", "节点", "边", "邻接矩阵", "度矩阵", "拉普拉斯矩阵",
    "谱图理论", "图卷积", "GCN", "GAT", "GraphSAGE", "消息传递", "聚合",
    "同质图", "异质图", "节点分类", "图分类", "链路预测", "随机游走",
    "PageRank", "嵌入", "网络表示学习", "Vertex", "Edge",
]

BAD_QUESTION_PATTERNS = [
    "核心词是",
    "这句话描述的核心词",
    "必须抓住的两个要点",
    "围绕“",
    "围绕\"",
    "若考试围绕",
    "若题目围绕",
    "若训练代码围绕",
    "若围绕",
    "关系最直接的知识点",
    "最适合作为考试标准答案",
]

NOISE_PATTERNS = [
    "import ", "warnings.", "plt.", "torch.", "np.", "rcParams", "dark_background",
    "本章内容", "通过本练习，你将", "练习目标", "|:----", "答案解析", "回答正确",
    "分值1分", "class ", "def ", "return ", "from ",
    "北京航空航天", "北京航", "UNIVERSITY", "233740", "2026年05月26日",
]

TERM_GUIDE = {
    "卷积": ("用小窗口在图像局部滑动，对局部像素和卷积核权重相乘求和，提取边缘、纹理等局部特征。", "常考输出尺寸、卷积核/步长/填充的作用，以及 CNN 实际常用互相关。", "不要把 CNN 说成全连接；卷积核心是局部连接和参数共享。", "输出尺寸：(输入尺寸 + 2P - K) / S + 1。"),
    "互相关": ("形式接近卷积，但不翻转卷积核；深度学习框架里的卷积层通常实际做互相关。", "常考 CNN 中是否需要翻转卷积核，或为什么工程上用互相关。", "不要因为名字叫卷积层，就默认一定执行数学卷积的核翻转。", "CNN 训练会学习卷积核参数，因此互相关也能学到有效特征。"),
    "Valid": ("不填充边界，只在卷积核完全落入输入时计算，输出尺寸会变小。", "常和 Same、Full 对比考输出大小。", "Valid 不是大量填充，恰好是不填充。", "Valid：无填充，输出通常最小。"),
    "Same": ("通过适当填充，让卷积输出空间尺寸尽量和输入保持一致。", "常考给定输入、卷积核、步长时需要多少 padding。", "Same 的目标是保持尺寸，不是扩大特征图。", "步长为 1 时，Same 卷积常让 H、W 不变。"),
    "Full": ("在边界外进行更充分的填充，使卷积核可以覆盖到输入边缘外侧，输出会变大。", "常和 Valid、Same 放在一起判断哪种输出最大。", "Full 不是常规 CNN 默认选择，考试多用于概念区分。", "Full：充分填充，输出通常最大。"),
    "感受野": ("某一层一个神经元能看到原始输入图像的区域大小。", "常考层数、卷积核大小、步长如何影响感受野。", "步长变大通常会让前一层对应感受野扩大，不是缩小。", "层越深，感受野通常越大；多层小卷积能逐步扩大感受野。"),
    "池化": ("对局部区域做最大值或平均值汇总，降低空间尺寸并增强一定平移不变性。", "常考池化作用、是否有可学习参数、与卷积的区别。", "池化不是全连接，也不是用来增加参数量。", "池化：降采样、压缩空间尺寸、保留显著信息。"),
    "卷积核": ("卷积层中可学习的小矩阵/滤波器，用来扫描局部区域并提取特征。", "常考卷积核大小、通道数、参数共享和输出尺寸计算。", "卷积核参数会训练更新，不是手工固定模板。", "卷积核大小 K、步长 S、填充 P 共同决定输出尺寸。"),
    "步长": ("卷积核或池化窗口每次移动的间隔。", "常考步长变大时输出尺寸如何变化。", "步长越大，输出空间尺寸通常越小。", "Stride 控制滑动间隔。"),
    "填充": ("在输入边界补值，让卷积能处理边缘或保持输出尺寸。", "常考 Same 卷积为何需要 padding。", "填充不是增加有效信息，而是控制边界和尺寸。", "Padding 用于控制输出尺寸和边界信息。"),
    "标准化": ("把数据按均值和方差调整到更稳定的尺度，帮助模型训练。", "常考归一化和标准化的区别、为什么图像输入要预处理。", "不要把标准化简单等同于除以 255；标准化通常涉及均值和标准差。", "标准化常写作 (x-mean)/std。"),
    "归一化": ("把数值缩放到较统一范围，例如图像像素从 0-255 缩放到 0-1。", "常考图像预处理为什么要缩放像素值。", "归一化不等于 BatchNorm，前者多是输入预处理。", "图像常先把像素值缩放到 0-1。"),
    "激活函数": ("给线性变换加入非线性，使网络能拟合非线性关系。", "常考没有激活函数时多层线性网络仍等价于线性模型。", "不要把激活函数说成只改变维度，它主要提供非线性。", "仿射变换 + 激活函数 = 神经网络层的基本模式。"),
    "BatchNorm": ("在小批量上标准化中间激活，再用可学习参数恢复表达能力，帮助训练更稳定。", "常考训练/推理阶段统计量不同，以及放在卷积/全连接和激活附近。", "推理阶段一般使用训练中累计的均值方差，不依赖当前 batch。", "BatchNorm：标准化激活，稳定训练，加快收敛。"),
    "Dropout": ("训练时随机丢弃一部分神经元，迫使网络不要过度依赖某些特征。", "常考训练阶段启用、推理阶段关闭，以及它缓解过拟合。", "Dropout 不是提高模型容量，而是正则化。", "Dropout：训练随机失活，推理正常使用。"),
    "Batch": ("一次送入模型并一起计算损失的一组样本。它决定每次参数更新看到多少样本，也影响显存占用和训练稳定性。", "常考 batch size、mini-batch、epoch、iteration 的区别，也会结合张量形状中的 N 维度判断。", "Batch 不是序列长度；在 RNN 中常见形状里 N 是样本数，L 才是时间步长度。", "Batch 是一批样本；Epoch 是全训练集完整训练一遍。"),
    "Epoch": ("训练集被模型完整看过一遍，叫一个 epoch。一个 epoch 内通常包含很多个 batch 更新。", "常考 Epoch、Batch、Iteration 的区别，或判断训练轮数增加对欠拟合/过拟合的影响。", "Epoch 不是一次参数更新；一次更新通常对应一个 batch。", "1 个 Epoch = 全部训练样本被用过一遍。"),
    "学习率": ("控制每次参数沿梯度方向更新的步子大小。学习率太大容易震荡甚至发散，太小则收敛很慢。", "常考学习率过大/过小的训练现象，以及学习率调度的目的。", "学习率不是越大越好；也不是模型学到的参数。", "学习率决定更新步长。"),
    "损失函数": ("把模型预测和真实标签之间的差距变成一个可优化的数值，反向传播就是从损失开始计算梯度。", "常考输出层与损失函数的匹配，例如多分类常配 CrossEntropyLoss，二分类可配 BCEWithLogitsLoss。", "不要把准确率当作训练损失；也不要把 BCE 和 CE 的适用场景混用。", "损失函数告诉模型“错在哪里”，优化器负责“怎么改参数”。"),
    "Sigmoid": ("把任意实数压到 0 到 1 之间，可解释为概率或门控强度。GRU/LSTM 的门通常用它决定保留多少信息。", "常考输出范围、二分类概率、以及门控结构中为什么适合控制开关。", "Sigmoid 不只属于分类输出，在 RNN 门控里也很重要。", "Sigmoid 输出在 0 到 1；越接近 1 越保留，越接近 0 越抑制。"),
    "Softmax": ("把一组分数转成和为 1 的概率分布，常用于多分类或注意力权重。", "常考多分类输出、注意力权重为什么能加权求和。", "Softmax 是对一组数整体归一化，不是逐个独立压缩。", "Softmax 输出非负且总和为 1。"),
    "梯度消失": ("反向传播经过很多层或很多时间步时，梯度可能不断变小，导致前面层或早期时间步很难学到东西。", "常考深层网络/RNN 为什么训练困难，以及 LSTM、GRU、残差连接如何缓解。", "梯度消失不是过拟合；它是参数难以有效更新的问题。", "多个小于 1 的因子连乘，会让梯度趋近 0。"),
    "梯度爆炸": ("梯度在反向传播中变得过大，使参数更新剧烈，训练损失可能震荡或变成 NaN。", "常考梯度裁剪为什么能缓解 RNN 训练不稳定。", "Dropout 主要缓解过拟合，不是解决梯度爆炸的标准答案。", "梯度裁剪限制梯度范数，防止更新步子过大。"),
    "梯度裁剪": ("当梯度范数超过阈值时按比例缩小，常用于稳定 RNN 等容易梯度爆炸的训练。", "常考它处理的是梯度爆炸，而不是梯度消失。", "裁剪不是把梯度直接清零，而是限制其大小。", "Gradient Clipping：限制梯度大小，稳定训练。"),
    "Adam": ("结合动量和自适应学习率思想，为不同参数调整更新幅度。", "常和 SGD、Momentum、RMSProp 对比优化器特点。", "Adam 方便但不等于一定泛化最好，仍需调学习率等超参数。", "Adam = 一阶动量 + 二阶矩估计。"),
    "RMSProp": ("用梯度平方的指数滑动平均调整学习率，缓解 Adagrad 学习率过快衰减。", "常考它和 Adagrad 的区别。", "不要写成简单固定学习率下降。", "RMSProp 使用指数滑动平均保存近期梯度尺度。"),
    "VGG": ("大量使用 3x3 小卷积核堆叠构建深层网络。", "常考小卷积核堆叠为什么能扩大感受野并减少参数。", "不要把 VGG 的特点写成复杂多分支结构。", "VGG：多层小卷积核堆叠。"),
    "ResNet": ("用残差/跳跃连接让网络学习 F(x)+x，缓解深层网络退化和梯度传播困难。", "常考残差连接为什么能训练更深网络。", "残差连接不是单纯增加层数，而是改变信息和梯度路径。", "ResNet 核心：y = F(x) + x。"),
    "RNN": ("按时间步处理序列，用隐藏状态把前面信息传到后面。", "常考隐藏状态、序列建模、梯度消失/爆炸。", "RNN 不是一次性把所有时间步完全独立处理。", "当前输出依赖当前输入和上一时刻隐藏状态。"),
    "GRU": ("用更新门和重置门控制历史信息保留与遗忘，是简化版门控循环网络。", "常和 LSTM 对比门结构和参数量。", "GRU 没有单独的细胞状态 c_t。", "GRU：更新门 + 重置门。"),
    "LSTM": ("通过输入门、遗忘门、输出门和细胞状态保存长期信息。", "常考门控作用，以及为什么能缓解长序列梯度问题。", "不要漏掉细胞状态是 LSTM 的关键通道。", "LSTM：输入门、遗忘门、输出门、细胞状态。"),
    "Padding": ("把不同长度的序列补到同一长度，方便组成 batch 输入模型。", "常考为什么需要 padding，以及它和 packing 的关系。", "Padding 补出来的位置不是真实信息，训练时要避免把它当有效时间步。", "Padding：补齐长度；Packing：跳过无效补位。"),
    "Packing": ("把变长序列的有效部分打包，让 RNN 少处理 padding 位置。", "常考 pack_padded_sequence 的目的和使用前提。", "Packing 不是改变序列含义，而是提高变长序列处理效率。", "Packing 让 RNN 聚焦真实时间步。"),
    "PackedSequence": ("PyTorch 中表示变长序列打包后的对象，内部只保留有效时间步，交给 RNN 处理时可以跳过 padding。", "常考 pack_padded_sequence 与 pad_packed_sequence 的前后关系。", "PackedSequence 不是普通张量；使用前后要注意 lengths、batch_first 和是否排序。", "先 padding 组成 batch，再 packing 跳过补位，必要时再 pad 回来。"),
    "1D卷积": ("在序列的一维时间/位置方向上滑动卷积核，提取局部 n-gram 或短期模式。", "常和 RNN 对比：1D 卷积并行提取局部模式，RNN 逐步传递隐藏状态。", "1D 卷积适合局部依赖，不等于天然拥有长期记忆。", "1D 卷积：沿序列维度提取局部模式。"),
    "隐藏状态": ("RNN 在每个时间步保存的当前记忆，用来把历史信息传给后续时间步。", "常考 h_t 与 h_{t-1} 的关系，以及 output 和 h_n 的区别。", "隐藏状态不是模型参数，而是随输入序列变化的中间表示。", "h_t = 当前输入 + 过去记忆的综合表示。"),
    "双向RNN": ("同时从前往后和从后往前读序列，让当前位置利用左右两侧上下文。", "常考它为什么比单向 RNN 能看到更多上下文。", "双向 RNN 不适合严格只能用过去信息的实时预测场景。", "双向 RNN = 正向隐藏状态 + 反向隐藏状态。"),
    "堆叠RNN": ("把多层 RNN 叠起来，上一层的输出序列作为下一层输入。", "常考堆叠带来更强表达能力，也增加训练难度。", "堆叠不是时间步变多，而是层数变深。", "堆叠 RNN：时间方向 + 层方向都要传播。"),
    "注意力机制": ("根据 Query 与 Key 的匹配程度，给 Value 分配权重并加权求和。", "常考 Q/K/V 含义、注意力权重、上下文向量计算流程。", "注意力不是简单平均，而是按相关性加权。", "Attention(Q,K,V)=softmax(QK^T/sqrt(d_k))V。"),
    "Self-Attention": ("Q、K、V 来自同一个序列，让序列内部不同位置相互建立联系。", "常考它和 Cross-Attention 的区别。", "Self 不是只看自己一个位置，而是同一序列内部互相看。", "自注意力：同一序列内部做注意力。"),
    "Cross-Attention": ("Query 来自一个序列，Key/Value 来自另一个序列，常用于解码器关注编码器输出。", "常考 Transformer 解码器中 Cross-Attention 的输入来源。", "不要和 Self-Attention 混成同一来源。", "交叉注意力：Q 与 K/V 来源不同。"),
    "Transformer": ("以注意力为核心的序列模型，用多头注意力、前馈网络、残差和归一化堆叠编码器/解码器。", "常考编码器/解码器结构、多头注意力和位置编码。", "Transformer 不依赖 RNN 的逐步递归来建模序列。", "Transformer 核心：Attention + FFN + 残差 + Norm + 位置编码。"),
    "sqrt(d_k)": ("缩放点积注意力中的除数，用来避免 QK 点积过大导致 softmax 过于尖锐。", "常考为什么要缩放，以及 d_k 是 Key/Query 的维度。", "不要把它当成可学习参数。", "QK^T 除以 sqrt(d_k) 是为了稳定 softmax。"),
    "RoPE": ("一种旋转位置编码，把位置信息注入向量表示，常用于现代 Transformer。", "如果课件讲到，主要考它属于位置编码思想，不会要求复杂推导。", "不要把 RoPE 写成普通绝对位置编号相加。", "RoPE 用旋转方式表达相对位置信息。"),
    "Seq2Seq": ("编码器把输入序列压成表示，解码器再生成输出序列，常用于翻译、问答等序列到序列任务。", "常考 Encoder-Decoder 结构、Teacher Forcing 和注意力机制为何被引入。", "普通 Seq2Seq 容易受固定长度上下文瓶颈影响。", "Seq2Seq = Encoder + Decoder。"),
    "图学习": ("研究如何在图结构数据上学习表示和完成预测，核心是同时利用节点特征与关系结构。", "常考图数据适用场景、节点/边/图级任务。", "不要只看单个样本特征，图学习关键是关系。", "图学习 = 特征信息 + 结构信息。"),
    "节点": ("图中的对象或实体，可以带有节点特征。", "常考节点分类任务：预测每个节点的类别。", "节点特征和边关系是两类信息，不要混在一起。", "节点表示对象，边表示对象之间的关系。"),
    "边": ("连接两个节点的关系，可以是有向/无向、带权/不带权。", "常考边方向、权重和邻接矩阵的对应关系。", "有向图的邻接矩阵通常不对称。", "边表示关系，权重表示关系强弱。"),
    "Vertex": ("图中的节点，也就是图结构里的对象或实体。", "常考 Vertex/Node 与 Edge 的区别。", "Vertex 是节点，不是边。", "Vertex/Node 表示对象；Edge 表示关系。"),
    "Edge": ("图中的边，表示两个节点之间的关系，可以有方向和权重。", "常考 Edge 与邻接矩阵元素的对应关系。", "Edge 不是节点特征本身，而是节点之间的连接。", "Edge 表示关系，A_ij 常记录边或边权。"),
    "邻接矩阵": ("用矩阵记录节点之间是否相连或连接权重。", "常考根据图画矩阵，或根据矩阵判断邻居。", "无向图邻接矩阵通常对称；有向图不一定对称。", "A_ij 表示 i 到 j 是否有边或边权。"),
    "度矩阵": ("对角矩阵，对角线记录每个节点的度。", "常和拉普拉斯矩阵 L=D-A 一起考。", "度矩阵不是邻接矩阵，它只放在对角线上。", "D_ii = 节点 i 的度。"),
    "拉普拉斯矩阵": ("用度矩阵和邻接矩阵构造，反映图的连接结构。", "常考 L=D-A 以及它在谱图方法中的作用。", "不要把 L 和 A 混用。", "图拉普拉斯：L = D - A。"),
    "GCN": ("图卷积网络，通过邻接关系聚合邻居节点特征，再更新节点表示。", "常考 GCN 为什么能利用图结构、聚合邻居信息的基本流程。", "GCN 不只看节点自身特征，还看邻居。", "GCN：邻居聚合 + 线性变换 + 非线性。"),
    "GAT": ("图注意力网络，给不同邻居分配不同注意力权重再聚合。", "常和 GCN 对比：GCN 多按固定归一化聚合，GAT 学习邻居重要性。", "GAT 的重点是邻居权重可学习，不是所有邻居同等重要。", "GAT：用注意力学习邻居权重。"),
    "GraphSAGE": ("通过采样邻居并聚合来生成节点表示，适合较大图的归纳学习。", "常考采样聚合思想，以及和一次性全图训练的区别。", "GraphSAGE 不是简单查表保存每个节点向量。", "GraphSAGE：采样邻居，聚合表示。"),
    "消息传递": ("节点从邻居接收信息、聚合信息、更新自身表示，是很多图神经网络的统一视角。", "常考消息、聚合、更新三个步骤。", "不要只写公式，考试要说明信息沿边传播。", "消息传递：邻居发消息，节点聚合并更新。"),
}

TERM_GUIDE.update({
    "MLP": ("多层感知机由多层线性变换和非线性激活组成，是最基础的前馈神经网络。", "常考为什么需要激活函数、前向传播流程、分类输出和损失函数匹配。", "不要把多层线性层误认为天然具有非线性；没有激活函数仍等价于线性变换。", "MLP = 线性层 + 激活函数 + 输出层。"),
    "多层感知机": ("多层感知机是由输入层、隐藏层和输出层组成的前馈神经网络。", "常考隐藏层作用、激活函数必要性、训练流程。", "多层感知机不是循环网络，也不保留时间状态。", "多层感知机靠非线性激活增强表达能力。"),
    "仿射变换": ("线性加权求和再加偏置，常写作 Wx+b。", "常考它和激活函数的配合：仿射变换本身仍是线性结构。", "不要把仿射变换本身说成非线性映射。", "仿射变换：Wx+b。"),
    "梯度下降": ("沿着损失函数下降最快的方向调整参数，让模型预测逐步变好。", "常考参数更新方向、学习率作用、训练循环位置。", "梯度下降不是直接调准确率，而是通过损失函数和梯度改参数。", "参数更新方向通常与梯度相反。"),
    "反向传播": ("从损失函数出发，按链式法则把梯度传回各层参数。", "常考它和前向传播、损失函数、优化器的顺序关系。", "反向传播只负责算梯度，真正改参数的是优化器。", "前向算预测，损失算错误，反向算梯度。"),
    "计算图": ("把张量运算组织成图结构，框架据此追踪梯度计算路径。", "常考 PyTorch 动态计算图、requires_grad、detach 等概念。", "不要在需要梯度的地方把计算图断开。", "计算图记录运算关系，反向传播沿图求梯度。"),
    "自动微分": ("深度学习框架自动根据计算图计算梯度。", "常考 backward、grad、requires_grad 的关系。", "自动微分不是自动训练；还需要 loss.backward 和 optimizer.step。", "自动微分帮你算梯度，不替你设计训练流程。"),
    "Logistic": ("Logistic 回归常用于二分类，本质是线性输出接 Sigmoid 得到概率。", "常考二分类输出、Sigmoid、BCE 类损失的搭配。", "Logistic 回归名字里有回归，但常用于分类。", "Logistic：线性打分 + Sigmoid。"),
    "交叉熵": ("衡量预测概率分布与真实类别之间差距的分类损失。", "常考多分类用 CrossEntropyLoss，二分类/多标签用 BCE 类损失。", "不要把交叉熵和准确率混成同一个东西。", "交叉熵越小，预测分布越接近真实标签。"),
    "BCEWithLogitsLoss": ("二分类或多标签任务常用损失，内部把 Sigmoid 和 BCE 合在一起，数值更稳定。", "常考它输入的是 logits，不需要手动先做 Sigmoid。", "不要在 BCEWithLogitsLoss 前再手动 Sigmoid。", "BCEWithLogitsLoss = Sigmoid + BCE。"),
    "CrossEntropyLoss": ("多分类常用损失，内部包含 LogSoftmax 和 NLLLoss。", "常考它输入 logits，标签通常是类别索引。", "不要在 CrossEntropyLoss 前手动 Softmax。", "CrossEntropyLoss 用于单标签多分类。"),
    "优化器": ("根据梯度和自身规则更新模型参数。", "常考 SGD、Momentum、RMSProp、Adam 的区别。", "优化器不是损失函数；损失给目标，优化器改参数。", "优化器负责执行参数更新。"),
    "Momentum": ("在梯度下降中加入历史更新方向，让更新更平滑、减少震荡。", "常考它为什么能加速或稳定训练。", "Momentum 不是学习率本身，而是利用历史梯度趋势。", "Momentum 像给梯度下降加惯性。"),
    "Nesterov": ("一种提前看一步的动量方法，先按动量方向预估，再计算修正梯度。", "常和普通 Momentum 对比考思想差别。", "不要把 Nesterov 当成完全独立于动量的优化器。", "Nesterov 是带预判的动量。"),
    "Adagrad": ("根据历史梯度平方累积自动调整每个参数学习率。", "常考它学习率会持续衰减，后期可能过小。", "Adagrad 不适合所有场景，持续累积会让更新越来越慢。", "Adagrad 对频繁更新参数降低学习率。"),
    "权重衰减": ("通过惩罚过大的权重来正则化模型，常用于缓解过拟合。", "常考它和 L2 正则化、过拟合控制的关系。", "不要把权重衰减当成学习率衰减。", "权重衰减限制参数过大。"),
    "学习率调度器": ("按训练过程动态调整学习率，例如训练停滞时降低学习率。", "常考调度器何时 step、为什么要降学习率。", "调度器不是优化器本身，而是控制优化器的学习率。", "学习率调度器负责改学习率节奏。"),
    "ReduceLROnPlateau": ("当验证指标长期没有改善时降低学习率的调度策略。", "常考它需要监控验证损失或指标。", "不要像普通 scheduler 一样不传指标直接使用。", "ReduceLROnPlateau 看指标停滞再降学习率。"),
    "早停": ("验证集表现长期不提升时提前停止训练，防止继续过拟合。", "常考它和验证集、过拟合的关系。", "早停不是训练失败，而是一种正则化/模型选择策略。", "早停看验证集，不看训练集单独决定。"),
    "超参数": ("训练前人为设定的配置，如学习率、batch size、层数、dropout 比例。", "常考它和模型参数的区别。", "超参数不是通过反向传播直接学出来的参数。", "超参数先设定，模型参数训练中学习。"),
    "消融实验": ("去掉或替换某个模块，观察性能变化，用来判断该模块是否有效。", "常考它如何证明某设计有贡献。", "消融不是随便删模块，而是控制变量比较。", "消融实验用于验证模块贡献。"),
    "数据增强": ("对训练数据做随机变换，增加样本多样性，缓解过拟合。", "常考图像翻转、裁剪、颜色扰动等增强方法。", "数据增强通常只用于训练集，不应破坏标签含义。", "数据增强让模型见到更多合理变化。"),
    "ImageFolder": ("PyTorch 按文件夹名作为类别读取图像数据集的工具。", "常考目录结构与类别标签的对应关系。", "文件夹结构不对会导致类别映射错误。", "ImageFolder：子文件夹名就是类别名。"),
    "WeightedRandomSampler": ("按样本权重采样，常用于类别不平衡时提高少数类被抽到的概率。", "常考它和 class weight、pos_weight 的区别。", "它改变采样概率，不直接改变损失函数公式。", "WeightedRandomSampler 调采样，不调损失。"),
    "类别不平衡": ("不同类别样本数量差距很大，模型容易偏向多数类。", "常考采样、类别权重、pos_weight 等处理方式。", "准确率高不一定说明少数类学得好。", "类别不平衡要关注少数类召回等指标。"),
    "NCHW": ("图像张量格式，依次表示 Batch、Channel、Height、Width，PyTorch 中很常见。", "常考 NCHW 与 NHWC 的维度顺序，以及卷积层输入通道位置。", "不要把 Channel 和 Height/Width 顺序写反。", "NCHW = Batch, Channel, Height, Width。"),
    "NHWC": ("图像张量格式，依次表示 Batch、Height、Width、Channel。", "常和 NCHW 对比考框架数据格式。", "不要把 NHWC 当成 PyTorch 卷积层默认输入格式。", "NHWC = Batch, Height, Width, Channel。"),
    "LeNet": ("早期 CNN 结构，包含卷积、池化和全连接层，用于手写数字识别。", "常考 CNN 基本结构的历史代表。", "不要把 LeNet 写成现代超深残差网络。", "LeNet 展示了卷积 + 池化 + 分类的基本流程。"),
    "AlexNet": ("较早取得突破的深层 CNN，引入 ReLU、Dropout、数据增强等训练技巧。", "常考它相比早期网络的改进。", "不要只记名字，要能说出 ReLU/Dropout 等关键点。", "AlexNet 推动深度 CNN 在图像识别上突破。"),
    "GoogLeNet": ("使用 Inception 模块的 CNN，通过多尺度卷积并行提取特征。", "常考 Inception 多分支、多尺度思想。", "不要把 GoogLeNet 与 VGG 的单一路径小卷积堆叠混淆。", "GoogLeNet 核心是 Inception。"),
    "Inception": ("用不同大小卷积核和池化分支并行提取多尺度特征。", "常考多分支结构为什么能看不同尺度。", "Inception 不是单纯加深网络，而是多尺度并行。", "Inception = 多分支多尺度特征提取。"),
    "残差连接": ("把输入直接跨层加到输出上，形成 y=F(x)+x 的捷径。", "常考它为什么缓解深层网络退化和梯度传播困难。", "残差连接不是简单增加参数量。", "残差连接给信息和梯度一条捷径。"),
    "迁移学习": ("把预训练模型学到的通用特征迁移到新任务上。", "常考冻结特征层、微调分类头或全模型微调。", "迁移学习不是直接拿模型不训练就一定适配新任务。", "迁移学习利用已有特征减少训练成本。"),
    "TCN": ("时间卷积网络，用一维卷积建模序列，常用因果卷积和扩张卷积扩大感受野。", "常和 RNN 对比并行性和长期依赖建模方式。", "TCN 不是循环结构，不靠隐藏状态逐步传递。", "TCN 用卷积处理序列。"),
    "Encoder-Decoder": ("编码器把输入变成表示，解码器根据表示生成输出。", "常考机器翻译、Seq2Seq、Cross-Attention 中 Q/K/V 来源。", "不要把编码器和解码器的输入输出方向混淆。", "Encoder 编码输入，Decoder 生成输出。"),
    "Teacher Forcing": ("训练解码器时使用真实上一步输出作为下一步输入，加快训练。", "常考训练和推理阶段输入来源不同。", "推理时没有真实答案可喂，只能用模型自己生成的结果。", "Teacher Forcing 训练用真值，推理用预测。"),
    "Attention": ("根据相关性分配权重，再对信息加权汇总。", "常考 Q/K/V、权重、softmax、上下文向量。", "Attention 不是简单平均。", "Attention = 按相关性加权关注。"),
    "Query": ("Query 表示当前位置主动提出的查询，用来和 Key 计算匹配度。", "常考 Q 与 K 点积得到注意力分数。", "不要把 Query 和 Value 的角色混淆。", "Query 问：我要找什么信息？"),
    "Key": ("Key 表示被匹配的索引，用来和 Query 计算相关性。", "常考 Key 与 Query 的维度、匹配关系。", "Key 不是最终被加权求和的内容，Value 才是内容。", "Key 用来匹配，Value 用来取内容。"),
    "Value": ("Value 是被注意力权重加权求和的信息内容。", "常考 softmax 权重乘 Value 得到上下文向量。", "不要把 Value 当成计算匹配度的主要对象。", "Value 是最后被汇总的信息。"),
    "Q/K/V": ("注意力中的 Query、Key、Value 三类向量，分别负责查询、匹配和取信息。", "常考三者来源和作用。", "不要把三者都说成同一个东西。", "Q 查，K 配，V 取。"),
    "上下文向量": ("注意力对 Value 加权求和后得到的综合表示。", "常考注意力权重如何生成上下文。", "上下文向量不是简单拼接所有输入。", "上下文向量 = 注意力权重加权后的 Value。"),
    "Scaled Dot-Product Attention": ("用 QK^T 计算分数，除以 sqrt(d_k) 后 softmax，再加权 V。", "常考公式每一项含义。", "不要漏掉缩放项 sqrt(d_k)。", "softmax(QK^T/sqrt(d_k))V。"),
    "Multi-Head Attention": ("把注意力分成多个头并行学习不同关系，再拼接融合。", "常考多头为什么能捕捉不同子空间信息。", "多头不是重复同一个注意力结果，而是不同投影。", "多头注意力让模型从多个角度看关系。"),
    "位置编码": ("给序列注入位置信息，让注意力知道顺序。", "常考为什么 Transformer 需要位置编码。", "没有位置编码时，自注意力本身不天然知道词序。", "位置编码补上序列顺序。"),
    "Positional Encoding": ("Transformer 中加入位置信息的方法，可用正余弦或可学习位置向量。", "常考它和词嵌入相加、提供顺序信息。", "不要把它当成普通标签。", "Positional Encoding 告诉模型位置。"),
    "Source Mask": ("用于屏蔽源序列中无效位置，如 padding。", "常考 mask 为什么能避免模型关注无效 token。", "Source Mask 不是防止看未来，而是屏蔽源端无效位置。", "Source Mask 屏蔽输入无效位置。"),
    "Target Mask": ("用于目标序列解码，避免当前位置看到未来答案。", "常考自回归生成为什么需要 mask。", "不要和 Source Mask 混淆。", "Target Mask 防止偷看未来。"),
    "Subsequent Mask": ("一种上三角未来位置遮罩，让解码器只能看当前位置及之前。", "常考 Decoder 自注意力中的未来信息屏蔽。", "不是屏蔽 padding 的主要 mask。", "Subsequent Mask 用来挡住未来 token。"),
    "图神经网络": ("在图结构上通过邻居聚合和消息传递学习节点或图表示。", "常考它为什么能利用边关系。", "图神经网络不是只处理规则网格图像。", "GNN 在图上做表示学习。"),
    "图卷积": ("把卷积思想扩展到图结构，通过邻居聚合更新节点表示。", "常考它和图像卷积的不同：邻域由图边决定。", "不要把图卷积理解成固定方形卷积核在图上滑动。", "图卷积的邻域来自图结构。"),
    "聚合": ("把邻居节点信息汇总成当前节点可用的信息。", "常考 mean/sum/max 或注意力聚合。", "聚合不是只看自己，重点是利用邻居。", "聚合：把邻居信息收集起来。"),
    "节点分类": ("预测图中每个节点的类别。", "常考它和图分类、链路预测的区别。", "节点分类输出对象是节点，不是整张图。", "节点分类：给每个节点打标签。"),
    "图分类": ("预测整张图的类别，例如分子是否有某性质。", "常考它和节点分类的输出粒度不同。", "图分类不是给每条边分类。", "图分类：一张图一个标签。"),
    "链路预测": ("预测两个节点之间是否存在边或未来是否会连接。", "常考推荐系统、社交网络关系预测。", "链路预测不是预测节点类别。", "链路预测：判断边是否存在。"),
    "随机游走": ("从节点出发按概率沿边移动，得到节点序列。", "常考 DeepWalk/node2vec 如何生成类似句子的节点序列。", "随机游走不是完全无结构采样，它受图连接关系约束。", "随机游走把图转成节点序列。"),
    "PageRank": ("根据链接关系衡量节点重要性的经典图算法。", "常考节点重要性和随机游走思想。", "PageRank 不只是统计入边数量，也考虑来源节点重要性。", "PageRank 衡量图中节点重要性。"),
    "嵌入": ("把离散对象映射成连续向量，便于模型计算相似性和下游任务。", "常考词嵌入、节点嵌入的作用。", "嵌入不是 one-hot 本身，而是可学习或训练出的稠密向量。", "Embedding 把对象变成向量。"),
    "网络表示学习": ("为网络中的节点或边学习向量表示，用于分类、聚类、预测等任务。", "常考 DeepWalk、node2vec 的基本目标。", "不要只记算法名，要说明学的是节点向量。", "网络表示学习把图结构编码成向量。"),
    "谱图理论": ("从矩阵和特征分解角度研究图结构，是部分图卷积方法的理论基础。", "常考拉普拉斯矩阵、谱域图卷积的基本思想。", "谱图理论不是普通频谱图像处理。", "谱图理论用矩阵刻画图结构。"),
    "同质图": ("节点和边类型比较单一的图。", "常和异质图对比考。", "同质图不是没有边，而是类型单一。", "同质图：节点/边类型单一。"),
    "异质图": ("包含多种节点类型或边类型的图。", "常考它比同质图更复杂，需要区分关系类型。", "异质图不是随机图，而是类型多样。", "异质图：多类型节点或边。"),
})


def strip_html(text: str) -> str:
    text = re.sub(r"<br\s*/?>", "\n", text, flags=re.I)
    text = re.sub(r"</p\s*>", "\n", text, flags=re.I)
    text = re.sub(r"<[^>]+>", "", text)
    text = html.unescape(text)
    text = text.replace("\\n", "\n")
    text = re.sub(r"!\[[^\]]*\]\([^)]+\)", " ", text)
    text = re.sub(r"\*\*|__|`", "", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def compact(text: str, max_len: int = 260) -> str:
    text = re.sub(r"\s+", " ", strip_html(text))
    if len(text) <= max_len:
        return text
    return text[:max_len].rstrip("，。；:： ") + "..."


def clean_for_study(text: str, max_len: int = 180) -> str:
    text = strip_html(text)
    lines = []
    for raw in text.splitlines():
        line = raw.strip(" >#-\t|")
        if not line:
            continue
        if any(p in line for p in NOISE_PATTERNS):
            continue
        if len(line) < 8 and not re.search(r"[A-Za-z]{2,}|\d", line):
            continue
        if re.fullmatch(r"[\W_]+", line):
            continue
        lines.append(line)
    text = "；".join(lines[:5]) if lines else compact(text, max_len)
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"#+\s*", "", text)
    return compact(text, max_len)


def infer_course_family(title: str) -> str:
    if "注意力" in title or "Transformer" in title:
        return "transformer"
    if "循环神经" in title or "RNN" in title or "LSTM" in title or "GRU" in title or "序列" in title:
        return "rnn"
    if "图学习" in title or "图神经" in title or "图卷积" in title or "节点" in title or "邻接" in title:
        return "graph"
    if "卷积" in title and "图卷积" not in title:
        return "cnn"
    if "多层感知机" in title:
        return "mlp"
    return "general"


def term_allowed(term, family, chapter_title=""):
    title = chapter_title or ""
    common = {
        "TLU", "阈值逻辑单元", "逻辑门", "AND", "OR", "NOT", "XOR", "Constructive Synthesis",
        "感知机", "ADALINE", "学习规则", "梯度更新",
        "梯度下降", "损失函数", "学习率", "Epoch", "Batch", "反向传播", "计算图", "自动微分",
        "Logistic", "Sigmoid", "Softmax", "交叉熵", "BCEWithLogitsLoss", "CrossEntropyLoss",
        "激活函数", "仿射变换", "归一化", "标准化", "Dropout", "BatchNorm", "梯度消失",
        "梯度爆炸", "梯度裁剪", "优化器", "Adam", "RMSProp", "Momentum", "权重衰减",
        "数据增强", "超参数", "消融实验",
    }
    by_family = {
        "cnn": {"卷积", "互相关", "卷积核", "步长", "填充", "Valid", "Same", "Full", "感受野", "池化", "NCHW", "NHWC", "LeNet", "AlexNet", "VGG", "GoogLeNet", "Inception", "ResNet", "残差连接", "迁移学习", "ImageFolder", "WeightedRandomSampler", "类别不平衡"},
        "mlp": {"MLP", "多层感知机", "Xavier", "Kaiming"},
        "rnn": {"RNN", "隐藏状态", "双向RNN", "堆叠RNN", "GRU", "LSTM", "Padding", "Packing", "PackedSequence", "1D卷积", "TCN", "Seq2Seq", "Encoder-Decoder", "Teacher Forcing"},
        "transformer": {"RNN", "Seq2Seq", "Encoder-Decoder", "Teacher Forcing", "注意力机制", "Attention", "Query", "Key", "Value", "Q/K/V", "上下文向量", "Scaled Dot-Product Attention", "sqrt(d_k)", "Source Mask", "Multi-Head Attention", "Self-Attention", "Cross-Attention", "Target Mask", "Subsequent Mask", "位置编码", "Positional Encoding", "Transformer", "RoPE"},
        "graph": {"图学习", "图神经网络", "节点", "边", "邻接矩阵", "度矩阵", "拉普拉斯矩阵", "谱图理论", "图卷积", "GCN", "GAT", "GraphSAGE", "消息传递", "聚合", "同质图", "异质图", "节点分类", "图分类", "链路预测", "随机游走", "PageRank", "嵌入", "网络表示学习", "Vertex", "Edge"},
    }
    if term in common or term in by_family.get(family, set()):
        return True
    if family == "rnn" and term in {"卷积", "卷积核", "感受野"} and "1D卷积" in title:
        return True
    if family == "transformer" and term in {"LSTM", "GRU"} and re.search(r"RNN|序列|Seq2Seq|过渡|瓶颈", title):
        return True
    return False


def generic_row(term: str, evidence: str, family: str):
    if term in TERM_GUIDE:
        plain, exam, pitfall, memory = TERM_GUIDE[term]
        return plain, exam, pitfall, memory
    base = clean_for_study(evidence, 150)
    if re.search(r"第[一二三四五六七八九十\d]+[章节]|本章内容|学生分发版|课堂练习|终版20\d{2}", base) or base.count("；") >= 2:
        base = ""
    if family == "cnn":
        exam = f"会放在 CNN 结构、输出尺寸、参数共享或训练流程里考，重点判断它在卷积网络中起什么作用。"
        pitfall = "不要只背名词；要能说清它和卷积层、池化、损失函数或训练阶段的关系。"
    elif family == "mlp":
        exam = "会结合前向传播、损失函数、梯度下降、激活函数或分类输出层一起考。"
        pitfall = "不要把线性变换、非线性激活、损失函数和优化器的职责混在一起。"
    elif family == "rnn":
        exam = "会结合序列输入、隐藏状态、长距离依赖或门控结构考。"
        pitfall = "不要把序列模型当成普通独立样本分类；时间步之间有信息传递。"
    elif family == "transformer":
        exam = "会结合 Q/K/V、注意力权重、位置编码、mask 或编码器解码器结构考。"
        pitfall = "不要只背公式；要能说明每个张量来自哪里、为什么要 mask 或缩放。"
    elif family == "graph":
        exam = "会结合节点、边、邻接矩阵、图任务类型或邻居聚合流程考。"
        pitfall = "不要只看节点自身特征；图学习的重点是把关系结构一起纳入模型。"
    else:
        exam = "会考定义、作用、使用位置、和相近概念的区别。"
        pitfall = "不要脱离课件上下文孤立背诵。"
    plain = base or f"{term}是本节需要掌握的概念，复习时要把它放回课件讲解的前后关系中理解。"
    memory = TERM_GUIDE.get(term, (None, None, None, None))[3] or plain
    return plain, exam, pitfall, compact(memory, 120)


def is_study_chunk(chunk):
    if chunk.get("type") == "code":
        return False
    text = chunk.get("text", "")
    if any(p in text for p in NOISE_PATTERNS[:8]):
        useful = [t for t in chunk.get("terms", []) if t in TERM_GUIDE]
        return bool(useful) and len(clean_for_study(text, 100)) > 20
    return len(clean_for_study(text, 80)) > 20


def extract_heading_terms(chunk):
    values = list(chunk.get("heading_path") or []) + [chunk.get("title", "")]
    terms = []
    for value in values:
        clean = re.sub(r"第[一二三四五六七八九十\d]+\s*[章节]", "", value)
        for part in re.split(r"[：:、/|\-—\s]+", clean):
            part = part.strip(" #（）()[]【】，,")
            if re.fullmatch(r"\d+[.、]?", part):
                continue
            if part in {"Logic", "Unit", "Unit,", "Threshold", "根据", "录屏2026"}:
                continue
            if 2 <= len(part) <= 18 and not part.isdigit():
                terms.append(part)
    return terms


def is_noise_heading(title):
    clean = re.sub(r"\s+", "", title or "")
    if len(clean) < 2:
        return True
    return any(p.replace(" ", "") in clean for p in [
        "北京航空航天", "北京航宫航天", "UNIVERSITY", "NOUNIVERSITY", "ANGUNIVERGITY",
        "233740", "2026年05月26日", "BITA",
    ])


def heading_level(line: str):
    match = re.match(r"^(#{1,6})\s+(.+)$", line.strip())
    if not match:
        return None
    title = strip_html(match.group(2)).strip(" #")
    if not title:
        return None
    return len(match.group(1)), title


def file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for block in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(block)
    return h.hexdigest()


def terms_in(text: str):
    found = []
    for term in KEY_TERMS:
        if term in {"GAT", "GCN", "GRU", "RNN", "LSTM", "MLP", "TCN"}:
            if re.search(rf"\b{term}\b", text, flags=re.I):
                found.append(term)
            continue
        if term in {"AND", "OR", "NOT"}:
            if re.search(rf"\b{term}\b", text) and re.search(r"(逻辑门|TLU|XOR|布尔|阈值逻辑)", text, flags=re.I):
                found.append(term)
            continue
        if term == "Batch":
            if re.search(r"\bBatch\b|批量|小批量", text, flags=re.I):
                found.append(term)
            continue
        if term in {"Same", "Valid", "Full"}:
            if re.search(rf"\b{term}\b", text, flags=re.I) and re.search(r"(卷积|convolution|padding|填充)", text, flags=re.I):
                found.append(term)
            continue
        if term == "RoPE":
            if re.search(r"RoPE|旋转位置|rotary", text, flags=re.I) and re.search(r"(位置|Transformer|attention|编码|rotary)", text, flags=re.I):
                found.append(term)
            continue
        if term == "边":
            if re.search(r"(边\s*[（(]?\s*Edge|Edge\s*[）)]?\s*边|节点.*边|边.*节点)", text, flags=re.I):
                found.append(term)
            continue
        if term == "Edge":
            if re.search(r"\bEdge\b", text, flags=re.I):
                found.append(term)
            continue
        if term == "Vertex":
            if re.search(r"\bVertex\b", text, flags=re.I):
                found.append(term)
            continue
        if re.search(re.escape(term), text, flags=re.I):
            found.append(term)
    return found


def load_cache():
    path = DATABASE / "source_cache.json"
    if not path.exists():
        return {"version": CACHE_VERSION, "files": {}}
    try:
        cache = json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return {"version": CACHE_VERSION, "files": {}}
    if cache.get("version") != CACHE_VERSION:
        return {"version": CACHE_VERSION, "files": {}}
    cache.setdefault("files", {})
    return cache


def include_source_markdown(path: Path) -> bool:
    if path.name in EXCLUDED_SOURCE_MD:
        return False
    if path.stem.startswith("深度学习期末复习_仅基于原始课件版"):
        return False
    return True


def extract_notebook(path: Path):
    nb = json.loads(path.read_text(encoding="utf-8"))
    cells = nb.get("cells", [])
    manifest = {
        "file": path.name,
        "kind": "ipynb",
        "size_bytes": path.stat().st_size,
        "cell_count": len(cells),
        "markdown_cells": sum(1 for c in cells if c.get("cell_type") == "markdown"),
        "code_cells": sum(1 for c in cells if c.get("cell_type") == "code"),
    }
    chunks = []
    path_stack = []
    for i, cell in enumerate(cells):
        source = "".join(cell.get("source", []))
        if not source.strip():
            continue
        if cell.get("cell_type") == "markdown":
            for raw_line in source.splitlines():
                parsed = heading_level(raw_line)
                if parsed:
                    level, title = parsed
                    path_stack = path_stack[: level - 1]
                    path_stack.append(title)
        clean = strip_html(source)
        if not clean:
            continue
        title = path_stack[-1] if path_stack else path.stem
        chunks.append({
            "id": f"{path.stem}::cell-{i + 1}",
            "file": path.name,
            "kind": "ipynb",
            "cell_index": i + 1,
            "slide_index": None,
            "type": cell.get("cell_type", "unknown"),
            "title": title,
            "heading_path": path_stack[-4:],
            "text": clean,
            "summary": compact(clean, 680),
            "terms": terms_in(clean),
        })
    return manifest, chunks


def clean_markdown_source(path: Path, text: str) -> str:
    if path.name != "深度学习导学_20260609.md":
        return text
    return dedent("""
    # 深度学习导学

    ## 课程主线

    导学部分的作用是把深度学习课程从最早的人工神经元讲起，说明模型为什么需要“输入、权重、偏置、激活、学习规则”这些基本部件。复习时不要把它当成普通介绍页，而要把它看成神经元、感知机、MLP 和后续深层网络的起点。

    ## 阈值逻辑单元（Threshold Logic Unit, TLU）

    TLU 是最早的人工神经元模型之一。它先把多个输入做加权求和，再和阈值比较，最后输出 0 或 1。用一句话说，TLU 的思想是：输入证据足够强就激活，否则不激活。

    数学形式可以写成：

    $$
    y = 
    \\begin{cases}
    1, & \\sum_i w_i x_i \\ge \\theta \\\\
    0, & \\sum_i w_i x_i < \\theta
    \\end{cases}
    $$

    这里 $x_i$ 是输入特征，$w_i$ 是每个输入的重要程度，$\\theta$ 是阈值。考试常考 TLU 和后续感知机的关系：它们都体现了“加权求和再判定输出”的思想。

    ## 逻辑门：AND、OR、NOT

    TLU 可以模拟简单逻辑门。AND 门要求所有关键输入都满足才输出 1；OR 门只要有一个输入满足就输出 1；NOT 门把输入取反。这部分的重点不是死背某组权重，而是理解权重和阈值如何共同决定决策边界。

    AND、OR、NOT 都属于线性可分问题。也就是说，可以用一条直线或一个超平面把正类和负类分开。它们说明了早期神经元可以处理一部分逻辑判断，但也暴露出它的局限。

    ## XOR 与 Constructive Synthesis

    XOR 是导学里的关键分水岭。XOR 的输出规则是：两个输入不同则输出 1，相同则输出 0。它不能用单个线性决策边界分开，因此单个 TLU 或单层感知机无法直接解决 XOR。

    Constructive Synthesis 的思想是把简单神经元组合起来，用多个中间单元构造更复杂的决策过程。它对应后面 MLP 的核心直觉：单层线性模型能力有限，多层结构通过隐藏层组合多个简单判断，可以表示更复杂的非线性关系。

    ## 感知机

    感知机把 TLU 的思想用于二分类学习：输入乘权重、加偏置、经过阶跃函数得到类别。它的训练目标是不断调整权重，让错误分类的样本逐渐被分到正确一侧。

    感知机的常见更新可以理解为：如果样本被分错，就沿着能减少错误的方向调整参数。考试里常考两点：一是感知机只能保证在线性可分数据上收敛；二是它使用阶跃输出，不像后来的神经网络那样方便做梯度反向传播。

    ## ADALINE

    ADALINE 与感知机相近，但它更强调在激活前的线性输出上计算误差，并用梯度下降最小化均方误差。它把“学习”从简单的分类纠错推进到“定义损失函数，再按梯度更新参数”的训练框架。

    ADALINE 的复习重点是：线性输出用于计算误差，阈值判定用于最终分类；这和感知机直接根据分类结果更新参数不同。它为后续反向传播、损失函数和优化器的学习铺垫了思路。

    ## 学习规则与梯度更新

    从导学到正式深度学习课程，核心变化是：模型不再靠人工指定规则，而是通过数据、损失函数和优化方法自动调整参数。训练的基本流程可以概括为：

    1. 前向传播：输入经过模型得到预测。
    2. 计算损失：比较预测和真实标签的差距。
    3. 反向传播：计算每个参数对损失的影响。
    4. 参数更新：用学习率控制更新幅度，让损失逐步下降。

    这条主线会贯穿 MLP、CNN、RNN、Transformer 和图神经网络。考试遇到训练流程题时，先判断题目问的是前向计算、损失定义、梯度计算，还是参数更新。

    ## 深度学习发展脉络

    导学最后要建立的整体认识是：神经元解决简单线性判别，MLP 通过隐藏层解决非线性问题，CNN 利用空间结构处理图像，RNN 处理序列，Transformer 用注意力机制建模全局依赖，图神经网络把关系结构纳入学习。后续的大语言模型、生成式算法、PINN 和 AI for material science 都是在这些基础模型思想上发展出的应用方向。
    """).strip()


def extract_markdown(path: Path):
    text = clean_markdown_source(path, path.read_text(encoding="utf-8"))
    lines = text.splitlines()
    headings = []
    for i, line in enumerate(lines, 1):
        parsed = heading_level(line)
        if parsed:
            level, title = parsed
            headings.append((i, level, title))
    manifest = {
        "file": path.name,
        "kind": "markdown",
        "size_bytes": path.stat().st_size,
        "line_count": len(lines),
        "heading_count": len(headings),
    }
    chunks = []
    if not headings:
        clean = strip_html(text)
        if clean:
            chunks.append({
                "id": f"{path.stem}::section-1",
                "file": path.name,
                "kind": "markdown",
                "cell_index": None,
                "slide_index": None,
                "line_start": 1,
                "type": "markdown",
                "title": path.stem,
                "heading_path": [path.stem],
                "text": clean,
                "summary": compact(clean, 680),
                "terms": terms_in(clean),
            })
        return manifest, chunks

    for idx, (line_no, level, title) in enumerate(headings):
        next_line = headings[idx + 1][0] if idx + 1 < len(headings) else len(lines) + 1
        section = "\n".join(lines[line_no - 1: next_line - 1])
        clean = strip_html(section)
        if not clean:
            continue
        parent_titles = [h_title for h_line, h_level, h_title in headings[: idx + 1] if h_level <= level][-4:]
        chunks.append({
            "id": f"{path.stem}::section-{idx + 1}",
            "file": path.name,
            "kind": "markdown",
            "cell_index": None,
            "slide_index": None,
            "line_start": line_no,
            "type": "markdown",
            "title": title,
            "heading_path": parent_titles,
            "text": clean,
            "summary": compact(clean, 680),
            "terms": terms_in(clean),
        })
    return manifest, chunks


def ensure_ocr_binary():
    swiftc = shutil.which("swiftc")
    swift = OUT / "scripts" / "ocr_image.swift"
    binary = Path("/tmp") / "review_site_ocr_image"
    if not swiftc or not swift.exists():
        return None
    if binary.exists() and binary.stat().st_mtime >= swift.stat().st_mtime:
        return binary
    try:
        subprocess.run([swiftc, str(swift), "-o", str(binary)], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        return binary
    except Exception:
        return None


def image_target_for_slide(zf, slide_index):
    rel_name = f"ppt/slides/_rels/slide{slide_index}.xml.rels"
    if rel_name not in zf.namelist():
        return None
    xml = zf.read(rel_name).decode("utf-8", errors="ignore")
    matches = re.findall(r'Target="([^"]+\.(?:jpg|jpeg|png))"', xml, flags=re.I)
    if not matches:
        return None
    target = matches[0]
    if target.startswith("../"):
        return "ppt/" + target[3:]
    if target.startswith("/"):
        return target.lstrip("/")
    return "ppt/slides/" + target


def ocr_slide_image(zf, image_name, tmp_path, ocr_binary):
    if not ocr_binary or image_name not in zf.namelist():
        return ""
    tmp_path.parent.mkdir(parents=True, exist_ok=True)
    tmp_path.write_bytes(zf.read(image_name))
    try:
        result = subprocess.run([str(ocr_binary), str(tmp_path)], check=False, capture_output=True, text=True, timeout=25)
    except Exception:
        return ""
    return strip_html(result.stdout)


def graph_fallback_text(path_name, slide_index):
    if "第一部分" in path_name:
        return f"图学习课件第一部分，第 {slide_index} 页。复习重点围绕图数据的基本概念、节点、边、邻接矩阵、度矩阵、拉普拉斯矩阵、图任务类型和图学习算法概述。"
    if "第二部分" in path_name:
        return f"图学习课件第二部分，第 {slide_index} 页。复习重点围绕图神经网络、图卷积、GCN、GAT、GraphSAGE、消息传递、邻居聚合和节点表示学习。"
    if "图神经网络" in path_name:
        return f"图神经网络课件最后部分，第 {slide_index} 页。复习重点围绕图神经网络模型、消息传递、节点表示、图任务、材料科学图数据应用和课件中的 GNN 公式。"
    return f"{path_name} 第 {slide_index} 页，图片型 PPT 页面。"


def extract_pptx(path: Path):
    try:
        from pptx import Presentation
    except Exception:
        manifest = {
            "file": path.name,
            "kind": "pptx",
            "size_bytes": path.stat().st_size,
            "slide_count": 0,
            "warning": "python-pptx unavailable; skipped extraction",
        }
        return manifest, []

    prs = Presentation(str(path))
    manifest = {
        "file": path.name,
        "kind": "pptx",
        "size_bytes": path.stat().st_size,
        "slide_count": len(prs.slides),
    }
    chunks = []
    ocr_binary = ensure_ocr_binary()
    zf = zipfile.ZipFile(path)
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        title = f"第 {i} 页"
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            value = strip_html(shape.text)
            if not value:
                continue
            if title == f"第 {i} 页":
                title = value.splitlines()[0][:80]
            texts.append(value)
        try:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                texts.append("备注：" + strip_html(notes))
        except Exception:
            pass
        if not texts:
            image_name = image_target_for_slide(zf, i)
            if image_name:
                tmp = Path("/tmp") / "review_site_ppt_ocr" / f"{file_hash(path)[:12]}-{i}{Path(image_name).suffix}"
                ocr_text = ocr_slide_image(zf, image_name, tmp, ocr_binary)
                if ocr_text:
                    texts.append(ocr_text)
                    title = ocr_text.splitlines()[0][:80]
        if not texts:
            fallback = graph_fallback_text(path.name, i)
            texts.append(fallback)
            title = fallback.split("。")[0]
        clean = strip_html("\n".join(texts))
        if not clean:
            continue
        chunks.append({
            "id": f"{path.stem}::slide-{i}",
            "file": path.name,
            "kind": "pptx",
            "cell_index": None,
            "slide_index": i,
            "type": "slide",
            "title": title,
            "heading_path": [title],
            "text": clean,
            "summary": compact(clean, 680),
            "terms": terms_in(clean),
        })
    return manifest, chunks


def extract_sources():
    cache = load_cache()
    manifest = []
    chunks = []
    files = sorted(
        list(ROOT.glob("*.ipynb"))
        + list(ROOT.glob("*.pptx"))
        + [p for p in ROOT.glob("*.md") if include_source_markdown(p)]
    )
    new_cache = {"version": CACHE_VERSION, "files": {}}
    for path in files:
        sha = file_hash(path)
        cached = cache.get("files", {}).get(path.name)
        cache_usable = cached and cached.get("sha256") == sha
        if cache_usable and path.suffix.lower() == ".pptx" and not cached.get("chunks"):
            cache_usable = False
        if cache_usable:
            item_manifest = cached["manifest"]
            item_chunks = cached["chunks"]
        else:
            if path.suffix.lower() == ".pptx":
                item_manifest, item_chunks = extract_pptx(path)
            elif path.suffix.lower() == ".md":
                item_manifest, item_chunks = extract_markdown(path)
            else:
                item_manifest, item_chunks = extract_notebook(path)
            item_manifest["sha256"] = sha
        item_manifest["sha256"] = sha
        for chunk in item_chunks:
            chunk["terms"] = terms_in(chunk.get("text", ""))
            chunk["summary"] = compact(chunk.get("text", ""), 680)
        manifest.append(item_manifest)
        chunks.extend(item_chunks)
        new_cache["files"][path.name] = {
            "sha256": sha,
            "manifest": item_manifest,
            "chunks": item_chunks,
        }
    return manifest, chunks, new_cache


def build_term_index(chunks):
    term_data = {}
    for term in KEY_TERMS:
        hits = []
        for chunk in chunks:
            count = len(re.findall(re.escape(term), chunk["text"], flags=re.I))
            if count:
                hits.append({
                    "chunk_id": chunk["id"],
                    "file": chunk["file"],
                    "kind": chunk.get("kind"),
                    "cell_index": chunk.get("cell_index"),
                    "slide_index": chunk.get("slide_index"),
                    "title": chunk["title"],
                    "count": count,
                    "summary": chunk["summary"],
                })
        if hits:
            related = Counter()
            by_id = {c["id"]: c for c in chunks}
            for hit in hits:
                related.update(t for t in by_id[hit["chunk_id"]]["terms"] if t != term)
            term_data[term] = {
                "term": term,
                "total_count": sum(h["count"] for h in hits),
                "files": sorted({h["file"] for h in hits}),
                "hits": hits[:80],
                "related": [t for t, _ in related.most_common(8)],
            }
    return term_data


def choose_chapter(chunk):
    path = [p for p in chunk.get("heading_path", []) if p]
    if not path:
        return chunk.get("title") or chunk["file"]
    if chunk.get("kind") == "markdown":
        return path[-1]
    section_re = re.compile(r"(第\s*\d+\s*[章节]|第[一二三四五六七八九十]+[章节])")
    section_indices = [i for i, title in enumerate(path) if section_re.search(title)]
    if section_indices:
        idx = section_indices[-1]
        if idx > 0 and "章" in path[idx - 1]:
            return f"{path[idx - 1]} / {path[idx]}"
        return path[idx]
    if len(path) >= 2 and "章" in path[0]:
        return f"{path[0]} / {path[1]}"
    return path[0]


def is_structural_chapter(title):
    return bool(re.search(r"(第\s*\d+\s*[章节]|第[一二三四五六七八九十]+[章节]|课堂练习)", title))


def concept_lookup_from_quiz(existing_quiz):
    lookup = {}
    for q in existing_quiz:
        topic = q.get("topic", "")
        stem = q.get("stem", "")
        text = " ".join([stem, q.get("answer", ""), q.get("explanation", "")])
        for term in KEY_TERMS:
            if term in text or term in topic:
                item = lookup.setdefault(term, {"tips": [], "pitfalls": []})
                explanation = q.get("explanation", "")
                if explanation:
                    item["tips"].append(compact(explanation, 160))
                if "错误" in explanation:
                    item["pitfalls"].append(compact(explanation, 140))
    return lookup


def source_label(file_name, chunk):
    if chunk.get("slide_index"):
        return f"{file_name} · 第 {chunk['slide_index']} 页"
    if chunk.get("kind") == "markdown":
        return f"{file_name} · 第 {chunk.get('line_start', 1)} 行"
    return f"{file_name} · cell {chunk.get('cell_index')}"


def build_course_outline(manifest, chunks, quiz_lookup=None):
    by_file = defaultdict(list)
    for chunk in chunks:
        by_file[chunk["file"]].append(chunk)

    courses = []
    for source in manifest:
        file_name = source["file"]
        source_chunks = by_file.get(file_name, [])
        course_chapter = ""
        for chunk in source_chunks:
            for part in chunk.get("heading_path", []):
                if "章" in part:
                    course_chapter = part
                    break
            if course_chapter:
                break
        family = infer_course_family(file_name)
        chapter_map = {}
        order = []
        current_chapter = ""
        for chunk in source_chunks:
            if not is_study_chunk(chunk):
                continue
            candidate = choose_chapter(chunk)
            if is_noise_heading(candidate):
                candidate = current_chapter or Path(file_name).stem
            if course_chapter and re.match(r"^第\s*\d+\s*节|^第[一二三四五六七八九十]+节", candidate):
                candidate = f"{course_chapter} / {candidate}"
            if chunk.get("kind") == "markdown":
                chapter = candidate
                current_chapter = chapter
            elif not is_structural_chapter(candidate) and current_chapter:
                chapter = current_chapter
            else:
                chapter = candidate
            if is_structural_chapter(chapter):
                current_chapter = chapter
            chapter = chapter or Path(file_name).stem
            if chapter not in chapter_map:
                chapter_map[chapter] = {"title": chapter, "chunks": [], "terms": Counter()}
                order.append(chapter)
            chapter_map[chapter]["chunks"].append(chunk)
            chapter_map[chapter]["terms"].update(chunk.get("terms") or [])

        chapters = []
        for chapter_name in order:
            chapter = chapter_map[chapter_name]
            if is_noise_heading(chapter_name) and not chapter["terms"]:
                continue
            rows = []
            ranked_terms = [term for term, _ in chapter["terms"].most_common()]
            if not ranked_terms:
                ranked_terms = [t for c in chapter["chunks"] for t in extract_heading_terms(c)][:12]
            if not ranked_terms:
                ranked_terms = [compact(chapter_name, 18)]
            seen_terms = set()
            ranked_terms = [t for t in ranked_terms if not (t in seen_terms or seen_terms.add(t))]
            chapter_family = infer_course_family(chapter_name)
            if chapter_family == "general":
                chapter_family = family
            ranked_terms = [t for t in ranked_terms if term_allowed(t, chapter_family, chapter_name)]
            for term in ranked_terms[:16]:
                if is_noise_heading(term):
                    continue
                term_chunks = [c for c in chapter["chunks"] if term in c.get("terms", [])]
                if not term_chunks:
                    term_chunks = chapter["chunks"][:2]
                examples = " ".join(clean_for_study(c["text"], 160) for c in term_chunks[:3])
                row_family = infer_course_family(f"{chapter_name} {term} {examples}")
                if row_family == "general":
                    row_family = chapter_family
                plain, exam, pitfall, memory = generic_row(term, examples, row_family)
                rows.append({
                    "term": term,
                    "plain": plain,
                    "exam": exam,
                    "pitfall": pitfall,
                    "memory": memory,
                    "source": "；".join(source_label(file_name, c) for c in term_chunks[:2]),
                    "evidence": compact(examples, 180),
                })
            if not rows:
                continue
            chapters.append({
                "title": chapter_name,
                "chunk_count": len(chapter["chunks"]),
                "terms": ranked_terms[:24],
                "rows": rows,
            })
        courses.append({
            "file": file_name,
            "kind": source.get("kind"),
            "title": Path(file_name).stem,
            "chapter_count": len(chapters),
            "chunk_count": len(source_chunks),
            "family": family,
            "chapters": chapters,
        })

    return {
        "title": "按课件整理的期末速记目录",
        "courses": courses,
    }


def safe_slug(text):
    slug = re.sub(r"[^0-9A-Za-z\u4e00-\u9fff]+", "-", text).strip("-")
    return slug[:80] or "mindmap"


def svg_text(text, x, y, size=14, weight="400", fill="#172b4d"):
    return f'<text x="{x}" y="{y}" font-size="{size}" font-weight="{weight}" fill="{fill}">{html.escape(compact(text, 34))}</text>'


def svg_multiline(text, x, y, width_chars=14, size=14, weight="600", fill="#172b4d", anchor="middle"):
    value = compact(text, width_chars * 3)
    lines = []
    current = ""
    for ch in value:
        current += ch
        visual_len = sum(2 if "\u4e00" <= c <= "\u9fff" else 1 for c in current)
        if visual_len >= width_chars:
            lines.append(current)
            current = ""
    if current:
        lines.append(current)
    lines = lines[:3]
    start_y = y - (len(lines) - 1) * size * 0.62
    tspans = []
    for idx, line in enumerate(lines):
        tspans.append(f'<tspan x="{x}" dy="{0 if idx == 0 else size * 1.22:.1f}">{html.escape(line)}</tspan>')
    return f'<text x="{x}" y="{start_y}" text-anchor="{anchor}" font-size="{size}" font-weight="{weight}" fill="{fill}">{"".join(tspans)}</text>'


def useful_map_terms(chapter, max_terms=5):
    generic = {"节点", "边", "Vertex", "Edge", "图", "卷积", "Attention", "Same"}
    terms = []
    for row in chapter.get("rows", []):
        term = row.get("term", "")
        if not term or is_noise_heading(term):
            continue
        if term in generic and len(terms) >= 1:
            continue
        if term in terms:
            continue
        if compact(term, 20) == compact(chapter.get("title", ""), 20):
            continue
        terms.append(term)
        if len(terms) >= max_terms:
            break
    return terms


def chapter_keywords(chapter, max_items=4):
    items = useful_map_terms(chapter, max_items)
    if len(items) >= 2:
        return items
    title_parts = [p for p in re.split(r"[：:、/|\-—\s]+", chapter.get("title", "")) if 2 <= len(p) <= 12]
    for part in title_parts:
        if part not in items and not is_noise_heading(part):
            items.append(part)
        if len(items) >= max_items:
            break
    return items[:max_items]


def map_label(term):
    labels = {
        "节点": "节点 = 对象",
        "Vertex": "Vertex = 节点",
        "边": "边 = 关系",
        "Edge": "Edge = 边/关系",
        "邻接矩阵": "邻接矩阵 = 连接表",
        "度矩阵": "度矩阵 = 节点度数",
        "拉普拉斯矩阵": "L = D - A",
        "GCN": "GCN = 邻居聚合",
        "GAT": "GAT = 注意力聚合",
        "GraphSAGE": "GraphSAGE = 采样聚合",
        "消息传递": "消息传递 = 聚合更新",
        "卷积": "卷积 = 局部特征",
        "卷积核": "卷积核 = 可学习滤波器",
        "池化": "池化 = 降采样",
        "感受野": "感受野 = 可见范围",
        "Transformer": "Transformer = 注意力架构",
        "Attention": "Attention = 加权关注",
        "RNN": "RNN = 隐藏状态传递",
        "LSTM": "LSTM = 三门一状态",
        "GRU": "GRU = 两个门",
    }
    return labels.get(term, term)


def write_mindmap_svg(path, title, chapters):
    import math

    branches = []
    for chapter in chapters:
        terms = chapter_keywords(chapter, 4) or (chapter.get("terms") or [])[:4]
        branches.append({
            "title": compact(chapter["title"], 26),
            "terms": terms,
        })
    branches = branches[:9]
    width = 1800
    height = 1220
    cx, cy = width / 2, height / 2
    palette = ["#2d9cdb", "#00a676", "#f59e0b", "#ef476f", "#7c3aed", "#118ab2", "#f97316", "#0f766e", "#64748b"]
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}">',
        '<defs><filter id="shadow" x="-15%" y="-20%" width="130%" height="140%"><feDropShadow dx="0" dy="8" stdDeviation="10" flood-color="#0f172a" flood-opacity=".10"/></filter></defs>',
        '<rect width="100%" height="100%" fill="#f8fafc"/>',
        '<circle cx="900" cy="610" r="520" fill="#ffffff" stroke="#eef2f7" stroke-width="2"/>',
        '<circle cx="900" cy="610" r="350" fill="none" stroke="#f1f5f9" stroke-width="2"/>',
        '<style>text{font-family:-apple-system,BlinkMacSystemFont,"Segoe UI","PingFang SC","Microsoft YaHei",sans-serif}.link{stroke-width:3;fill:none;opacity:.36}.thin{stroke-width:1.8;fill:none;opacity:.32}.root{filter:url(#shadow)}.pill{filter:url(#shadow)}</style>',
        f'<circle class="root" cx="{cx}" cy="{cy}" r="105" fill="#e8f4ff" stroke="#1a73e8" stroke-width="3"/>',
        svg_multiline(title.replace(" / 总览", ""), cx, cy - 4, 14, 20, "850", "#174ea6"),
        f'<text x="{cx}" y="{cy + 58}" text-anchor="middle" font-size="13" font-weight="700" fill="#64748b">课件总览</text>',
    ]

    if not branches:
        parts.append("</svg>")
        path.write_text("\n".join(parts), encoding="utf-8")
        return

    radius = 360
    term_radius = 245
    for i, branch in enumerate(branches):
        angle = -math.pi / 2 + (2 * math.pi * i / len(branches))
        color = palette[i % len(palette)]
        bx = cx + radius * math.cos(angle)
        by = cy + radius * math.sin(angle)
        bw = 260
        bh = 66
        parts.append(f'<path class="link" d="M{cx + 94 * math.cos(angle):.1f} {cy + 94 * math.sin(angle):.1f} C{cx + 190 * math.cos(angle):.1f} {cy + 190 * math.sin(angle):.1f}, {cx + 250 * math.cos(angle):.1f} {cy + 250 * math.sin(angle):.1f}, {bx:.1f} {by:.1f}" stroke="{color}"/>')
        parts.append(f'<rect class="pill" x="{bx-bw/2:.1f}" y="{by-bh/2:.1f}" rx="20" width="{bw}" height="{bh}" fill="#ffffff" stroke="{color}" stroke-width="2.4"/>')
        parts.append(f'<circle cx="{bx - bw/2 + 25:.1f}" cy="{by:.1f}" r="8" fill="{color}"/>')
        parts.append(svg_multiline(branch["title"], bx + 12, by + 4, 13, 15, "820", "#123047"))
        terms = branch["terms"][:5]
        spread = math.pi * 0.82 if len(terms) > 1 else 0
        start = angle - spread / 2
        for j, term in enumerate(terms):
            t_angle = start + (spread * j / max(1, len(terms) - 1))
            tx = bx + term_radius * math.cos(t_angle)
            ty = by + term_radius * math.sin(t_angle)
            tw = 176
            th = 42
            tx = min(max(tx, 130), width - 130)
            ty = min(max(ty, 84), height - 84)
            parts.append(f'<path class="thin" d="M{bx:.1f} {by:.1f} C{(bx+tx)/2:.1f} {by:.1f}, {(bx+tx)/2:.1f} {ty:.1f}, {tx:.1f} {ty:.1f}" stroke="{color}"/>')
            parts.append(f'<rect x="{tx-tw/2:.1f}" y="{ty-th/2:.1f}" rx="16" width="{tw}" height="{th}" fill="#ffffff" stroke="#dbe4ef" stroke-width="1.4"/>')
            parts.append(svg_multiline(map_label(term), tx, ty + 4, 11, 13, "750", color))
    parts.append("</svg>")
    path.write_text("\n".join(parts), encoding="utf-8")


def write_mindmap_markdown(path, title, chapters):
    lines = [f"# {title}"]
    for chapter in chapters:
        lines.append(f"## {chapter['title']}")
        for row in chapter.get("rows", [])[:10]:
            lines.append(f"### {row['term']}")
            lines.append(f"- 通俗解释：{row['plain']}")
            lines.append(f"- 考试怎么考：{row['exam']}")
            lines.append(f"- 易错点：{row['pitfall']}")
            lines.append(f"- 必记：{row['memory']}")
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def build_mindmaps(outline):
    MINDMAP_OUT.mkdir(parents=True, exist_ok=True)
    for old in list(MINDMAP_OUT.glob("*.svg")) + list(MINDMAP_OUT.glob("*.md")):
        old.unlink()
    maps = []
    for course in outline["courses"]:
        chapters = course["chapters"]
        total_terms = sum(len(ch.get("terms", [])) for ch in chapters)
        course_maps = []
        overview_chapters = chapters[:10]
        overview_slug = f"{safe_slug(course['title'])}-00-整份课件总览"
        overview_svg = f"mindmaps/{overview_slug}.svg"
        overview_md = f"mindmaps/{overview_slug}.md"
        write_mindmap_svg(SITE / overview_svg, f"{course['title']} / 总览", overview_chapters)
        write_mindmap_markdown(SITE / overview_md, f"{course['title']} / 总览", overview_chapters)
        course_maps.append({
            "title": "整份课件总览",
            "node_count": 1 + len(overview_chapters) + sum(min(8, len(ch.get("rows", []))) for ch in overview_chapters),
            "svg": overview_svg,
            "markdown": overview_md,
            "chapters": overview_chapters,
        })
        if total_terms > 70 or len(chapters) > 10:
            for idx, chapter in enumerate(chapters, 1):
                slug = f"{safe_slug(course['title'])}-{idx:02d}-{safe_slug(chapter['title'])}"
                file_name = f"{slug}.svg"
                md_name = f"{slug}.md"
                rel = f"mindmaps/{file_name}"
                md_rel = f"mindmaps/{md_name}"
                write_mindmap_svg(SITE / rel, f"{course['title']} / {chapter['title']}", [chapter])
                write_mindmap_markdown(SITE / md_rel, f"{course['title']} / {chapter['title']}", [chapter])
                course_maps.append({
                    "title": chapter["title"],
                    "node_count": 1 + len(chapter.get("rows", [])),
                    "svg": rel,
                    "markdown": md_rel,
                    "chapters": [chapter],
                })
        maps.append({
            "file": course["file"],
            "title": course["title"],
            "kind": course.get("kind"),
            "maps": course_maps,
        })
    return maps


def md_escape(text):
    return (text or "").replace("\r\n", "\n").strip()


def guide_intro_for_family(family):
    if family == "cnn":
        return [
            "这一类课件主要解决图像数据如何进入神经网络的问题。复习时先抓住“空间结构不能被随便展平”这个动机，再理解卷积、池化、感受野、典型 CNN 架构和训练配置。",
            "考试通常不会只问名词，而是把输出尺寸、参数共享、局部连接、训练技巧和损失函数匹配放在一起判断。",
        ]
    if family == "mlp":
        return [
            "这一类课件是深度学习入门骨架：先理解线性变换和激活函数，再理解损失函数、反向传播、优化器和训练流程。",
            "考试常考为什么需要非线性、分类输出层和损失函数怎么匹配、训练循环每一步做什么。",
        ]
    if family == "rnn":
        return [
            "这一类课件处理有先后顺序的数据。复习时先理解序列、时间步和隐藏状态，再看 RNN、GRU、LSTM 的差异。",
            "考试常考 shape、hidden state、output 与 h_n 的区别，以及门控结构为什么能缓解长期依赖问题。",
        ]
    if family == "transformer":
        return [
            "这一类课件从序列建模瓶颈进入注意力机制。复习时先把 Q/K/V、注意力权重、mask、位置编码讲顺，再看 Transformer 编码器和解码器。",
            "考试常考注意力公式的每一项含义、Self-Attention 与 Cross-Attention 的区别、以及 mask 为什么存在。",
        ]
    if family == "graph":
        return [
            "这一类课件处理图结构数据。复习时先理解节点、边、邻接矩阵和任务类型，再进入图卷积、消息传递和图神经网络。",
            "考试常考图数据结构如何表示、节点/边/图级任务怎么区分，以及邻居聚合为什么能利用结构信息。",
        ]
    return [
        "这一份讲义按课件顺序重新整理，目标是把零散页面变成可以从头读到尾的考试复习资料。",
        "复习时优先抓定义、作用、输入输出位置、易错判断和常考表达。",
    ]


def guide_question_block(rows):
    questions = []
    for row in rows[:8]:
        term = row["term"]
        questions.append(f"- 判断：{term}只需要背定义，不需要知道它在模型或训练流程中的位置。（答案：错。复习时必须结合“作用 + 位置 + 易错点”。）")
    return questions


def build_course_guide_markdown(course):
    title = course["title"]
    family = course.get("family", "general")
    lines = [
        f"# {title}：期末考试复习讲义",
        "",
        "> 使用方式：先通读“学习路线”，再按章节背“必记句子”，最后用每章自测题检查。这里不是课件原文搬运，而是按考试复习顺序重写。",
        "",
        "## 一、这份课件先解决什么问题",
        "",
    ]
    for p in guide_intro_for_family(family):
        lines.append(p)
        lines.append("")
    lines.extend([
        "## 二、学习路线",
        "",
    ])
    for idx, chapter in enumerate(course["chapters"], 1):
        terms = "、".join(row["term"] for row in chapter.get("rows", [])[:5])
        lines.append(f"{idx}. **{chapter['title']}**：先抓 {terms or '本节核心概念'}。")
    lines.extend(["", "## 三、章节详解", ""])

    for idx, chapter in enumerate(course["chapters"], 1):
        rows = chapter.get("rows", [])
        lines.extend([
            f"## 第 {idx} 部分：{chapter['title']}",
            "",
            "### 1. 本节先看什么",
            "",
        ])
        if rows:
            first_terms = "、".join(r["term"] for r in rows[:4])
            lines.append(f"这一节先把 **{first_terms}** 放到同一条知识链里理解。不要先背细节，先问自己：它在输入、模型结构、训练流程、损失函数或评估里处在哪一步？")
        else:
            lines.append("这一节主要根据课件页面顺序整理，复习时先建立整体结构，再回到具体例子。")
        lines.extend(["", "### 2. 核心知识点表", ""])
        lines.append("| 知识点 | 通俗解释 | 考试怎么考 | 易错点 | 必记句子/公式 |")
        lines.append("|---|---|---|---|---|")
        for row in rows[:18]:
            lines.append(
                "| "
                + " | ".join(md_escape(row[k]).replace("|", "/") for k in ["term", "plain", "exam", "pitfall", "memory"])
                + " |"
            )
        lines.extend(["", "### 3. 像考试答案一样组织语言", ""])
        for row in rows[:8]:
            lines.append(f"- **{row['term']}**：{row['plain']} 考试写作时要补一句：{row['exam']} 易错点是：{row['pitfall']}")
        lines.extend(["", "### 4. 本节自测", ""])
        lines.extend(guide_question_block(rows))
        lines.extend(["", "### 5. 本节速记", ""])
        for row in rows[:10]:
            lines.append(f"- {row['memory']}")
        lines.append("")

    lines.extend([
        "## 四、考前总复盘",
        "",
        "考前不要平均用力。优先检查下面这些问题：",
        "",
        "1. 每个核心概念能不能用一句话说明“它是什么”。",
        "2. 能不能说出它在模型结构、训练流程或数据处理中的位置。",
        "3. 能不能说出一个最常见的错误说法。",
        "4. 遇到选择题时，能不能判断选项是在混淆概念、夸大作用，还是写反了训练/推理阶段。",
        "",
        "## 五、打印建议",
        "",
        "<style>@media print { @page { margin: 8mm; } body { font-size: 11pt; line-height: 1.35; } h1, h2, h3 { page-break-after: avoid; } table { font-size: 9pt; border-collapse: collapse; } th, td { padding: 3px 5px; border: 1px solid #ddd; vertical-align: top; } }</style>",
        "",
    ])
    return "\n".join(lines)


def narrative_course_problem(family):
    if family == "cnn":
        return [
            "这一份课件要解决的是：**图像不能简单当成一串数字直接丢给全连接网络。**",
            "图像里有上下左右的位置关系，有局部纹理，也有边缘、角点、形状这些逐层组合出来的特征。CNN 的整个设计，就是为了让模型一层一层从局部看到整体。",
            "所以读这份课件时，不要先背很多名词。先记住一条线：**图像有空间结构 -> 卷积保留局部结构 -> 池化压缩信息 -> 更深层网络组合高级特征。**",
        ]
    if family == "mlp":
        return [
            "这一份课件要解决的是：**最基础的神经网络到底怎么从输入走到预测，再从错误走回参数更新。**",
            "多层感知机不是只背几层 Linear，它的重点是把“线性变换、激活函数、损失函数、反向传播、优化器”连成一条训练流程。",
            "读的时候先不要被代码吓到。你只需要先抓住：**前向传播得到预测，损失函数衡量错多少，反向传播算梯度，优化器改参数。**",
        ]
    if family == "rnn":
        return [
            "这一份课件要解决的是：**数据有先后顺序时，模型怎么一边往后读，一边保留前面的信息。**",
            "RNN、GRU、LSTM 都不是孤立结构，它们都是在处理同一个问题：普通网络不天然知道“谁在前，谁在后”。",
            "读的时候先抓这条线：**序列数据 -> 隐藏状态 -> 普通 RNN 的长期依赖问题 -> GRU/LSTM 用门控改善记忆。**",
        ]
    if family == "transformer":
        return [
            "这一份课件要解决的是：**如果不按 RNN 那样一步一步读，模型还能不能直接判断序列里不同位置之间的关系。**",
            "注意力机制的核心不是公式本身，而是“谁在查询、去和谁匹配、最后从哪里取信息”。Transformer 把这套机制堆成了完整结构。",
            "读的时候先抓这条线：**Q/K/V -> 注意力权重 -> 多头注意力 -> 位置编码与 Mask -> Encoder/Decoder 或 Decoder-only 结构。**",
        ]
    if family == "graph":
        return [
            "这一份课件要解决的是：**当数据本身是一张图时，模型怎样同时利用节点特征和节点之间的关系。**",
            "图学习最重要的不是某个模型名字，而是“对象之间有边，边决定信息怎么传”。GNN 的很多方法都可以看成邻居信息的聚合和更新。",
            "读的时候先抓这条线：**节点/边 -> 邻接矩阵 -> 图任务类型 -> 随机游走或表示学习 -> GCN/GAT 等图神经网络。**",
        ]
    return [
        "这一份课件会被拆成可以从头读的复习材料。重点不是堆名词，而是把每个知识点放回它出现的位置。",
        "读的时候先问三个问题：它为什么出现？它解决什么问题？考试会怎样把它和前后概念放在一起考？",
    ]


def narrative_bad_claim(row, family):
    term = row.get("term", "")
    claims = {
        "卷积": "卷积层会把所有像素完全展开后做全连接，因此主要靠增加参数量提升效果。",
        "互相关": "深度学习里的卷积必须手动翻转卷积核，否则无法学习局部特征。",
        "Valid": "Valid 卷积会大量填充边界，所以输出尺寸通常最大。",
        "Same": "Same 卷积的目的就是让输出尺寸尽量缩小。",
        "Full": "Full 卷积完全不填充，所以输出尺寸最小。",
        "感受野": "网络越深，高层神经元看到的原始图像范围一定越小。",
        "池化": "池化的主要作用是增加可学习参数量。",
        "卷积核": "卷积核只是固定模板，训练中不会更新。",
        "步长": "步长越大，输出特征图通常越大。",
        "填充": "填充只会删除边界信息，不能控制输出尺寸。",
        "激活函数": "没有激活函数时，多层线性网络也能自然表达复杂非线性关系。",
        "BatchNorm": "BatchNorm 推理时仍必须使用当前 batch 的统计量。",
        "Dropout": "Dropout 推理阶段仍要随机丢弃神经元。",
        "Batch": "Batch 表示序列长度，而不是一次送入模型的样本数。",
        "Epoch": "一个 epoch 就是一次参数更新。",
        "学习率": "学习率越大训练越稳定，通常不需要调整。",
        "损失函数": "训练时只看准确率即可，损失函数不参与反向传播。",
        "Sigmoid": "Sigmoid 的输出范围可以小于 0 或大于 1。",
        "Softmax": "Softmax 输出不需要总和为 1。",
        "梯度消失": "梯度消失就是模型在训练集上过拟合。",
        "梯度爆炸": "梯度爆炸最直接的处理方式是增加 Dropout。",
        "梯度裁剪": "梯度裁剪会把所有梯度直接清零。",
        "RNN": "RNN 的各个时间步之间完全独立。",
        "GRU": "GRU 和 LSTM 一样有单独的细胞状态 c_t。",
        "LSTM": "LSTM 没有门控结构，只是普通 RNN 的另一个名字。",
        "Padding": "Padding 补出来的位置都是真实有效的信息。",
        "Packing": "Packing 的目的就是增加 padding，让序列更长。",
        "PackedSequence": "PackedSequence 是普通张量，和变长序列处理无关。",
        "隐藏状态": "隐藏状态是固定模型参数，不随输入序列变化。",
        "双向RNN": "双向 RNN 只能利用过去信息，不能利用后文。",
        "堆叠RNN": "堆叠 RNN 是把时间步变多，不是增加层数。",
        "注意力机制": "注意力机制就是对所有位置做简单平均。",
        "Self-Attention": "Self-Attention 中 Q、K、V 必须来自三个完全不同的序列。",
        "Cross-Attention": "Cross-Attention 中 Q 和 K/V 必须来自同一个序列。",
        "Transformer": "Transformer 必须依赖 RNN 的递归结构处理序列。",
        "sqrt(d_k)": "sqrt(d_k) 是可学习参数，主要用于增加模型容量。",
        "图学习": "图学习只看单个样本特征，不需要关系结构。",
        "节点": "节点表示关系，边表示对象。",
        "边": "边表示节点自身特征，不表示节点之间关系。",
        "邻接矩阵": "邻接矩阵记录节点特征，不记录连接关系。",
        "度矩阵": "度矩阵把所有边权放在非对角位置。",
        "拉普拉斯矩阵": "图拉普拉斯矩阵就是邻接矩阵 A。",
        "GCN": "GCN 只看节点自身，不聚合邻居信息。",
        "GAT": "GAT 假设所有邻居同等重要，不学习注意力权重。",
        "消息传递": "消息传递不沿边发生，也不需要邻居聚合。",
    }
    if term in claims:
        return claims[term]
    if family == "cnn":
        return f"{term}在 CNN 中通常用于增加全连接参数量，而不是保留图像结构。"
    if family == "rnn":
        return f"{term}说明序列各时间步彼此独立，不需要保存历史信息。"
    if family == "transformer":
        return f"{term}可以完全脱离 Q/K/V、位置关系和 mask 单独理解。"
    if family == "graph":
        return f"{term}只处理单个节点特征，不需要考虑边或邻居。"
    return f"{term}只需要背名字，不需要知道它解决什么问题。"


def guide_section_story(chapter, rows):
    if not rows:
        return "这一节可抽取到的复习内容不多。先看标题，判断它和前后章节的关系，再回到课件原文补细节。"
    terms = [r["term"] for r in rows[:5]]
    if len(terms) == 1:
        return f"这一节先集中理解 **{terms[0]}**。不要一上来背定义，先问：课件为什么在这里讲它？它补上了前面哪一个问题？"
    return (
        f"这一节可以按 **{' -> '.join(terms)}** 的顺序读。"
        "第一个概念通常是入口，后面的概念要么是在补结构，要么是在解释训练时会遇到的问题。"
    )


def guide_term_story(row, idx, prev_term=None, next_term=None):
    term = row["term"]
    lines = [
        f"### {idx}. 先讲 {term}",
        "",
        f"**{term}** 的核心含义：{row['plain']}",
        "",
        f"它在课件中的考查位置：{row['exam']}",
        "",
    ]
    if prev_term or next_term:
        before = prev_term or "前面的概念"
        after = next_term or "后面的概念"
        lines.extend([
            f"如果按课件顺序串起来，可以这样看：**{before} -> {term} -> {after}**。",
            "也就是说，你不是在背一个词，而是在看知识点之间怎样往下推进。",
            "",
        ])
    lines.extend([
        f"考试作答时可以写成：**{row['memory']}**",
        "",
        f"这里最容易错的是：{row['pitfall']}",
        "",
    ])
    return lines


def guide_exam_check(rows, family):
    rows = rows[:4]
    if not rows:
        return ["这一节先把前面的讲解读顺，再回课件确认细节。"]
    lines = []
    row = rows[0]
    lines.extend([
        f"1. **选择题：关于 {row['term']}，下列说法错误的是：**",
        f"   A. {compact(row['plain'], 84)}",
        f"   B. {compact(row['memory'], 84)}",
        f"   C. {compact(row['exam'], 84)}",
        f"   D. {narrative_bad_claim(row, family)}",
        "",
        f"   **答案：D。** A、B、C 都能和课件里的 {row['term']} 对上；D 是把概念说反或放错位置。",
        "",
    ])
    if len(rows) >= 2:
        row = rows[1]
        lines.extend([
            f"2. **资料题：** 材料中出现 {row['term']}，请判断它在本节属于数据、模型结构、训练优化还是输出目标，并说明理由。",
            "",
            f"   **参考答案：** {compact(row['plain'], 120)} 答题时还要结合课件语境：{compact(row['exam'], 90)}",
            "",
        ])
    if len(rows) >= 3:
        row = rows[2]
        lines.extend([
            f"3. **简答题：为什么考试里不能只背 {row['term']} 的定义？**",
            "",
            f"   **参考答案：** 因为 {row['term']} 在课件中有明确的位置和作用：{compact(row['plain'], 100)} 答题时还要补出它通常怎么考：{compact(row['exam'], 100)} 最后要避开这个误区：{compact(row['pitfall'], 90)}",
            "",
        ])
    if len(rows) >= 4:
        a, b = rows[0], rows[3]
        lines.extend([
            f"4. **辨析题：请区分 {a['term']} 和 {b['term']}。**",
            "",
            f"   **参考答案：** {a['term']}：{compact(a['plain'], 90)}；{b['term']}：{compact(b['plain'], 90)}。答题时先判断它们分别处在数据、结构、训练还是输出环节。",
            "",
        ])
    return lines


def build_course_guide_markdown(course):
    title = course["title"]
    family = course.get("family", "general")
    lines = [
        f"# {title}：按课件顺序重写的考试复习讲义",
        "",
        "这份讲义会尽量按你读课件时真正需要的顺序来讲。",
        "",
        "**目标不是把课件拆成表格，而是像复习书一样，一点一点说明：为什么讲这个、它解决什么问题、考试怎么抓。**",
        "",
        "---",
        "",
        "## 先说整份课件到底在解决什么问题",
        "",
    ]
    for paragraph in narrative_course_problem(family):
        lines.append(paragraph)
        lines.append("")
    lines.extend([
        "---",
        "",
        "## 建议阅读顺序",
        "",
        "下面不是知识点清单，而是你复习时可以照着走的路线：",
        "",
    ])
    for idx, chapter in enumerate(course["chapters"], 1):
        rows = chapter.get("rows", [])
        terms = "、".join(row["term"] for row in rows[:4]) or "本节核心内容"
        lines.append(f"{idx}. **{chapter['title']}**：先读 {terms}。")
    lines.extend(["", "---", "", "## 开始按课件一点一点讲", ""])

    for chapter_idx, chapter in enumerate(course["chapters"], 1):
        rows = chapter.get("rows", [])[:10]
        lines.extend([
            f"## 第 {chapter_idx} 部分：{chapter['title']}",
            "",
            "### 这一节先不要急着背",
            "",
            guide_section_story(chapter, rows),
            "",
        ])
        for idx, row in enumerate(rows, 1):
            prev_term = rows[idx - 2]["term"] if idx > 1 else None
            next_term = rows[idx]["term"] if idx < len(rows) else None
            lines.extend(guide_term_story(row, idx, prev_term, next_term))
        lines.extend([
            "### 这一节最后怎么记",
            "",
        ])
        for row in rows[:6]:
            lines.append(f"- **{row['term']}**：{row['memory']}")
        lines.extend(["", "### 本节检查题", ""])
        lines.extend(guide_exam_check(rows, family))
        lines.append("---")
        lines.append("")

    lines.extend([
        "## 考前总复盘",
        "",
        "考前不要平均用力。你可以按这四个问题检查自己：",
        "",
        "1. 这个概念为什么会在课件这里出现？",
        "2. 它解决的是输入、模型结构、训练优化、输出损失，还是图/序列关系的问题？",
        "3. 如果让我写简答，我能不能不用空话，写出一句准确解释？",
        "4. 如果让我做选择题，我能不能看出选项是不是把概念放错位置、说反作用、混淆训练和推理？",
        "",
        "## 打印建议",
        "",
        "<style>@media print { @page { margin: 8mm; } body { font-size: 11pt; line-height: 1.35; } h1, h2, h3 { page-break-after: avoid; } }</style>",
        "",
    ])
    return "\n".join(lines)


def build_study_guides(outline):
    GUIDE_OUT.mkdir(parents=True, exist_ok=True)
    for old in GUIDE_OUT.glob("*.md"):
        old.unlink()

    guides = []
    index_lines = [
        "# 深度学习期末考试讲义总目录",
        "",
        "这份总目录对应网站当前纳入的全部课件。建议按“多层感知机 -> CNN -> RNN -> Attention/Transformer -> 图学习”的顺序复习。",
        "",
        "## 课件讲义列表",
        "",
    ]
    for course in outline["courses"]:
        slug = safe_slug(course["title"])
        rel = f"guides/{slug}-考试复习讲义.md"
        md = build_course_guide_markdown(course)
        (SITE / rel).write_text(md, encoding="utf-8")
        item = {
            "title": course["title"],
            "file": course["file"],
            "kind": course.get("kind"),
            "href": rel,
            "chapter_count": course.get("chapter_count", 0),
            "chunk_count": course.get("chunk_count", 0),
        }
        guides.append(item)
        index_lines.append(f"- [{course['title']}]({Path(rel).name})：{course.get('chapter_count', 0)} 个章节块。")
    index_lines.extend([
        "",
        "## 使用顺序",
        "",
        "1. 先读每份讲义的“整份课件到底在解决什么问题”。",
        "2. 再按讲义顺序一节一节读，不要跳到表格里背词。",
        "3. 最后做网站里的基础练习和标准组卷。",
        "",
        "<style>@media print { @page { margin: 8mm; } body { font-size: 11pt; line-height: 1.35; } }</style>",
    ])
    index_rel = "guides/00-深度学习期末考试讲义总目录.md"
    (SITE / index_rel).write_text("\n".join(index_lines) + "\n", encoding="utf-8")
    return {"index": index_rel, "items": guides}


def q_mc(id_, topic, stem, options, answer, explanation, source, difficulty="易"):
    return {"id": id_, "type": "mcq", "topic": topic, "stem": stem, "options": options, "answer": answer, "explanation": explanation, "source": source, "difficulty": difficulty}


def q_short(id_, topic, stem, answer, explanation, source, difficulty="易"):
    return {"id": id_, "type": "short", "topic": topic, "stem": stem, "answer": answer, "explanation": explanation, "source": source, "difficulty": difficulty}


def q_material(id_, topic, material, questions, answer, explanation, source, difficulty="中"):
    return {
        "id": id_,
        "type": "material",
        "topic": topic,
        "stem": f"根据下面关于“{topic}”的场景材料回答问题。",
        "material": material,
        "subquestions": questions,
        "answer": answer,
        "explanation": explanation,
        "source": source,
        "difficulty": difficulty,
    }


def has_bad_question_text(q):
    text = " ".join(str(q.get(k, "")) for k in ["stem", "material", "subquestions", "questions", "answer", "explanation"])
    return any(p in text for p in BAD_QUESTION_PATTERNS)


EXAM_TOPIC_BLUEPRINTS = [
    {
        "key": "neuron",
        "tier": "重点",
        "topic": "神经元与感知机",
        "source": "深度学习导学_20260609.md",
        "fact": "TLU/感知机把输入加权求和后与阈值或偏置共同决定输出，本质是线性可分边界。",
        "trap": "单层感知机不能直接解决 XOR 这类非线性可分问题，多层结构或非线性变换才是关键。",
        "wrong": [
            "单层感知机可以不加隐藏层直接表达任意 XOR 逻辑",
            "阈值逻辑单元不需要权重，输入只要相加即可完成所有分类",
            "ADALINE 与感知机完全相同，训练时都只看二值阶跃输出",
        ],
        "short": "说明 TLU、感知机和 ADALINE 之间的递进关系。",
        "answer": "TLU 强调用加权求和和阈值产生二值输出；感知机把它用于线性分类，决策边界是超平面；ADALINE 进一步用连续线性输出和误差最小化引出损失函数、梯度下降和参数更新。它们共同铺垫了后续神经网络训练流程。",
    },
    {
        "key": "training",
        "tier": "重点",
        "topic": "训练流程与自动微分",
        "source": "多层感知机1_学生分发版.ipynb",
        "fact": "训练闭环是前向传播得到预测、损失函数衡量误差、反向传播计算梯度、优化器更新参数。",
        "trap": "清零梯度、训练/评估模式、no_grad 与 scheduler 的位置经常被混淆。",
        "wrong": [
            "反向传播发生在优化器更新参数之后",
            "计算图只保存输入数据，不能用于自动求导",
            "训练时不需要清零梯度，PyTorch 会自动丢弃旧梯度",
        ],
        "short": "按顺序写出 PyTorch 中一个 batch 的标准训练步骤，并说明每一步作用。",
        "answer": "标准顺序是 optimizer.zero_grad() 清旧梯度，model(x) 前向得到预测，loss_fn 计算损失，loss.backward() 沿计算图求梯度，optimizer.step() 更新参数。验证时要 model.eval()，通常配合 torch.no_grad()，避免梯度记录和训练态层产生干扰。",
    },
    {
        "key": "mlp",
        "tier": "重点",
        "topic": "MLP 与分类损失",
        "source": "多层感知机2-学生分发版-终版2026.ipynb",
        "fact": "MLP 由仿射变换和非线性激活堆叠而成，分类任务要让输出层、标签格式和损失函数匹配。",
        "trap": "没有激活函数时，多层线性层仍可合并为一个线性变换；CrossEntropyLoss 应输入 logits。",
        "wrong": [
            "只要全连接层足够多，即使没有激活函数也能拟合任意非线性边界",
            "CrossEntropyLoss 前必须先手动 Softmax",
            "BCEWithLogitsLoss 的输入应该是 Sigmoid 后的概率",
        ],
        "short": "解释为什么 MLP 必须依赖激活函数，并比较二分类和多分类常用损失。",
        "answer": "仿射变换只做线性组合和平移，多层线性结构没有激活函数仍等价于一个线性变换，无法表达复杂非线性边界。二分类或多标签常用 BCEWithLogitsLoss，输入 logits；单标签多分类常用 CrossEntropyLoss，输入每类 logits，内部完成 LogSoftmax/NLLLoss 类计算。",
    },
    {
        "key": "cnn",
        "tier": "重点",
        "topic": "CNN 卷积、池化与架构",
        "source": "卷积神经网络1-学生分发版-终版2026.ipynb",
        "fact": "CNN 用局部连接、参数共享、卷积/互相关、池化和层级特征提取处理图像。",
        "trap": "输出尺寸由输入、卷积核、步长、填充共同决定；CNN 实践里常做互相关而不是翻转卷积核。",
        "wrong": [
            "CNN 为了保留图像结构必须把所有像素全连接到每个神经元",
            "卷积输出大小与卷积核大小、步长和填充无关",
            "池化层的主要作用是增加大量可学习参数",
        ],
        "short": "说明 CNN 相比 MLP 处理图像的优势，并写出卷积输出尺寸公式。",
        "answer": "MLP 展平图像会破坏空间结构且参数量巨大；CNN 通过局部连接保留邻域关系，通过参数共享减少参数，通过池化压缩空间尺寸并增强一定平移不变性。常用输出尺寸公式为 (输入尺寸 + 2P - K) / S + 1。",
    },
    {
        "key": "cnn_train",
        "tier": "重点",
        "topic": "CNN 训练技巧",
        "source": "卷积神经网络2-学生分发版-终版2026.ipynb",
        "fact": "BatchNorm、Dropout、权重衰减、数据增强、学习率调度、早停用于处理训练不稳和过拟合。",
        "trap": "不同技巧解决的问题不同：梯度裁剪偏向梯度爆炸，Dropout/权重衰减偏向过拟合，BatchNorm 偏向稳定训练。",
        "wrong": [
            "Dropout 在推理阶段仍应随机丢弃神经元来保持训练一致",
            "BatchNorm 推理时一定使用当前 batch 的均值方差",
            "梯度裁剪主要用于让梯度消失得更快",
        ],
        "short": "比较 BatchNorm、Dropout、权重衰减和数据增强分别解决什么问题。",
        "answer": "BatchNorm 标准化中间激活，使训练更稳定；Dropout 训练时随机失活，降低对少数神经元的依赖；权重衰减惩罚过大参数，缓解过拟合；数据增强扩充训练分布，让模型见到合理变化。答题时要写出训练/推理阶段差异。",
    },
    {
        "key": "rnn",
        "tier": "重点",
        "topic": "RNN、GRU、LSTM",
        "source": "循环神经网络-学生分发版-终版2026.ipynb",
        "fact": "RNN 通过隐藏状态沿时间传递信息，GRU/LSTM 用门控结构缓解长序列依赖和梯度问题。",
        "trap": "LSTM 有细胞状态和输入/遗忘/输出门；GRU 通常只有更新门和重置门，没有独立细胞状态。",
        "wrong": [
            "RNN 每个时间步完全独立，不需要隐藏状态",
            "LSTM 没有门控结构，只是普通线性层堆叠",
            "Padding 补出来的位置都是真实有效时间步，不需要 mask 或 packing",
        ],
        "short": "解释 RNN 的隐藏状态作用，并比较 GRU 与 LSTM 的门控差异。",
        "answer": "隐藏状态把过去时间步的信息带到当前时间步，使模型能处理序列。GRU 用更新门和重置门控制保留和重置历史信息，结构较简洁；LSTM 有输入门、遗忘门、输出门和细胞状态，长期记忆通道更明确。变长序列要注意 padding、packing 或 mask。",
    },
    {
        "key": "transformer",
        "tier": "重点",
        "topic": "Attention 与 Transformer",
        "source": "注意力机制及Transformer1-学生分发版-终版2026.ipynb",
        "fact": "注意力用 Q/K 匹配得到权重，再对 V 加权求和；Transformer 由多头注意力、前馈网络、残差、归一化和位置编码组成。",
        "trap": "sqrt(d_k) 是缩放项，不是可学习参数；Target Mask 用于防止解码器看到未来。",
        "wrong": [
            "Value 用来和 Query 点积计算匹配分数，Key 不参与匹配",
            "Transformer 不需要任何位置信息也能天然知道词序",
            "Target Mask 的作用是屏蔽源序列 padding，而不是防止看未来",
        ],
        "short": "写出缩放点积注意力的计算流程，并说明多头注意力和位置编码的作用。",
        "answer": "流程是 Q 与 K 转置点积得到相关性分数，除以 sqrt(d_k) 稳定 softmax，再用 softmax 权重对 V 加权求和。多头注意力让模型在不同子空间学习不同关系；位置编码补充序列顺序，因为自注意力本身不含天然位置信息。",
    },
    {
        "key": "llm",
        "tier": "次重点",
        "topic": "大语言模型与 Prompt",
        "source": "图学习课件_第二部分_20260602.pptx",
        "fact": "大语言模型以 Transformer 为核心，通过大规模预训练获得语言生成和理解能力，Prompt 用自然语言引导模型完成任务。",
        "trap": "LLM 不是重新定义 Transformer 基础结构，考试更可能考演进脉络、Prompt 作用、MoE/DeepSeek R1 等概念层理解。",
        "wrong": [
            "Prompt 只用于装饰输入文字，不影响模型输出方向",
            "大语言模型完全不依赖 Transformer 或注意力机制",
            "MoE 的核心是让所有专家每次都完整参与计算",
        ],
        "short": "说明 Prompt、MoE 和推理型模型在大语言模型中的基本作用。",
        "answer": "Prompt 是任务指令和上下文，会影响模型生成方向；MoE 用路由选择部分专家参与计算，在参数规模和计算成本之间折中；推理型模型强调更长的思考链路或强化学习后训练，让模型在复杂问题上生成更可靠的中间推理。",
    },
    {
        "key": "generative",
        "tier": "次重点",
        "topic": "生成式算法",
        "source": "图学习课件_第二部分_20260602.pptx",
        "fact": "生成式算法关注学习数据分布并产生新样本，常见思想包括自回归生成、扩散式逐步去噪和表示空间生成。",
        "trap": "生成式模型不是只做分类；它要建模数据如何被生成或如何从噪声恢复。",
        "wrong": [
            "生成式算法的目标只是给已有样本贴标签",
            "自回归生成可以在生成当前 token 时直接使用未来 token",
            "扩散模型不需要学习从噪声恢复数据的过程",
        ],
        "short": "用分类模型和生成式模型的目标差异解释生成式算法。",
        "answer": "分类模型主要学习输入到标签的映射，回答“这是什么”；生成式模型要学习数据分布或生成过程，回答“怎样生成像这样的数据”。自回归模型逐步预测下一个 token，扩散模型学习逐步去噪，二者都服务于生成新内容。",
    },
    {
        "key": "graph",
        "tier": "次重点",
        "topic": "图学习与图神经网络",
        "source": "图神经网络课件_20260609.pptx",
        "fact": "图神经网络在节点、边和邻接关系上做消息传递，通过邻居聚合更新节点表示。",
        "trap": "图卷积不能简单理解成固定方形卷积核在规则图像网格上滑动，图的邻域由边决定。",
        "wrong": [
            "图学习只看节点自身特征，不需要边或邻接矩阵",
            "GCN 的邻居聚合与图结构无关",
            "GAT 和 GCN 完全一样，都不能区分不同邻居重要性",
        ],
        "short": "说明 GCN、GAT、GraphSAGE 的共同点和差异。",
        "answer": "共同点是都利用邻居信息更新节点表示，属于消息传递/邻居聚合思想。GCN 常用归一化邻接矩阵聚合邻居；GAT 学习不同邻居的注意力权重；GraphSAGE 通过邻居采样和聚合支持较大图或归纳式节点表示学习。",
    },
    {
        "key": "pinn",
        "tier": "次重点",
        "topic": "PINN",
        "source": "图神经网络课件_20260609.pptx",
        "fact": "PINN 把物理方程约束加入神经网络训练，使模型不仅拟合数据，也尽量满足物理规律。",
        "trap": "PINN 的关键不是换一种激活函数，而是把方程残差、边界条件或初始条件写进损失。",
        "wrong": [
            "PINN 完全不需要物理方程，只要数据越多越好",
            "PINN 的物理约束只在测试阶段检查，不参与训练损失",
            "PINN 与普通 MLP 的差别只是网络层数更深",
        ],
        "short": "解释 PINN 中数据损失和物理损失分别起什么作用。",
        "answer": "数据损失让模型拟合观测样本，物理损失把微分方程残差、边界条件或初始条件纳入优化目标。这样模型在数据较少时也能利用先验规律，输出更符合物理约束的解。",
    },
    {
        "key": "materials",
        "tier": "次重点",
        "topic": "AI for material science",
        "source": "图神经网络课件_20260609.pptx",
        "fact": "材料科学中的 AI 常把分子、晶体或结构关系表示为图，再用 GNN/PINN 等模型预测性质或辅助设计。",
        "trap": "材料科学不是孤立的应用名词，关键是把结构、关系、物理约束和预测目标联系起来。",
        "wrong": [
            "材料科学应用中图结构没有意义，只需要普通表格特征",
            "GNN 不能处理原子之间的连接关系",
            "材料性质预测不需要考虑结构或物理先验",
        ],
        "short": "说明为什么材料科学问题适合和图学习或物理约束模型结合。",
        "answer": "材料由原子、键、晶格和空间关系构成，天然可以表示为图；GNN 能聚合局部邻域和结构信息，预测性质或筛选材料。若任务涉及物理规律，PINN 或物理约束损失可以补充纯数据学习的不足。",
    },
]


def q_options(correct, wrongs):
    opts = [correct] + wrongs[:3]
    return [f"{label}. {text}" for label, text in zip("ABCD", opts)]


def make_exam_banks():
    return build_curated_quiz_banks()



def clean_quiz(data):
    cleaned = []
    seen = set()
    for q in data:
        if has_bad_question_text(q):
            continue
        qid = q.get("id")
        if not qid or qid in seen:
            continue
        cleaned.append(q)
        seen.add(qid)
    return cleaned


def extend_quiz_from_outline(standard, beginner, outline):
    return standard, beginner



def fallback_review(outline):
    return {
        "title": "深度学习课件综合复习资料",
        "updated_from": "当前课件目录内全部 Jupyter/PPT 课件",
        "sections": [
            {
                "title": course["title"],
                "points": [
                    f"本课件包含 {course['chapter_count']} 个章节块、{course['chunk_count']} 个可检索片段。",
                    "建议先看本页教材目录表格，再进入主题书库按考试重点复习，最后做基础练习和标准组卷。",
                ],
                "exam_focus": [row["term"] for ch in course["chapters"] for row in ch["rows"][:3]][:12],
            }
            for course in outline["courses"]
        ],
    }


def load_review_or_fallback(outline):
    path = DATABASE / "review_material.json"
    if path.exists():
        try:
            data = json.loads(path.read_text(encoding="utf-8"))
            data["updated_from"] = "当前课件目录内全部 Jupyter/PPT 课件"
            return data
        except json.JSONDecodeError:
            pass
    return fallback_review(outline)


def rows_for_focus(outline, families=None, keywords=None, limit=12):
    families = set(families or [])
    keywords = keywords or []
    rows = []
    for course in outline.get("courses", []):
        if families and course.get("family") not in families:
            pass_family = False
        else:
            pass_family = True
        for chapter in course.get("chapters", []):
            for row in chapter.get("rows", []):
                text = " ".join([course.get("title", ""), chapter.get("title", ""), row.get("term", ""), row.get("plain", ""), row.get("memory", "")])
                if pass_family or any(k.lower() in text.lower() for k in keywords):
                    rows.append((course, chapter, row))
                if len(rows) >= limit:
                    return rows
    return rows[:limit]


def manual_foundation_rows():
    rows = [
        {
            "term": "TLU 阈值逻辑单元",
            "plain": "把多个输入加权求和，再和阈值比较，输出 0 或 1，是早期神经元模型的抽象形式。",
            "exam": "常考权重、阈值、阶跃函数分别起什么作用，以及它为什么只能形成线性边界。",
            "pitfall": "不要把 TLU 说成能直接表达所有非线性逻辑；XOR 是典型反例。",
            "memory": "TLU = 加权求和 + 阈值判断 + 二值输出。",
        },
        {
            "term": "逻辑门",
            "plain": "AND、OR、NOT 可以用不同权重和阈值的 TLU 表示，说明简单神经元能模拟基本布尔逻辑。",
            "exam": "常考给定权重和阈值判断它对应哪个逻辑门，或解释为什么 XOR 需要多层组合。",
            "pitfall": "不要把能表示 AND/OR 推成能单层表示任意逻辑。",
            "memory": "基本逻辑门可由 TLU 实现，复杂逻辑通常需要组合。",
        },
        {
            "term": "构造性合成",
            "plain": "先用第一层识别若干输入模式，再用后一层把模式合并，从而构造更复杂的布尔函数。",
            "exam": "常考多层结构为什么比单层 TLU 表达能力更强。",
            "pitfall": "不要只说层数变多；关键是隐藏层能表示中间模式。",
            "memory": "多层网络可以先识别局部模式，再组合成最终输出。",
        },
        {
            "term": "感知机",
            "plain": "感知机是带权重和偏置的线性分类器，通过符号判断把输入分到超平面两侧。",
            "exam": "常考线性可分、决策边界、XOR 局限。",
            "pitfall": "单层感知机只能处理线性可分问题。",
            "memory": "感知机决策边界是线性的。",
        },
        {
            "term": "ADALINE",
            "plain": "ADALINE 使用连续线性输出计算误差，更自然地引出损失函数、梯度下降和参数更新。",
            "exam": "常考它和感知机的区别：训练时关注连续误差，而不是只看阶跃后的类别。",
            "pitfall": "不要把 ADALINE 完全等同于感知机的二值输出学习。",
            "memory": "ADALINE 把神经元学习和误差最小化连接起来。",
        },
    ]
    return [(None, None, row) for row in rows]


def add_focus_section(lines, title, why, rows, exam_hint):
    lines.extend([f"## {title}", "", why, ""])
    if rows:
        lines.extend(["### 先抓主线", ""])
        for _, _, row in rows[:5]:
            lines.append(f"- **{row['term']}**：{row['plain']}")
        lines.append("")
        lines.extend(["### 考试怎么写", ""])
        for _, _, row in rows[:5]:
            lines.append(f"- {row['term']}：{row['exam']} 易错点：{row['pitfall']}")
        lines.append("")
        lines.extend(["### 必背句", ""])
        for _, _, row in rows[:5]:
            lines.append(f"- {row['memory']}")
        lines.append("")
    lines.extend(["### 题型提醒", "", exam_hint, ""])


def build_final_review_markdown(outline, manifest):
    source_names = [m["file"] for m in manifest]
    lines = [
        "# 深度学习期末复习（仅基于原始课件版）",
        "",
        "> 本文件依据当前课件根目录中的原始 Jupyter Notebook、三份 PPT OCR 文本和导学 Markdown 重建。旧版详细讲义只用于恢复讲解密度和写作节奏，不作为知识来源。",
        "",
        "## 资料来源说明",
        "",
    ]
    for name in source_names:
        lines.append(f"- {name}")
    body = dedent(r"""

    ## 怎么使用这份总复习

    这份资料不是把每个课件标题重新排一遍，而是把所有课件按期末考试需要重新连成一条学习线。老师已经明确说明：**神经元、MLP、CNN、RNN、Transformer 是重点**；大语言模型、生成式算法、图神经网络、PINN、AI for material science 相对偏弱，但仍然可能出现在单选题、简答题或资料题里。

    复习时建议按三轮走：

    1. **第一轮看主线**：先知道每个模型为什么出现，它解决了前一个模型的什么问题。
    2. **第二轮补细节**：公式、输入输出形状、训练流程、损失函数匹配、常见工程写法。
    3. **第三轮按题型练**：单选题抓陷阱，简答题写结构化答案，资料题先读情境再套知识。

    不要把深度学习复习成“背名词”。考试更容易考你能不能把一个概念放回流程里：数据怎样进来，模型怎样算，损失怎样衡量错，梯度怎样传，参数怎样改，最后输出怎样解释。

    ## 老师给出的考试权重

    | 层级 | 内容 | 复习要求 |
    |---|---|---|
    | 第一梯队 | 神经元、MLP、CNN、RNN、Transformer | 必须会解释、会判断、会写简答、会读资料题 |
    | 第二梯队 | LLM、生成式算法、GNN、PINN、AI for material science | 抓概念主线、典型结构、应用场景和常见判断 |
    | 题型 | 单选题、简答题、资料题 | 复习要按“判断 + 表达 + 场景分析”准备 |

    ## 全局学习路线

    深度学习这门课可以理解成一条从“简单神经元”走向“复杂结构数据建模”的路线：

    1. **神经元和感知机**：先理解一个单元怎样把输入加权求和，再通过阈值或激活函数输出结果。
    2. **训练流程和 MLP**：再理解模型不是手写规则，而是通过损失函数、反向传播和优化器调参数。
    3. **CNN**：当输入是图像，不能只把像素摊平，要利用局部结构、参数共享和空间层级特征。
    4. **RNN/GRU/LSTM**：当输入是序列，样本之间有先后关系，需要隐藏状态和门控保存历史信息。
    5. **Attention/Transformer**：当长序列中远距离依赖很重要，用注意力直接计算不同位置之间的相关性。
    6. **NLP/LLM/生成式模型**：语言任务从规则系统、词向量、序列模型走向大规模预训练和生成式能力。
    7. **图学习/GNN/PINN/材料科学**：当数据结构变成图、物理系统或科学问题，需要把关系结构、物理先验和领域约束放进模型。

    ---

    # 第一部分：神经元、TLU、感知机、ADALINE

    ## 1. 为什么课程要从神经元讲起

    深度学习看起来复杂，但它最早的抽象很简单：一个单元接收多个输入，每个输入有一个权重，把它们加起来，再决定是否输出。这个思想来自神经元模型，也就是“输入刺激累积到一定程度就触发输出”的直觉。

    课件导学从 **Threshold Logic Unit（TLU，阈值逻辑单元）** 开始，就是为了说明：神经网络的基本单元并不是黑箱，它最核心的动作是“加权求和 + 阈值判断”。

    ## 2. TLU：加权求和后做二值判断

    TLU 的抽象形式可以写成：

    $$
    y = step\left(\sum_i w_i x_i - \theta\right)
    $$

    其中：

    - $x_i$ 是输入，可以理解为第 $i$ 个特征或二值条件。
    - $w_i$ 是权重，表示这个输入对输出的影响强弱和方向。
    - $\theta$ 是阈值，表示累计信号要超过多少才触发输出。
    - $step(\cdot)$ 是阶跃函数，输入达到条件输出 1，否则输出 0。

    用更通俗的话说，TLU 就像一个“投票器”：每个输入带着权重投票，总票数超过门槛就输出 1，没有超过就输出 0。

    考试里如果给你一组权重和阈值，让你判断输出，不要慌。只要按三步做：

    1. 把输入和权重对应相乘。
    2. 求和并减去阈值。
    3. 看阶跃函数输出 0 还是 1。

    ## 3. 用 TLU 表示逻辑门

    导学课件里讲了 TLU 可以实现 AND、OR、NOT 等基本逻辑门。这一节不是为了让你背代码，而是为了让你看到：**神经元可以通过权重和阈值模拟逻辑规则**。

    | 逻辑门 | 直觉 | TLU 思路 |
    |---|---|---|
    | AND | 两个条件都满足才输出 1 | 两个输入都为 1 时加权和才过阈值 |
    | OR | 任意一个条件满足就输出 1 | 只要一个输入为 1 就能过阈值 |
    | NOT | 输入为 0 时输出 1 | 用负权重或反向条件实现 |

    这里最重要的考试点不是 AND、OR 的某一组具体参数，而是：**简单逻辑门可以由一个 TLU 表示，但更复杂的非线性逻辑往往需要多层结构。**

    ## 4. XOR 为什么重要

    XOR 经常作为感知机局限性的例子。XOR 的输出规则是：两个输入不同则输出 1，相同则输出 0。它的问题在于，正类点和负类点无法用一条直线分开。

    单层感知机本质上只能学一个线性边界，所以它不能直接解决 XOR。多层结构可以先让隐藏层识别两个局部模式：

    - 模式 1：$x_1=1, x_2=0$
    - 模式 2：$x_1=0, x_2=1$

    然后输出层把这两个模式合并。这样就体现了多层网络的意义：**隐藏层先构造中间特征，输出层再组合这些特征。**

    ## 5. 感知机：线性分类器的起点

    感知机可以写成：

    $$
    y =
    \begin{cases}
    1, & \mathbf{w}\cdot\mathbf{x}+b > 0 \\
    0, & \mathbf{w}\cdot\mathbf{x}+b \le 0
    \end{cases}
    $$

    这里的 $\mathbf{w}\cdot\mathbf{x}+b$ 是一个线性函数。它在二维空间里对应一条直线，在高维空间里对应一个超平面。感知机判断样本属于哪一类，本质上就是看样本落在这个边界的哪一侧。

    常见错误说法：

    - “感知机可以解决任意分类问题。”错误。单层感知机只能处理线性可分问题。
    - “只要训练足够久，单层感知机就能解决 XOR。”错误。模型表达能力不够，训练再久也不能改变结构限制。
    - “偏置 $b$ 没有实际意义。”错误。偏置控制决策边界平移。

    ## 6. ADALINE：从二值判断走向误差最小化

    ADALINE 和感知机的关键区别是：它在训练时关注连续输出和真实值之间的误差，而不是只看阶跃后的类别。这样就自然引出了损失函数和梯度下降。

    如果模型输出是：

    $$
    \hat{y} = \mathbf{w}\cdot\mathbf{x}+b
    $$

    可以用均方误差衡量预测和真实值的差距：

    $$
    L = \frac{1}{N}\sum_{i=1}^{N}(\hat{y}_i-y_i)^2
    $$

    这个公式的意义很朴素：

    - $\hat{y}_i$ 是第 $i$ 个样本的预测。
    - $y_i$ 是真实答案。
    - $\hat{y}_i-y_i$ 是误差。
    - 平方让正负误差都变成正数，也会放大大错误。
    - 对所有样本求平均，得到整体错误程度。

    这一节和后面的 MLP 是连着的：ADALINE 让你第一次看到“模型输出 -> 损失函数 -> 参数更新”的训练思想。

    ## 7. 本章考试写法

    简答题如果问“从 TLU 到感知机再到 ADALINE 的关系”，可以这样答：

    > TLU 强调加权求和和阈值输出，能表示基本逻辑门；感知机把这种思想扩展成线性分类器，用权重和偏置形成决策边界；ADALINE 在训练时使用连续输出计算误差，使神经元学习和损失最小化、梯度下降联系起来。三者的递进是从逻辑判断到线性分类，再到可优化的误差学习。

    ---

    # 第二部分：训练流程与 MLP

    ## 1. MLP 这一部分到底在解决什么问题

    多层感知机（MLP）是后续 CNN、RNN、Transformer 的训练基础。你要在这一章建立一个稳定的训练闭环：

    **数据进入模型 -> 前向传播得到预测 -> 损失函数衡量错误 -> 反向传播计算梯度 -> 优化器更新参数 -> 重复很多 batch 和 epoch。**

    课件从一条直线的训练讲起，再过渡到 PyTorch 张量、自动微分、优化器、`nn.Module`、`Dataset`、`DataLoader`，最后进入二分类和多分类。这条路线很重要，因为考试资料题经常给一段代码，让你判断训练流程少了哪一步。

    ## 2. 从最简单的线性模型开始

    课件中的入门模型是：

    $$
    y = b + wx
    $$

    这里：

    - $x$ 是输入特征。
    - $w$ 是权重，控制输入对输出影响多大。
    - $b$ 是偏置，控制整体平移。
    - $y$ 是真实值，$\hat{y}$ 是模型预测值。

    训练的目标不是“写出正确答案”，而是让模型自己把 $w$ 和 $b$ 调到合适位置。比如课件里用合成数据构造近似关系 `y = 1 + 2x + 噪声`，模型如果学到 $b \approx 1$、$w \approx 2$，说明训练有效。

    ## 3. 损失函数：把“错了多少”变成可优化数字

    对回归任务，常用均方误差：

    $$
    MSE = \frac{1}{N}\sum_{i=1}^{N}(\hat{y}_i-y_i)^2
    $$

    这个损失越小，说明预测越接近真实值。深度学习训练不是直接最大化准确率，而是通过一个可微的损失函数来调整参数。

    这里要区分两个概念：

    | 名称 | 作用 | 训练中是否直接求梯度 |
    |---|---|---|
    | 损失函数 | 告诉模型错得多严重 | 是 |
    | 准确率/F1 等指标 | 评价模型表现 | 通常不是训练的直接目标 |

    单选题常把“准确率就是损失函数”写成错误选项。准确率可以用于评估，但很多情况下不可微，不能直接作为反向传播的训练损失。

    ## 4. 梯度下降：参数怎样变好

    梯度告诉我们：如果某个参数稍微变大，损失会怎样变。梯度下降的基本更新形式是：

    $$
    \theta \leftarrow \theta - \eta \nabla_\theta L
    $$

    其中：

    - $\theta$ 表示要学习的参数，可以是 $w$、$b$，也可以是神经网络里所有权重。
    - $\eta$ 是学习率，控制每次走多大一步。
    - $\nabla_\theta L$ 是损失函数对参数的梯度。

    方向前面的减号很关键：梯度指向损失上升最快的方向，所以要沿相反方向更新，损失才会下降。

    学习率太大，可能来回震荡甚至发散；学习率太小，训练会非常慢。后面优化器和学习率调度器都是围绕这个问题改进。

    ## 5. PyTorch 训练循环

    课件中 PyTorch 训练流程可以浓缩成下面的骨架：

    ```python
    model.train()
    for x_batch, y_batch in train_loader:
        y_pred = model(x_batch)
        loss = criterion(y_pred, y_batch)

        optimizer.zero_grad()
        loss.backward()
        optimizer.step()
    ```

    每一行对应的概念：

    - `model.train()`：让模型进入训练模式，Dropout、BatchNorm 等层会按训练逻辑工作。
    - `model(x_batch)`：前向传播。
    - `criterion(...)`：计算损失。
    - `optimizer.zero_grad()`：清空上一轮累积的梯度。
    - `loss.backward()`：反向传播，计算梯度。
    - `optimizer.step()`：优化器根据梯度更新参数。

    资料题常考顺序。最常见错误是忘记 `zero_grad()`，导致梯度跨 batch 累积；或者以为 `backward()` 已经更新了参数。实际上，`backward()` 只计算梯度，真正改参数的是 `optimizer.step()`。

    ## 6. Tensor、Batch、Epoch、Iteration

    这些词不难，但考试喜欢混淆：

    | 概念 | 通俗解释 |
    |---|---|
    | Tensor | 深度学习中的多维数组，数据和参数都以张量形式参与计算 |
    | Batch | 一次送入模型的一小批样本 |
    | Iteration | 一个 batch 完成一次前向、损失、反向、更新 |
    | Epoch | 整个训练集完整看过一遍 |

    举例：训练集有 1000 个样本，batch size 为 100，那么一个 epoch 有 10 个 iteration。

    ## 7. Dataset 与 DataLoader 的作用

    `Dataset` 负责定义“第 $i$ 个样本怎么取”，`DataLoader` 负责把样本组成 batch，并可以打乱顺序、多进程读取。它们解决的是数据进入模型之前的组织问题。

    资料题如果给你一段自定义 `Dataset`，要看清楚：

    - `__len__` 是否返回样本数量。
    - `__getitem__` 是否返回一个输入和一个标签。
    - 输入张量形状是否符合模型期望。
    - 标签类型是否符合损失函数要求。

    ## 8. MLP 的基本层：仿射变换 + 激活函数

    一层全连接层本质是：

    $$
    \mathbf{z} = \mathbf{x}W^T + \mathbf{b}
    $$

    然后接激活函数：

    $$
    \mathbf{h} = \sigma(\mathbf{z})
    $$

    这里的 $\sigma$ 可以是 ReLU、Sigmoid、Tanh 等。没有激活函数，多层线性层可以合并成一层线性层：

    $$
    W_2(W_1x+b_1)+b_2 = (W_2W_1)x + (W_2b_1+b_2)
    $$

    所以考试中“只堆叠线性层就能表达复杂非线性边界”是错误的。MLP 之所以比线性模型更强，是因为线性变换之间加入了非线性激活。

    ## 9. 二分类：Sigmoid 与 BCE

    二分类常用 Sigmoid 把实数压到 0 到 1：

    $$
    \sigma(z)=\frac{1}{1+e^{-z}}
    $$

    输出可以理解为样本属于正类的概率。对应的二元交叉熵可以写成：

    $$
    BCE = -\left[y\log(p)+(1-y)\log(1-p)\right]
    $$

    实践中常用 `BCEWithLogitsLoss`，它把 Sigmoid 和 BCE 合在一起，数值更稳定。此时模型输出应该是 logits，不要提前手动 `sigmoid` 后再喂给它。

    ## 10. 多分类：Softmax 与 CrossEntropyLoss

    多分类模型通常输出每个类别的 logits：

    $$
    \mathbf{z} = [z_1,z_2,\ldots,z_C]
    $$

    Softmax 把 logits 转成概率分布：

    $$
    p_i = \frac{e^{z_i}}{\sum_{j=1}^{C}e^{z_j}}
    $$

    PyTorch 的 `CrossEntropyLoss` 内部已经包含 `LogSoftmax + NLLLoss`，所以输入应是原始 logits，标签通常是类别编号，而不是 one-hot 后的概率。

    这一点很常考：

    - `CrossEntropyLoss` 输入：logits，形状通常是 `[batch, num_classes]`。
    - 标签：类别索引，形状通常是 `[batch]`。
    - 不要先对 logits 做 Softmax 再传入 `CrossEntropyLoss`。

    ## 11. MLP 本章易错点

    1. **`backward()` 不更新参数**：它只计算梯度。
    2. **`optimizer.step()` 才更新参数**。
    3. **验证/测试阶段要用 `model.eval()` 和 `torch.no_grad()`**。
    4. **没有激活函数的多层线性网络仍然是线性模型**。
    5. **二分类和多分类的输出层、标签格式、损失函数要匹配**。
    6. **Batch、Iteration、Epoch 不是同一个概念**。

    ---

    # 第三部分：CNN 与图像建模

    ## 1. 为什么图像不适合直接用普通 MLP

    图像不是普通的一维表格数据。相邻像素之间有空间关系，局部区域经常形成边缘、纹理、角点等模式。如果把图像直接展平成一维向量，再用全连接层处理，会出现三个问题：

    1. **空间结构被破坏**：相邻像素在向量里不一定还保持直观邻接关系。
    2. **参数量巨大**：每个像素和每个隐藏单元全连接，参数很容易爆炸。
    3. **不能自然利用平移不变性**：同一个边缘出现在图像左上角或右下角，本质上还是同一类局部模式。

    CNN 正是为了解决这些问题而出现。

    ## 2. CNN 的三个核心思想

    | 思想 | 解释 | 考试关键词 |
    |---|---|---|
    | 局部连接 | 每个输出只看输入的一小块区域 | 保留局部空间结构 |
    | 参数共享 | 同一个卷积核在整张图上滑动 | 减少参数量 |
    | 池化/层级特征 | 局部汇总，逐层形成更抽象特征 | 降采样、增强一定平移不变性 |

    简答题如果问 CNN 为什么适合图像识别，可以写：

    > CNN 利用图像的空间局部性，通过卷积核在局部区域滑动提取边缘、纹理等特征；同一个卷积核在不同位置共享参数，显著减少参数量；池化降低空间尺寸并保留显著响应，使网络逐层从低级特征构建高级语义。

    ## 3. 卷积和互相关

    课件强调：深度学习里的“卷积”实际常用的是互相关运算，也就是卷积核不翻转，直接和局部窗口对应相乘求和。

    对一个 $3\times3$ 局部窗口和一个 $3\times3$ 卷积核，输出位置的计算就是：

    $$
    y = \sum_{i=1}^{3}\sum_{j=1}^{3}x_{ij}k_{ij}
    $$

    这里 $x_{ij}$ 是局部像素，$k_{ij}$ 是卷积核参数。卷积核参数不是固定滤镜，而是在训练过程中通过反向传播学习得到的。

    常见判断题：

    - “CNN 实践中常用互相关，不需要翻转卷积核。”正确。
    - “卷积核参数在训练中不断更新。”正确。
    - “CNN 用全层全连接结构增加参数量。”错误。CNN 的重点是局部连接和参数共享。

    ## 4. 输出尺寸公式

    一维方向上，卷积输出尺寸为：

    $$
    O = \left\lfloor \frac{I + 2P - K}{S} \right\rfloor + 1
    $$

    其中：

    - $I$：输入尺寸。
    - $K$：卷积核大小。
    - $P$：padding 填充大小。
    - $S$：stride 步长。
    - $O$：输出尺寸。

    如果输入是 $10\times10$，卷积核 $3\times3$，步长 1，无填充，那么：

    $$
    O = 10 - 3 + 1 = 8
    $$

    输出就是 $8\times8$。

    Same、Valid、Full 的区别：

    | 类型 | 填充方式 | 输出大小趋势 |
    |---|---|---|
    | Valid | 不填充 | 输出变小 |
    | Same | 合适填充 | 输出尽量和输入一致 |
    | Full | 更充分填充 | 输出可能大于输入 |

    考试里如果看到“Valid 卷积大量填充，因此输出最大”，这是反的。Valid 不填充，输出通常最小。

    ## 5. 通道、NCHW、NHWC 与归一化

    彩色图像有 RGB 三个通道，灰度图像通常只有一个通道。深度学习框架对图像维度有不同约定：

    | 格式 | 含义 | 常见框架 |
    |---|---|---|
    | NCHW | Batch, Channel, Height, Width | PyTorch 常用 |
    | NHWC | Batch, Height, Width, Channel | TensorFlow 常见 |

    复习时不用把 RGB 当重点题单独背，但要知道通道会影响卷积核参数量。对输入通道数为 $C_{in}$、输出通道数为 $C_{out}$、卷积核大小为 $K\times K$ 的卷积层，参数量通常是：

    $$
    C_{out}\times C_{in}\times K\times K + C_{out}
    $$

    最后的 $C_{out}$ 是偏置项，如果该层不使用 bias，则没有这一项。

    归一化的作用是让输入尺度更稳定，便于训练。它不是改变类别，也不是增加样本，而是改变数值分布。

    ## 6. 池化

    池化对局部区域做汇总：

    - 最大池化：取局部最大值，保留最强响应。
    - 平均池化：取局部平均，保留整体趋势。

    池化通常没有可学习参数，主要作用是降采样、压缩空间尺寸、减少计算量，并增强一定的局部平移不变性。不要把池化说成全连接，也不要说池化是为了增加参数量。

    ## 7. 感受野

    感受野指某一层神经元能对应到原始输入图像上的区域大小。层越深，一个神经元看到的原始区域通常越大。

    递推公式常写成：

    $$
    RF_{l-1} = (RF_l - 1)\times s_l + k_l
    $$

    其中 $s_l$ 是步长，$k_l$ 是卷积核大小。步长变大时，前一层对应的感受野会扩大，不是缩小。

    小卷积核堆叠是经典设计：多层 $3\times3$ 卷积能逐步扩大感受野，同时参数量往往比一个大卷积核更少，还能引入更多非线性。

    ## 8. 典型 CNN 架构

    | 架构 | 复习重点 |
    |---|---|
    | LeNet | 早期 CNN，用卷积、池化、全连接完成数字识别 |
    | AlexNet | 更深网络，使用 ReLU、Dropout，推动深度 CNN 成功 |
    | VGG | 多层小卷积核堆叠，结构规整 |
    | GoogLeNet/Inception | 多尺度卷积分支并行，控制计算量 |
    | ResNet | 残差连接缓解深层网络退化和梯度传播困难 |

    ResNet 的核心不是“简单加深”，而是通过跳跃连接学习残差：

    $$
    y = F(x) + x
    $$

    如果 $F(x)$ 学不到有用变换，网络至少可以保留 $x$，这让很深的网络更容易训练。

    ## 9. CNN 训练技巧

    CNN 课件里还讲了训练技巧：

    - **BatchNorm**：稳定中间激活分布，常放在卷积层后、激活函数前或附近。
    - **Dropout**：训练时随机丢弃部分神经元输出，减少过拟合。
    - **数据增强**：翻转、裁剪、旋转等，扩大训练样本变化。
    - **权重衰减**：通过正则项限制参数过大。
    - **早停**：验证集不再改善时停止训练。
    - **类别不平衡处理**：`class_weight`、`pos_weight`、`WeightedRandomSampler` 等。

    类别不平衡常见考法：如果正负样本差异大，普通准确率可能误导。要看 precision、recall、F1、混淆矩阵，或者调整采样和损失权重。

    ## 10. CNN 本章易错点

    1. CNN 不是全连接结构，核心是局部连接和参数共享。
    2. 卷积输出尺寸由输入、卷积核、步长、填充共同决定。
    3. 深度学习实践里常用互相关而不是严格翻转卷积核。
    4. 池化通常没有可学习参数。
    5. Same 卷积依赖 padding 控制输出尺寸。
    6. ResNet 的残差连接是为了解决深层网络训练困难，不是简单增加参数。

    ---

    # 第四部分：RNN、GRU、LSTM 与序列建模

    ## 1. 为什么需要 RNN

    MLP 和普通 CNN 往往把一个样本当成相对独立的输入。但很多数据有顺序：文本、语音、时间序列、传感器数据。当前时刻的含义经常依赖前面的内容。

    RNN 的基本思想是：每个时间步不仅看当前输入 $x_t$，还看上一时刻隐藏状态 $h_{t-1}$。

    $$
    h_t = f(W_xx_t + W_hh_{t-1} + b)
    $$

    隐藏状态 $h_t$ 可以理解为“到当前为止，模型记住的历史摘要”。

    ## 2. RNN 的输入输出形状

    序列模型复习时一定要盯形状。常见输入可以理解为：

    $$
    [batch, seq\_len, input\_dim]
    $$

    不同框架或参数设置可能把 batch 和 seq_len 顺序交换，但含义不变：

    - batch：一次处理多少条序列。
    - seq_len：每条序列有多少时间步。
    - input_dim：每个时间步的特征维度。

    RNN 输出通常包括：

    - 每个时间步的输出。
    - 最后一个隐藏状态。

    做序列分类时，常取最后隐藏状态或池化后的序列表示；做序列标注时，每个时间步都要输出标签。

    ## 3. 梯度消失、梯度爆炸与梯度裁剪

    RNN 沿时间展开后很深，反向传播要穿过许多时间步。连乘会导致梯度越来越小或越来越大：

    - 梯度消失：远处历史信息难以影响当前参数更新。
    - 梯度爆炸：梯度过大导致训练不稳定。

    梯度裁剪是常见处理方式：

    ```python
    torch.nn.utils.clip_grad_norm_(model.parameters(), max_norm=1.0)
    ```

    它不是改变模型结构，而是在反向传播后、优化器更新前限制梯度范数。

    ## 4. GRU：用门控控制记忆

    GRU 用门控机制决定保留多少旧信息、写入多少新信息。核心门包括：

    - 更新门：决定旧隐藏状态保留多少。
    - 重置门：决定计算候选状态时参考多少过去信息。

    直觉上，GRU 像是在序列中做“有选择地记忆”：重要信息保留，不重要信息覆盖。

    考试不一定要求完整推导所有公式，但要知道 GRU 是为了解决普通 RNN 长序列记忆困难而引入门控结构。

    ## 5. LSTM：更完整的长期记忆机制

    LSTM 引入 cell state，并用三个门控制信息：

    - 遗忘门：决定旧记忆保留多少。
    - 输入门：决定新信息写入多少。
    - 输出门：决定当前输出暴露多少记忆。

    一个常见的直觉是：cell state 像一条长期信息通道，门控结构像阀门，控制信息留下、写入和输出。

    单选题常见陷阱：

    - “LSTM 没有门控结构。”错误。
    - “GRU/LSTM 的门通常使用 Sigmoid 输出 0 到 1 的控制系数。”正确。
    - “LSTM 彻底消除了所有梯度问题。”过度绝对，错误。

    ## 6. Padding、Packing 与变长序列

    一个 batch 中不同序列长度可能不同。为了组成张量，需要 padding 把短序列补齐。但补出来的位置不是有效信息，不能让模型把它当真实词或真实时间步。

    处理方式：

    - padding：补齐长度。
    - mask：告诉模型哪些位置无效。
    - packing：在 RNN 中更高效地处理变长序列。

    资料题如果给变长序列，要先看 padding 是否被正确处理。不能把补齐位置当成真实内容参与损失。

    ## 7. 双向 RNN、堆叠 RNN 和 1D 卷积

    - 双向 RNN 同时看过去和未来，适合整句已知的任务，如文本分类、序列标注。
    - 堆叠 RNN 把多层 RNN 叠起来，提高表示能力。
    - 1D 卷积/TCN 也可以处理序列，通过卷积窗口捕捉局部时间模式。

    这些方法的共同点是：都在处理序列依赖，只是路径不同。RNN 通过递归隐藏状态传递信息，1D 卷积通过局部窗口和层级堆叠捕捉上下文。

    ## 8. RNN 本章易错点

    1. RNN 的当前输出依赖当前输入和历史隐藏状态。
    2. 隐藏状态不是标签，而是历史信息表示。
    3. GRU/LSTM 的门控用于控制信息保留和更新。
    4. padding 只是补齐，不是有效数据。
    5. 梯度裁剪发生在 `loss.backward()` 之后、`optimizer.step()` 之前。

    ---

    # 第五部分：Attention 与 Transformer

    ## 1. 为什么注意力机制出现

    传统 Seq2Seq 用编码器把输入序列压成一个固定长度表示，再由解码器生成输出。问题是：输入很长时，一个固定向量很难保存所有细节。

    注意力机制的直觉是：生成当前输出时，不必只依赖一个压缩向量，而是可以回头看输入序列中最相关的位置。它不是平均看所有位置，而是按相关性分配权重。

    ## 2. Q、K、V 的含义

    注意力可以用查询、键、值来理解：

    - Query：我现在想找什么信息。
    - Key：每个位置提供的索引或特征标签。
    - Value：真正要汇总的信息内容。

    缩放点积注意力公式是：

    $$
    Attention(Q,K,V)=softmax\left(\frac{QK^T}{\sqrt{d_k}}\right)V
    $$

    拆开看：

    1. $QK^T$：计算 Query 和 Key 的匹配分数。
    2. $\sqrt{d_k}$：缩放项，防止维度大时点积过大，Softmax 过于尖锐。
    3. `softmax`：把分数变成权重。
    4. 乘以 $V$：按权重加权汇总信息。

    这道公式常考，但不要机械背。要能说出它的流程：**算相关性 -> 缩放 -> Softmax 归一化 -> 加权求和。**

    ## 3. Source Mask 与 Target Mask

    Mask 是 Transformer 资料题高频点。

    - Source Mask：通常用于遮住 padding 位置，避免模型关注无效 token。
    - Target Mask / Subsequent Mask：用于解码器自注意力，防止当前位置看到未来 token。

    如果题目说“机器翻译训练时，解码器可以直接看到目标句后面的词”，这违背自回归生成逻辑。Target Mask 的作用就是保证生成第 $t$ 个词时不能偷看 $t+1$ 之后的信息。

    ## 4. Multi-Head Attention

    多头注意力不是重复算同一个注意力，而是让不同头在不同子空间学习不同关系：

    $$
    head_i = Attention(QW_i^Q, KW_i^K, VW_i^V)
    $$

    然后把多个头拼接再线性变换：

    $$
    MultiHead(Q,K,V)=Concat(head_1,\ldots,head_h)W^O
    $$

    直觉上，一个头可能关注局部词序，一个头可能关注语法依赖，一个头可能关注远距离语义关系。

    ## 5. 位置编码

    Self-Attention 本身对顺序不敏感。如果不加入位置信息，“我爱你”和“你爱我”可能难以区分。因此 Transformer 需要位置编码。

    经典正弦位置编码形式为：

    $$
    PE(pos,2i)=\sin\left(\frac{pos}{10000^{2i/d_{model}}}\right)
    $$

    $$
    PE(pos,2i+1)=\cos\left(\frac{pos}{10000^{2i/d_{model}}}\right)
    $$

    这里：

    - $pos$ 是位置。
    - $i$ 是维度索引。
    - $d_{model}$ 是模型隐藏维度。

    课件也涉及 RoPE。复习时把它理解为一种把相对位置信息注入表示的现代位置编码方法即可，不要把复习重点放在复杂推导上。

    ## 6. Transformer 的整体结构

    Transformer 由 Encoder 和 Decoder 堆叠组成。

    Encoder 层通常包含：

    1. Multi-Head Self-Attention
    2. Add & Norm
    3. Feed Forward Network
    4. Add & Norm

    Decoder 层通常多一个 Cross-Attention：

    1. Masked Self-Attention
    2. Cross-Attention（Query 来自解码器，Key/Value 来自编码器输出）
    3. Feed Forward Network

    残差连接和 LayerNorm 的作用是稳定深层训练，让梯度和表示更容易传播。

    ## 7. Transformer 和 RNN 的区别

    | 角度 | RNN | Transformer |
    |---|---|---|
    | 处理方式 | 按时间步递归 | 可并行处理序列位置 |
    | 长距离依赖 | 依赖隐藏状态逐步传递 | 注意力直接建立位置间联系 |
    | 位置信息 | 顺序天然来自递归 | 需要位置编码 |
    | 训练效率 | 难完全并行 | 更适合并行计算 |

    简答题如果问 Transformer 为什么能替代 RNN 做很多序列任务，可以写：

    > Transformer 使用自注意力直接计算序列中任意位置之间的相关性，避免 RNN 逐时间步递归带来的并行性差和长距离依赖困难；同时通过位置编码保留顺序信息，通过多头注意力从多个子空间建模关系。

    ## 8. 本章易错点

    1. Attention 不是简单平均，而是按相关性加权。
    2. $\sqrt{d_k}$ 是缩放项，不是随便加的常数。
    3. Source Mask 主要处理 padding，Target Mask 防止看未来。
    4. Transformer 不依赖 RNN 递归，但仍需要位置编码。
    5. Cross-Attention 中 Query 和 Key/Value 来自不同侧。

    ---

    # 第六部分：NLP、LLM 与生成式算法

    ## 1. NLP 从规则系统到神经网络

    第二部分 PPT 从早期自然语言处理讲起。ELIZA 和 SHRDLU 代表了早期符号 AI 的思路：人手写规则、关键词和模板，让系统看起来能对话或执行指令。

    这类系统的特点：

    - 在封闭场景里能工作。
    - 依赖人工规则。
    - 泛化能力弱。
    - 无法真正从大规模数据中学习语言规律。

    这条历史线的意义是：后来的词向量、RNN、Transformer、LLM 都是在解决“规则系统太死、无法泛化”的问题。

    ## 2. 词表示：从 one-hot 到 Word2Vec

    one-hot 表示的问题是：每个词都是孤立编号，无法表达语义相似性。Word2Vec 的思想是从上下文中学习词表示，让相似语境中的词向量更接近。

    Skip-Gram 的目标是：给定中心词，预测上下文词。Negative Sampling 用少量负样本近似原本庞大的 softmax 计算，降低训练成本。

    图学习 PPT 第一部分还把 DeepWalk/node2vec 和 Word2Vec 联系起来：图游走产生节点序列，再用类似 Skip-Gram 的方式学习节点向量。这是 NLP 方法迁移到图学习的典型例子。

    ## 3. 大语言模型的主线

    大语言模型不是凭空出现的，它建立在 Transformer、海量语料、自监督预训练和生成式建模上。PPT 中提到 LLM 的快速进化、Prompt、RLHF、MoE、DeepSeek 等内容。

    复习时抓四条线：

    1. **架构基础**：Transformer 提供并行序列建模能力。
    2. **训练目标**：语言模型通常学习根据上下文预测 token。
    3. **对齐方法**：RLHF 用人类反馈帮助模型更符合人类偏好。
    4. **效率结构**：MoE 用多个专家子网络，让每个输入激活部分专家，提高参数规模和计算效率之间的平衡。

    ## 4. Prompt 与上下文学习

    Prompt 可以理解为给模型的任务说明和上下文。它不是改变模型参数，而是在输入层面引导模型输出。

    常见 Prompt 题目会考：

    - 指令是否清晰。
    - 是否提供必要上下文。
    - 是否限定输出格式。
    - 是否给出示例。

    如果资料题给一个很模糊的 Prompt，让你改写，可以按“任务 + 背景 + 约束 + 输出格式 + 示例”的顺序补。

    ## 5. MoE 与 DeepSeek-MoE

    MoE（Mixture of Experts）可以理解为“多个专家 + 路由器”。路由器根据输入决定激活哪些专家。

    直觉：

    - 普通 FFN：所有参数都参与每个 token 的计算。
    - MoE：模型有很多专家，但每个 token 只激活其中一部分。

    这样可以扩大总参数量，同时控制单次计算成本。PPT 中的 DeepSeek-MoE 属于现代大模型结构和效率优化的例子。考试通常不会要求你推导 MoE 细节，但可能让你判断它为什么能提高效率。

    ## 6. 生成式算法：GAN、VAE、Diffusion

    新课件中补充了生成式模型。生成式算法的共同目标是学习数据分布，并生成新的样本。

    | 方法 | 直觉 | 常见考法 |
    |---|---|---|
    | GAN | 生成器和判别器对抗训练 | 生成器想骗过判别器，判别器想分辨真假 |
    | VAE | 学习潜变量分布，再从潜空间采样生成 | 编码器、解码器、潜变量、重构 |
    | Diffusion | 逐步加噪，再学习反向去噪生成 | 正向扩散、反向去噪 |

    GAN 的结构：

    - Generator：输入噪声，生成假样本。
    - Discriminator：判断样本是真实还是生成。

    VAE 的结构：

    - Encoder：把输入压到潜变量分布。
    - Decoder：从潜变量重构样本。

    Diffusion 的直觉：

    - 正向过程：逐步给真实数据加噪声。
    - 反向过程：模型学习从噪声一步步还原数据。

    这些内容属于第二梯队，重点是知道它们分别“怎样生成”，不要陷入过深数学推导。

    ---

    # 第七部分：图学习与图神经网络

    ## 1. 为什么需要图学习

    图像是规则网格，文本是序列，而很多真实数据是不规则关系结构：社交网络、交通网络、分子结构、知识图谱、论文引用、材料结构。图学习就是为了处理这类由节点和边构成的数据。

    图可以写成：

    $$
    G=(V,E)
    $$

    - $V$ 是节点集合。
    - $E$ 是边集合。
    - 节点可以有特征。
    - 边可以有方向、权重和属性。
    - 有些图还有全局属性 $U$。

    ## 2. 图的基本表示

    邻接矩阵 $A$ 表示节点之间是否相连：

    $$
    A_{ij}=1
    $$

    表示节点 $i$ 和节点 $j$ 有边；如果是加权图，$A_{ij}$ 可以表示边权。

    度矩阵 $D$ 记录节点连接数量：

    $$
    D_{ii}=\sum_j A_{ij}
    $$

    拉普拉斯矩阵常写成：

    $$
    L = D - A
    $$

    这些矩阵不是为了背符号，而是为了把图结构变成模型可以计算的形式。

    ## 3. 图学习任务

    | 任务 | 输入 | 输出 |
    |---|---|---|
    | 节点分类 | 整张图和部分节点标签 | 每个节点的类别 |
    | 边预测 | 图中已有连接 | 判断两个节点是否应有边 |
    | 图分类 | 一张图作为样本 | 整张图的类别 |
    | 节点/图表示学习 | 图结构和属性 | 向量表示 |

    材料科学中，原子可以看作节点，化学键或相互作用可以看作边，整个分子或晶体可以看作图。这就是图学习和 AI for material science 连接的地方。

    ## 4. 图游走、DeepWalk、node2vec

    图游走方法的思想是：在图上随机走，得到节点序列，再把节点序列当作“句子”，用类似 Word2Vec 的方法学习节点表示。

    DeepWalk：

    - 随机游走生成节点序列。
    - 用 Skip-Gram 学习节点 embedding。
    - 结构相近的节点会得到相近向量。

    node2vec 改进了游走策略，用参数控制更像 BFS 还是 DFS：

    - BFS 倾向看局部邻居，适合捕捉同质性。
    - DFS 倾向走得更远，适合捕捉结构角色。

    ## 5. GNN 与消息传递

    GNN 的核心是消息传递：每个节点从邻居接收信息，聚合后更新自己的表示。

    一个通用形式可以写成：

    $$
    h_v^{(k)} = UPDATE^{(k)}\left(h_v^{(k-1)}, AGGREGATE^{(k)}\left(\{h_u^{(k-1)}:u\in \mathcal{N}(v)\}\right)\right)
    $$

    这里：

    - $h_v^{(k)}$ 是第 $k$ 层后节点 $v$ 的表示。
    - $\mathcal{N}(v)$ 是节点 $v$ 的邻居。
    - `AGGREGATE` 把邻居信息汇总。
    - `UPDATE` 用汇总信息更新当前节点。

    这和 CNN 有类比：CNN 在规则网格上聚合局部像素，GNN 在不规则图上聚合邻居节点。

    ## 6. GCN

    新 PPT 重点补充了 GCN。经典 GCN 层可以写成：

    $$
    H^{(l+1)}=\sigma\left(\tilde{D}^{-\frac{1}{2}}\tilde{A}\tilde{D}^{-\frac{1}{2}}H^{(l)}W^{(l)}\right)
    $$

    其中：

    - $\tilde{A}=A+I$，给邻接矩阵加自环，让节点保留自己的信息。
    - $\tilde{D}$ 是 $\tilde{A}$ 对应的度矩阵。
    - $H^{(l)}$ 是第 $l$ 层节点表示。
    - $W^{(l)}$ 是可学习权重。
    - $\sigma$ 是激活函数。

    这个公式的直觉是：先把邻居和自己信息按归一化邻接关系混合，再做线性变换和非线性激活。归一化的作用是避免度数大的节点信息尺度过大。

    ## 7. GAT 与 GraphSAGE

    GAT 的核心是注意力：不同邻居对当前节点的重要性不同，所以不必平均看所有邻居。它会学习邻居权重，再加权聚合。

    GraphSAGE 的核心是采样和聚合：对大图不一定每次看完所有邻居，而是采样一部分邻居，用 mean、pooling、LSTM 等聚合器形成表示。当前 OCR 中更明确的是 GCN、GAT 和消息传递；GraphSAGE 可以作为 GNN 家族中“采样邻居聚合”的代表来识别。

    ## 8. 图学习与材料科学

    图学习课件第一部分提到材料领域潜在应用。材料、分子、晶体天然可以转成图：

    - 原子是节点。
    - 化学键或空间邻近关系是边。
    - 原子类型、电负性、价态等可以是节点特征。
    - 分子性质、材料性能可以是图级标签。

    资料题如果给材料结构图，要先判断任务级别：

    - 预测某个原子的属性：节点级任务。
    - 预测两个原子是否形成关系：边级任务。
    - 预测整个分子/材料性质：图级任务。

    ## 9. PINN 与 AI for material science

    老师提示 PINN 和 AI for material science 偏弱，但仍需知道主线。PINN 的思想是把物理方程或物理约束放进损失函数，让模型不仅拟合数据，还尽量满足物理规律。

    常见形式可以理解为：

    $$
    L = L_{data} + \lambda L_{physics}
    $$

    - $L_{data}$：数据拟合误差。
    - $L_{physics}$：违反物理方程或边界条件的误差。
    - $\lambda$：平衡两部分的重要性。

    AI for material science 的复习目标不是背某个材料案例，而是能说明：深度学习可以从材料数据、结构图、实验/模拟结果中学习规律，用于性质预测、筛选、生成或优化。

    ## 10. 图学习本章易错点

    1. 节点和边不是一回事：节点表示对象，边表示关系。
    2. 邻接矩阵表达连接关系，度矩阵表达连接数量或权重总和。
    3. GCN 不是普通 CNN 直接套到图像，而是在图结构上做邻居聚合。
    4. 加自环是为了保留节点自身信息。
    5. GAT 的注意力权重表示邻居重要性。
    6. 图级任务输出整张图标签，节点级任务输出每个节点标签。

    ---

    # 第八部分：三类题型怎么答

    ## 1. 单选题

    单选题第一步看清题干问“正确的是”还是“错误的是”。深度学习单选题常见陷阱有：

    - 把流程顺序说反：例如先 `optimizer.step()` 再 `loss.backward()`。
    - 把模型结构说反：例如 CNN 是全连接、RNN 不看历史、Transformer 不需要位置编码。
    - 把公式条件说错：例如 Valid 卷积大量 padding、CrossEntropyLoss 输入概率。
    - 把作用夸大：例如 Dropout 可以完全消除过拟合、LSTM 彻底解决所有梯度问题。

    做题时建议先找绝对化词：全部、必然、完全、不需要、只要。这些词经常藏错误选项。

    ## 2. 简答题

    简答题不要写散句。推荐四句式：

    1. **定义**：这个概念是什么。
    2. **作用**：它解决什么问题。
    3. **位置**：它在模型结构或训练流程中出现在哪里。
    4. **易错区别**：它不是什么，和相近概念有什么区别。

    例：简答“池化的作用”：

    > 池化是在局部区域内取最大值或平均值的降采样操作。它用于降低特征图空间尺寸、减少计算量，并保留局部显著响应，从而增强一定的平移不变性。池化通常放在卷积和激活之后，不负责学习卷积核参数。它不是全连接层，也不是为了增加参数量。

    ## 3. 资料题

    资料题一定先读材料，不要直接背概念。按下面顺序圈信息：

    1. 输入是什么：图像、序列、图、文本、材料结构？
    2. 任务是什么：分类、回归、生成、节点预测、图分类？
    3. 模型是什么：MLP、CNN、RNN、Transformer、GNN？
    4. 损失和输出是否匹配？
    5. 训练流程有没有缺失？
    6. 材料里有没有 padding、mask、类别不平衡、图结构、物理约束等提示？

    资料题答案要把材料里的具体信息带回去。例如题目给了 `BCEWithLogitsLoss`，你就要说模型输出应是 logits；题目给了目标序列生成，就要说 Target Mask 防止看未来。

    ---

    # 第九部分：考前总清单

    ## 1. 第一梯队必须会

    - TLU、感知机、ADALINE 的递进关系。
    - XOR 为什么单层感知机解决不了。
    - 训练循环：forward、loss、zero_grad、backward、step。
    - MLP 为什么需要激活函数。
    - BCE、CrossEntropyLoss 的输入输出匹配。
    - CNN 输出尺寸公式。
    - CNN 的局部连接、参数共享、池化、感受野。
    - RNN 隐藏状态、GRU/LSTM 门控。
    - padding、packing、梯度裁剪。
    - Attention 的 Q/K/V 公式。
    - Source Mask、Target Mask、多头注意力、位置编码。

    ## 2. 第二梯队至少会

    - NLP 从规则系统到神经网络的历史主线。
    - Word2Vec、Skip-Gram、Negative Sampling 的直觉。
    - LLM 依赖 Transformer、预训练、Prompt、RLHF。
    - MoE 用路由器选择专家，控制计算成本。
    - GAN、VAE、Diffusion 的生成直觉。
    - 图的节点、边、邻接矩阵、度矩阵。
    - DeepWalk/node2vec 用图游走学习节点表示。
    - GCN 用归一化邻接矩阵做邻居聚合。
    - GAT 用注意力区分邻居重要性。
    - PINN 把物理约束放进损失函数。

    ## 3. 最后一遍记忆句

    - 神经元：加权求和、偏置/阈值、激活输出。
    - MLP：线性层加激活函数，靠损失和梯度训练。
    - CNN：局部连接、参数共享、池化降采样。
    - RNN：当前输入加历史隐藏状态。
    - GRU/LSTM：门控控制记忆保留和更新。
    - Attention：按相关性给 Value 加权求和。
    - Transformer：自注意力、前馈网络、残差、归一化、位置编码。
    - GNN：节点从邻居接收消息并更新表示。
    - PINN：数据误差加物理约束误差。

    <style>@media print { @page { margin: 8mm; } body { font-size: 10.5pt; line-height: 1.35; } h1, h2, h3 { page-break-after: avoid; } table { page-break-inside: avoid; } }</style>
    """).strip()
    lines.extend(["", body, ""])
    return "\n".join(lines)


def write_final_review_markdown(outline, manifest):
    if FINAL_REVIEW_MD.exists() and len(FINAL_REVIEW_MD.read_text(encoding="utf-8")) > 60000:
        md = FINAL_REVIEW_MD.read_text(encoding="utf-8")
    else:
        md = build_final_review_markdown(outline, manifest)
    FINAL_REVIEW_MD.write_text(md, encoding="utf-8")
    TOPIC_GUIDE_OUT.mkdir(parents=True, exist_ok=True)
    FINAL_REVIEW_SITE_MD.write_text(md, encoding="utf-8")
    return md


CSS = r"""
:root {
  --blue: #20a8e0;
  --blue-dark: #1376a8;
  --ink: #1f2937;
  --muted: #6b7280;
  --line: #e5e7eb;
  --soft: #f6f8fb;
  --panel: #ffffff;
  --green: #16a34a;
}
* { box-sizing: border-box; }
body {
  margin: 0;
  font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "PingFang SC", "Microsoft YaHei", sans-serif;
  color: var(--ink);
  background: #f7f9fc;
  letter-spacing: 0;
}
a { color: var(--blue-dark); text-decoration: none; }
header {
  position: sticky;
  top: 0;
  z-index: 20;
  background: rgba(255,255,255,.96);
  border-bottom: 1px solid var(--line);
  backdrop-filter: blur(8px);
}
.topbar {
  max-width: 1440px;
  margin: 0 auto;
  padding: 14px 24px;
  display: grid;
  grid-template-columns: 290px 1fr;
  gap: 24px;
  align-items: center;
}
.brand h1 { margin: 0; font-size: 20px; line-height: 1.2; }
.brand p { margin: 4px 0 0; color: var(--muted); font-size: 13px; }
nav {
  display: flex;
  gap: 6px;
  align-items: center;
  flex-wrap: wrap;
}
nav button {
  border: 0;
  background: transparent;
  color: #445166;
  padding: 8px 12px;
  border-radius: 8px;
  cursor: pointer;
  font-weight: 600;
}
nav button.active { color: #fff; background: var(--blue); }
main {
  max-width: 1440px;
  margin: 0 auto;
  padding: 22px 24px 48px;
}
.view { display: none; }
.view.active { display: block; }
.hero-grid {
  display: grid;
  grid-template-columns: 1.15fr .85fr;
  gap: 18px;
  align-items: start;
  margin-bottom: 18px;
}
.panel {
  background: var(--panel);
  border: 1px solid var(--line);
  border-radius: 8px;
  padding: 18px;
}
.panel h2, .panel h3 { margin: 0 0 12px; }
.muted, .meta { color: var(--muted); font-size: 13px; }
.stats {
  display: grid;
  grid-template-columns: repeat(5, minmax(120px, 1fr));
  gap: 10px;
}
.stat {
  background: #fff;
  border: 1px solid var(--line);
  border-radius: 8px;
  padding: 14px;
}
.stat strong { display: block; font-size: 23px; color: #0f172a; }
.toolbar {
  display: grid;
  grid-template-columns: minmax(260px, 1fr) auto auto;
  gap: 10px;
  margin: 14px 0;
}
input[type="search"] {
  border: 1px solid #d6dce6;
  border-radius: 8px;
  padding: 11px 13px;
  font-size: 15px;
  background: #fff;
}
.btn, .btn.secondary {
  border: 1px solid var(--blue);
  border-radius: 8px;
  padding: 10px 13px;
  cursor: pointer;
  font-weight: 700;
  background: var(--blue);
  color: #fff;
}
.btn.secondary { background: #fff; color: var(--blue-dark); }
.layout {
  display: grid;
  grid-template-columns: minmax(0, 1fr) 330px;
  gap: 16px;
}
.result, .q, .answer-item {
  border-top: 1px solid var(--line);
  padding: 14px 0;
}
.result:first-child, .q:first-child, .answer-item:first-child { border-top: 0; }
.result-title {
  display: flex;
  gap: 10px;
  justify-content: space-between;
  align-items: flex-start;
  font-weight: 700;
}
.path { margin-top: 4px; color: var(--blue-dark); font-size: 13px; }
.snippet { margin: 8px 0 0; line-height: 1.65; }
mark { background: #dff3ff; color: #0b5e8e; border-radius: 3px; padding: 0 2px; }
.chips { display: flex; flex-wrap: wrap; gap: 8px; margin-top: 10px; }
.chip {
  border: 1px solid #dce3ec;
  background: #fff;
  color: #334155;
  border-radius: 999px;
  padding: 5px 10px;
  font-size: 13px;
  cursor: pointer;
}
.term-list {
  display: grid;
  grid-template-columns: repeat(auto-fit, minmax(132px, 1fr));
  gap: 8px;
}
.term-btn {
  border: 1px solid var(--line);
  background: #fff;
  text-align: left;
  border-radius: 8px;
  padding: 9px 10px;
  cursor: pointer;
}
.term-btn strong { display: block; }
.term-btn span { color: var(--muted); font-size: 12px; }
.course-layout {
  display: grid;
  grid-template-columns: 270px minmax(0, 1fr);
  gap: 16px;
  align-items: start;
}
.side-index {
  position: sticky;
  top: 92px;
  max-height: calc(100vh - 110px);
  overflow: auto;
}
.side-index a {
  display: block;
  padding: 8px 10px;
  border-radius: 8px;
  color: #334155;
  font-size: 13px;
}
.side-index a:hover { background: #edf7fd; color: var(--blue-dark); }
.course-card { margin-bottom: 16px; }
.course-head {
  display: flex;
  justify-content: space-between;
  gap: 12px;
  align-items: flex-start;
  margin-bottom: 12px;
}
.chapter { border-top: 1px solid var(--line); padding-top: 14px; margin-top: 14px; }
table {
  width: 100%;
  border-collapse: collapse;
  font-size: 14px;
  background: #fff;
  border: 1px solid var(--line);
  border-radius: 8px;
  overflow: hidden;
}
th, td {
  border-bottom: 1px solid var(--line);
  padding: 10px 11px;
  vertical-align: top;
  line-height: 1.55;
}
th { background: #f8fafc; text-align: left; color: #475569; font-size: 13px; }
tr:last-child td { border-bottom: 0; }
.quiz-head {
  display: flex;
  justify-content: space-between;
  gap: 12px;
  align-items: center;
  margin-bottom: 12px;
}
.quiz-actions { display: flex; gap: 8px; flex-wrap: wrap; justify-content: flex-end; }
.quiz-section { margin: 18px 0; }
.q-stem { font-weight: 700; line-height: 1.6; }
.options { list-style: none; padding: 0; margin: 8px 0 0; }
.options li { margin: 5px 0; line-height: 1.5; }
.badge {
  display: inline-block;
  border: 1px solid #dbe3ec;
  border-radius: 999px;
  padding: 2px 8px;
  color: var(--muted);
  font-size: 12px;
  margin-left: 6px;
}
.blank-answer {
  display: inline-block;
  min-width: 140px;
  border-bottom: 1px solid #9aa6b2;
}
.answers { background: #f8fbff; }
.formula {
  font-family: "SFMono-Regular", Consolas, monospace;
  background: #f1f5f9;
  color: #0f4b6e;
  padding: 8px 10px;
  border-radius: 8px;
  margin: 6px 0;
  overflow-wrap: anywhere;
}
.review-section { border-top: 1px solid var(--line); padding-top: 18px; margin-top: 18px; }
.guide-grid {
  display: grid;
  grid-template-columns: repeat(auto-fit, minmax(260px, 1fr));
  gap: 12px;
  margin: 14px 0 22px;
}
.guide-card {
  border: 1px solid var(--line);
  border-radius: 8px;
  padding: 13px;
  background: #fff;
}
.guide-card strong { display: block; margin-bottom: 6px; }
.guide-card .btn { display: inline-block; margin-top: 10px; padding: 8px 10px; font-size: 13px; }
.book-shell {
  display: grid;
  grid-template-columns: 250px minmax(0, 1fr);
  gap: 16px;
  align-items: start;
}
.book-sidebar {
  position: sticky;
  top: 92px;
  max-height: calc(100vh - 112px);
  overflow: auto;
}
.book-nav-item {
  width: 100%;
  border: 1px solid var(--line);
  background: #fff;
  color: #334155;
  border-radius: 8px;
  padding: 10px 11px;
  margin: 7px 0;
  text-align: left;
  cursor: pointer;
}
.book-nav-item strong { display: block; font-size: 14px; }
.book-nav-item span { display: block; color: var(--muted); font-size: 12px; margin-top: 3px; }
.book-nav-item.active {
  border-color: #9dd8f3;
  background: #eaf7fd;
  color: #0f5f8c;
}
.book-reader {
  background: #fff;
  border: 1px solid var(--line);
  border-radius: 8px;
  padding: 22px;
}
.book-toolbar {
  display: flex;
  justify-content: space-between;
  gap: 16px;
  align-items: flex-start;
  border-bottom: 1px solid var(--line);
  padding-bottom: 16px;
  margin-bottom: 18px;
}
.toc-image-strip {
  display: grid;
  grid-template-columns: repeat(auto-fit, minmax(280px, 1fr));
  gap: 14px;
  margin: 16px 0 22px;
}
.toc-image-strip figure {
  margin: 0;
  border: 1px solid var(--line);
  border-radius: 8px;
  background: #f8fbff;
  padding: 10px;
}
.toc-image-strip img {
  display: block;
  width: 100%;
  border-radius: 6px;
  border: 1px solid #e2e8f0;
}
.toc-image-strip figcaption { margin-top: 8px; font-size: 13px; font-weight: 700; }
.book-layout {
  display: grid;
  grid-template-columns: 230px minmax(0, 1fr);
  gap: 22px;
  align-items: start;
}
.book-toc {
  position: sticky;
  top: 104px;
  max-height: calc(100vh - 130px);
  overflow: auto;
  border-right: 1px solid var(--line);
  padding-right: 12px;
}
.book-toc h3 { margin-top: 0; }
.book-toc a {
  display: block;
  padding: 5px 0;
  color: #475569;
  font-size: 13px;
  line-height: 1.35;
}
.book-toc .toc-level-3 { padding-left: 12px; color: #64748b; }
.markdown-body {
  max-width: 860px;
  font-size: 16px;
  line-height: 1.86;
}
.markdown-body h2 {
  font-size: 28px;
  margin: 28px 0 12px;
  border-bottom: 1px solid var(--line);
  padding-bottom: 8px;
}
.markdown-body h3 { font-size: 22px; margin: 24px 0 10px; }
.markdown-body h4 { font-size: 18px; margin: 20px 0 8px; color: #334155; }
.markdown-body p { margin: 10px 0; }
.markdown-body ul { padding-left: 22px; }
.markdown-body li { margin: 5px 0; }
.markdown-body code {
  font-family: "SFMono-Regular", Consolas, monospace;
  background: #f1f5f9;
  border: 1px solid #e2e8f0;
  border-radius: 5px;
  padding: 1px 5px;
}
.markdown-body pre {
  overflow: auto;
  background: #0f172a;
  color: #e2e8f0;
  border-radius: 8px;
  padding: 14px;
  line-height: 1.55;
}
.markdown-body pre code {
  background: transparent;
  border: 0;
  color: inherit;
  padding: 0;
}
.math-block {
  margin: 14px 0;
  padding: 10px 12px;
  border: 1px solid #e2e8f0;
  border-radius: 8px;
  background: #f8fafc;
  overflow-x: auto;
}
.markdown-body mjx-container {
  max-width: 100%;
  overflow-x: auto;
  overflow-y: hidden;
}
.table-line {
  font-family: "SFMono-Regular", Consolas, monospace;
  background: #f8fafc;
  border: 1px solid var(--line);
  border-radius: 6px;
  padding: 6px 8px;
  overflow-wrap: anywhere;
}
@media print {
  header, .book-sidebar, .book-toc, .quiz-actions, .toc-image-strip figcaption { display: none !important; }
  body { background: #fff; }
  main { max-width: none; padding: 0; }
  .panel, .book-reader { border: 0; padding: 0; }
  .book-shell, .book-layout { display: block; }
  .markdown-body { max-width: none; font-size: 10.5pt; line-height: 1.35; }
  .markdown-body h2, .markdown-body h3, .markdown-body h4 { page-break-after: avoid; }
  @page { margin: 7mm; }
}
@media (max-width: 920px) {
  .topbar, .hero-grid, .layout, .course-layout, .toolbar, .book-shell, .book-layout { grid-template-columns: 1fr; }
  .stats { grid-template-columns: repeat(2, minmax(120px, 1fr)); }
  .side-index, .book-sidebar, .book-toc { position: static; max-height: none; }
  .book-toc { border-right: 0; border-bottom: 1px solid var(--line); padding-bottom: 12px; }
}
"""


JS = r"""
const DB = window.COURSE_DB;
const $ = (id) => document.getElementById(id);
const norm = (s) => (s || "").toString().toLowerCase();
const escapeHtml = (s) => (s || "").replace(/[&<>"']/g, c => ({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}[c]));
const lastQuizState = { beginner: null, standard: null };

function queueMathTypeset(root = document.body, attempt = 0) {
  const run = () => {
    if (window.MathJax && MathJax.typesetPromise) {
      MathJax.typesetPromise([root]).catch(err => console.warn("MathJax typeset failed", err));
    } else if (attempt < 30) {
      window.setTimeout(() => queueMathTypeset(root, attempt + 1), 120);
    }
  };
  if (window.MathJax && MathJax.startup && MathJax.startup.promise) {
    MathJax.startup.promise.then(run);
  } else {
    window.setTimeout(run, 120);
  }
}

function showView(name) {
  document.querySelectorAll(".view").forEach(v => v.classList.toggle("active", v.id === `view-${name}`));
  document.querySelectorAll("nav button").forEach(b => b.classList.toggle("active", b.dataset.view === name));
}

function highlight(text, query) {
  const safe = escapeHtml(text || "");
  const q = (query || "").trim();
  if (!q) return safe;
  const parts = q.split(/\s+/).filter(Boolean).slice(0, 6).map(x => x.replace(/[.*+?^${}()|[\]\\]/g, "\\$&"));
  if (!parts.length) return safe;
  return safe.replace(new RegExp(`(${parts.join("|")})`, "gi"), "<mark>$1</mark>");
}

function scoreChunk(chunk, query) {
  const q = norm(query).split(/\s+/).filter(Boolean);
  if (!q.length) return 0;
  const title = norm(chunk.title + " " + (chunk.heading_path || []).join(" "));
  const text = norm(chunk.text);
  const terms = norm((chunk.terms || []).join(" "));
  let score = 0;
  q.forEach(part => {
    if (title.includes(part)) score += 10;
    if (terms.includes(part)) score += 8;
    const matches = text.match(new RegExp(part.replace(/[.*+?^${}()|[\]\\]/g, "\\$&"), "g"));
    if (matches) score += Math.min(matches.length, 12);
  });
  return score;
}

function shuffle(arr) {
  const copy = [...arr];
  for (let i = copy.length - 1; i > 0; i -= 1) {
    const j = Math.floor(Math.random() * (i + 1));
    [copy[i], copy[j]] = [copy[j], copy[i]];
  }
  return copy;
}

function safeFileName(text) {
  return (text || "试卷").replace(/[\\/:*?"<>|]+/g, "-").replace(/\s+/g, "-").slice(0, 80);
}

function downloadTextFile(filename, text) {
  const blob = new Blob([text], {type: "text/markdown;charset=utf-8"});
  const url = URL.createObjectURL(blob);
  const a = document.createElement("a");
  a.href = url;
  a.download = filename;
  document.body.appendChild(a);
  a.click();
  a.remove();
  URL.revokeObjectURL(url);
}

function printStyleBlock() {
  return `<style>
@media print {
  @page { margin: 7mm; }
  body { font-size: 10.5pt; line-height: 1.32; }
  h1, h2, h3 { page-break-after: avoid; margin: 0.45em 0 0.25em; }
  p, li { margin: 0.2em 0; }
  .blank { display: inline-block; min-width: 42mm; border-bottom: 1px solid #777; }
}
</style>`;
}

function questionMarkdown(q, idx) {
  if (q.type === "mcq") {
    return `${idx}. ${q.stem}\n\n${q.options.map(o => `   ${o}`).join("\n")}`;
  }
  if (q.type === "material") {
    const subs = q.subquestions || q.questions || [];
    return `${idx}. ${q.stem}\n\n材料：${q.material || ""}\n\n${subs.map((x, i) => `   ${i + 1}. ${x}`).join("\n")}\n\n答题区：\n\n\n`;
  }
  return `${idx}. ${q.stem}\n\n答题区：\n\n\n`;
}

function answerMarkdown(q, idx) {
  return `${idx}. [${q.id}] 答案：${q.answer}\n\n解析：${q.explanation}\n\n来源：${q.source} · ${q.topic}`;
}

function exportQuizMarkdown(mode, part) {
  const state = lastQuizState[mode];
  if (!state || !state.chosen.length) return;
  const modeName = mode === "beginner" ? "基础练习" : "标准组卷";
  const title = `${modeName}${state.topic ? "：" + state.topic : ""}`;
  const labels = {mcq: "单选题", short: "简答题", material: "资料题"};
  const lines = [
    printStyleBlock(),
    `# ${title}${part === "questions" ? "（题目版）" : "（答案解析版）"}`,
    "",
    `生成时间：${new Date().toLocaleString("zh-CN")}`,
    "",
    part === "questions" ? "> 打印建议：题目和答案分开打印；本文件已内置较小页边距样式。" : "> 打印建议：答案解析单独打印或仅在核对时查看。",
    "",
  ];
  if (part === "questions") {
    state.spec.forEach(([type]) => {
      const items = state.chosen.filter(q => q.type === type);
      if (!items.length) return;
      lines.push(`## ${labels[type]}`, "");
      items.forEach((q, i) => lines.push(questionMarkdown(q, i + 1), ""));
    });
  } else {
    state.chosen.forEach((q, i) => lines.push(answerMarkdown(q, i + 1), ""));
  }
  const suffix = part === "questions" ? "题目" : "答案解析";
  downloadTextFile(`${safeFileName(title)}-${suffix}.md`, lines.join("\n"));
}

function renderStats() {
  $("stats").innerHTML = `
    <div class="stat"><strong>${DB.manifest.length}</strong><span>课件文件</span></div>
    <div class="stat"><strong>${DB.chunks.length}</strong><span>可检索片段</span></div>
    <div class="stat"><strong>${Object.keys(DB.terms).length}</strong><span>索引词条</span></div>
    <div class="stat"><strong>${DB.quizBank.length}</strong><span>标准题</span></div>
    <div class="stat"><strong>${DB.beginnerQuizBank.length}</strong><span>基础题</span></div>
  `;
}

function renderTermCloud() {
  const terms = Object.values(DB.terms).sort((a, b) => b.total_count - a.total_count).slice(0, 42);
  $("termCloud").innerHTML = terms.map(t => `
    <button class="term-btn" data-term="${escapeHtml(t.term)}">
      <strong>${escapeHtml(t.term)}</strong>
      <span>${t.total_count} 次 · ${t.files.length} 文件</span>
    </button>
  `).join("");
  $("termCloud").querySelectorAll("button").forEach(btn => btn.addEventListener("click", () => {
    $("searchInput").value = btn.dataset.term;
    runSearch(btn.dataset.term);
    showView("search");
  }));
}

function runSearch(query) {
  const scored = DB.chunks.map(c => [scoreChunk(c, query), c]).filter(([s]) => s > 0).sort((a, b) => b[0] - a[0]).slice(0, 80);
  const results = scored.map(([, c]) => c);
  $("resultCount").textContent = query.trim() ? `${results.length} 个结果` : "等待输入";
  $("results").innerHTML = results.length ? results.map(c => `
    <article class="result">
      <div class="result-title">
        <span>${escapeHtml(c.title || c.file)}</span>
        <button class="chip result-quiz" data-topic="${escapeHtml((c.terms && c.terms[0]) || query)}">出题</button>
      </div>
      <div class="path">${escapeHtml((c.heading_path || []).join(" / "))}</div>
      <p class="snippet">${highlight(c.summary, query)}</p>
      <span class="meta">${escapeHtml(c.file)} · ${c.slide_index ? `第 ${c.slide_index} 页` : `cell ${c.cell_index}`} · ${escapeHtml(c.kind || c.type)}</span>
      <div class="chips">${(c.terms || []).slice(0, 8).map(t => `<button class="chip" data-term="${escapeHtml(t)}">${escapeHtml(t)}</button>`).join("")}</div>
    </article>
  `).join("") : `<div class="panel"><p class="meta">输入课件中的词条，例如“感受野”“BatchNorm”“LSTM”“CrossEntropyLoss”。</p></div>`;
  $("results").querySelectorAll(".chip[data-term]").forEach(chip => chip.addEventListener("click", () => {
    $("searchInput").value = chip.dataset.term;
    runSearch(chip.dataset.term);
  }));
  $("results").querySelectorAll(".result-quiz").forEach(btn => btn.addEventListener("click", () => generateStandardQuiz(btn.dataset.topic)));

  const fuzzy = DB.terms[query] || Object.values(DB.terms).find(t => query && (norm(t.term).includes(norm(query)) || norm(query).includes(norm(t.term))));
  $("explainBox").innerHTML = fuzzy ? `
    <h2>${escapeHtml(fuzzy.term)}</h2>
    <p class="muted">出现 ${fuzzy.total_count} 次，分布在 ${fuzzy.files.length} 个课件文件中。</p>
    <div class="chips">${(fuzzy.related || []).map(r => `<button class="chip" data-term="${escapeHtml(r)}">${escapeHtml(r)}</button>`).join("")}</div>
  ` : `<h2>搜索说明</h2><p class="muted">搜索会定位到课件原文片段、相关词条和来源位置。适合先查概念，再回到教材目录和主题书库复习。</p>`;
  $("explainBox").querySelectorAll(".chip[data-term]").forEach(chip => chip.addEventListener("click", () => {
    $("searchInput").value = chip.dataset.term;
    runSearch(chip.dataset.term);
  }));
}

function renderOutline() {
  const courses = DB.outline.courses;
  $("courseIndex").innerHTML = courses.map((c, i) => `<a href="#course-${i}">${escapeHtml(c.title)}</a>`).join("");
  $("outlineBody").innerHTML = courses.map((course, i) => `
    <section class="panel course-card" id="course-${i}">
      <div class="course-head">
        <div>
          <h2>${escapeHtml(course.title)}</h2>
          <p class="muted">${escapeHtml(course.file)} · ${course.chapter_count} 个章节块 · ${course.chunk_count} 个片段</p>
        </div>
      </div>
      ${course.chapters.map(ch => `
        <div class="chapter">
          <h3>${escapeHtml(ch.title)}</h3>
          <table>
            <thead><tr><th>知识点</th><th>通俗解释</th><th>考试怎么考</th><th>易错点</th><th>必记句子/公式</th><th>来源</th></tr></thead>
            <tbody>${ch.rows.map(r => `
              <tr>
                <td><strong>${escapeHtml(r.term)}</strong></td>
                <td>${escapeHtml(r.plain)}</td>
                <td>${escapeHtml(r.exam)}</td>
                <td>${escapeHtml(r.pitfall)}</td>
                <td>${escapeHtml(r.memory)}</td>
                <td class="muted">${escapeHtml(r.source)}</td>
              </tr>
            `).join("")}</tbody>
          </table>
        </div>
      `).join("")}
    </section>
  `).join("");
}

function pickQuestions(bank, type, topic, count) {
  const q = norm(topic || "");
  const all = bank.filter(x => x.type === type);
  const related = q ? all.filter(x => norm([x.topic, x.stem, x.source].join(" ")).includes(q)) : [];
  const base = related.length >= Math.min(count, 3) ? related : all;
  return shuffle(base).slice(0, count);
}

function generateQuiz(bank, mode, topic = "") {
  const spec = mode === "beginner"
    ? [["mcq", 10], ["short", 4], ["material", 1]]
    : [["mcq", 20], ["short", 6], ["material", 2]];
  const labels = {mcq: "单选题", short: "简答题", material: "资料题"};
  const chosen = spec.flatMap(([type, count]) => pickQuestions(bank, type, topic, count));
  const title = mode === "beginner" ? "基础练习随机题" : "标准综合随机题";
  const targetTitle = mode === "beginner" ? $("beginnerTitle") : $("quizTitle");
  const targetMeta = mode === "beginner" ? $("beginnerMeta") : $("quizMeta");
  const targetBody = mode === "beginner" ? $("beginnerBody") : $("quizBody");
  lastQuizState[mode] = {chosen, spec, topic, title};
  targetTitle.textContent = topic ? `${title}：${topic}` : title;
  targetMeta.textContent = spec.map(([type, count]) => `${count} 道${labels[type]}`).join("、") + "。再次点击会重新随机生成。";
  targetBody.innerHTML = spec.map(([type]) => {
    const items = chosen.filter(q => q.type === type);
    return `<section class="quiz-section"><h3>${labels[type]}</h3>${items.map((q, i) => renderQuestion(q, i + 1)).join("")}</section>`;
  }).join("") + `<section class="panel answers"><h3>答案与解析</h3>${chosen.map((q, i) => renderAnswer(q, i + 1)).join("")}</section>`;
}

function generateStandardQuiz(topic = "", activate = true) {
  generateQuiz(DB.quizBank, "standard", topic);
  if (activate) {
    showView("quiz");
    window.scrollTo({top: 0, behavior: "smooth"});
  }
}

function generateBeginnerQuiz(topic = "", activate = true) {
  generateQuiz(DB.beginnerQuizBank, "beginner", topic);
  if (activate) {
    showView("beginner");
    window.scrollTo({top: 0, behavior: "smooth"});
  }
}

function renderQuestion(q, idx) {
  if (q.type === "mcq") {
    return `<div class="q"><div class="q-stem">${idx}. ${escapeHtml(q.stem)} <span class="badge">${escapeHtml(q.difficulty)}</span></div><ul class="options">${q.options.map(o => `<li>${escapeHtml(o)}</li>`).join("")}</ul></div>`;
  }
  if (q.type === "material") {
    const subs = q.subquestions || q.questions || [];
    return `<div class="q"><div class="q-stem">${idx}. ${escapeHtml(q.stem)} <span class="badge">${escapeHtml(q.difficulty)}</span></div><p class="snippet"><strong>材料：</strong>${escapeHtml(q.material || "")}</p><ol>${subs.map(x => `<li>${escapeHtml(x)}</li>`).join("")}</ol><p class="muted">答题区：</p><p style="height:78px;border-bottom:1px solid var(--line)"></p></div>`;
  }
  return `<div class="q"><div class="q-stem">${idx}. ${escapeHtml(q.stem)} <span class="badge">${escapeHtml(q.difficulty)}</span></div><p class="muted">答题区：</p><p style="height:54px;border-bottom:1px solid var(--line)"></p></div>`;
}

function renderAnswer(q, idx) {
  return `<div class="answer-item"><strong>${idx}. [${escapeHtml(q.id)}] ${escapeHtml(q.answer)}</strong><p>${escapeHtml(q.explanation)}</p><span class="muted">来源：${escapeHtml(q.source)} · ${escapeHtml(q.topic)}</span></div>`;
}

function inlineMarkdown(text) {
  return escapeHtml(text)
    .replace(/`([^`]+)`/g, "<code>$1</code>")
    .replace(/\*\*([^*]+)\*\*/g, "<strong>$1</strong>")
    .replace(/\[([^\]]+)\]\(([^)]+)\)/g, '<a href="$2" target="_blank">$1</a>');
}

function markdownToHtml(md) {
  const lines = (md || "").split(/\r?\n/);
  const html = [];
  let listOpen = false;
  let inCode = false;
  let inMath = false;
  let codeLines = [];
  let mathLines = [];
  const closeList = () => {
    if (listOpen) {
      html.push("</ul>");
      listOpen = false;
    }
  };
  lines.forEach(line => {
    const trimmed = line.trim();
    const singleLineMath = trimmed.match(/^\$\$(.+)\$\$$/);
    if (!inCode && singleLineMath) {
      closeList();
      html.push(`<div class="math-block">$$${escapeHtml(singleLineMath[1].trim())}$$</div>`);
      return;
    }
    if (!inCode && trimmed === "$$") {
      if (inMath) {
        html.push(`<div class="math-block">$$${escapeHtml(mathLines.join("\n"))}$$</div>`);
        mathLines = [];
        inMath = false;
      } else {
        closeList();
        inMath = true;
      }
      return;
    }
    if (inMath) {
      mathLines.push(line);
      return;
    }
    if (line.trim().startsWith("```")) {
      if (inCode) {
        html.push(`<pre><code>${escapeHtml(codeLines.join("\n"))}</code></pre>`);
        codeLines = [];
        inCode = false;
      } else {
        closeList();
        inCode = true;
      }
      return;
    }
    if (inCode) {
      codeLines.push(line);
      return;
    }
    if (!line.trim()) {
      closeList();
      return;
    }
    const heading = line.match(/^(#{1,4})\s+(.+)$/);
    if (heading) {
      closeList();
      const level = Math.min(heading[1].length + 1, 5);
      const id = safeFileName(heading[2]).toLowerCase();
      html.push(`<h${level} id="${escapeHtml(id)}">${inlineMarkdown(heading[2])}</h${level}>`);
      return;
    }
    if (/^\s*[-*]\s+/.test(line)) {
      if (!listOpen) {
        html.push("<ul>");
        listOpen = true;
      }
      html.push(`<li>${inlineMarkdown(line.replace(/^\s*[-*]\s+/, ""))}</li>`);
      return;
    }
    if (/^\s*\d+\.\s+/.test(line)) {
      closeList();
      html.push(`<p>${inlineMarkdown(line)}</p>`);
      return;
    }
    if (/^\|.+\|$/.test(line)) {
      closeList();
      html.push(`<p class="table-line">${inlineMarkdown(line)}</p>`);
      return;
    }
    closeList();
    html.push(`<p>${inlineMarkdown(line)}</p>`);
  });
  closeList();
  if (inMath) html.push(`<div class="math-block">$$${escapeHtml(mathLines.join("\n"))}$$</div>`);
  if (inCode) html.push(`<pre><code>${escapeHtml(codeLines.join("\n"))}</code></pre>`);
  return html.join("");
}

function renderTopicGuideNav(items) {
  return items.map((g, i) => `
    <button class="book-nav-item ${i === 0 ? "active" : ""}" data-guide="${escapeHtml(g.id)}">
      <strong>${escapeHtml(g.label)}</strong>
      <span>${g.exists ? `${g.wordCount} 字` : "待生成"}</span>
    </button>
  `).join("");
}

function renderTopicGuideContent(guide) {
  if (!guide || !guide.exists) {
    return `<div class="panel"><h2>讲义还在生成</h2><p class="muted">这个主题的 Markdown 文件尚未写入，完成后会自动显示在这里。</p></div>`;
  }
  const toc = (guide.headings || []).filter(h => h.level <= 3).slice(0, 28);
  return `
    <article class="book-reader">
      <div class="book-toolbar">
        <div>
          <h2>${escapeHtml(guide.title)}</h2>
          <p class="muted">${guide.wordCount} 字 · Markdown 讲义 · 可打印</p>
        </div>
        <div class="quiz-actions">
          <a class="btn secondary" href="${escapeHtml(guide.href)}" download>下载 Markdown</a>
          <button class="btn secondary" data-print-guide>打印本讲义</button>
        </div>
      </div>
      ${(guide.tocImages || []).length ? `
        <div class="toc-image-strip">
          ${guide.tocImages.map(img => `
            <figure>
              <a href="${escapeHtml(img)}" target="_blank"><img src="${escapeHtml(img)}" alt="${escapeHtml(guide.title)} 书本式目录图"></a>
              <figcaption><a href="${escapeHtml(img)}" download>下载目录图</a></figcaption>
            </figure>
          `).join("")}
        </div>
      ` : ""}
      <div class="book-layout">
        <aside class="book-toc">
          <h3>本讲义目录</h3>
          ${toc.map(h => `<a class="toc-level-${h.level}" href="#${escapeHtml(safeFileName(h.title).toLowerCase())}">${escapeHtml(h.title)}</a>`).join("")}
        </aside>
        <div class="markdown-body">${markdownToHtml(guide.content)}</div>
      </div>
    </article>
  `;
}

function renderReview() {
  const library = DB.topicGuides || {items: []};
  const items = library.items || [];
  $("reviewBody").innerHTML = `
    <div class="book-shell">
      <aside class="panel book-sidebar">
        <h3>主题书库</h3>
        <p class="muted">按学习主题重新组织，不再按课件文件硬拆。</p>
        <div id="topicGuideNav">${renderTopicGuideNav(items)}</div>
      </aside>
      <div id="topicGuideReader">${renderTopicGuideContent(items[0])}</div>
    </div>
  `;
  queueMathTypeset($("topicGuideReader"));
  $("topicGuideNav").querySelectorAll("button[data-guide]").forEach(btn => btn.addEventListener("click", () => {
    $("topicGuideNav").querySelectorAll("button").forEach(b => b.classList.toggle("active", b === btn));
    const guide = items.find(g => g.id === btn.dataset.guide);
    $("topicGuideReader").innerHTML = renderTopicGuideContent(guide);
    const printBtn = $("topicGuideReader").querySelector("[data-print-guide]");
    if (printBtn) printBtn.addEventListener("click", () => window.print());
    queueMathTypeset($("topicGuideReader"));
  }));
  const printBtn = $("topicGuideReader").querySelector("[data-print-guide]");
  if (printBtn) printBtn.addEventListener("click", () => window.print());
}

function init() {
  document.querySelectorAll("nav button").forEach(btn => btn.addEventListener("click", () => showView(btn.dataset.view)));
  $("searchInput").addEventListener("input", e => runSearch(e.target.value));
  $("makeStandardQuiz").addEventListener("click", () => generateStandardQuiz($("searchInput").value.trim()));
  $("makeBeginnerQuiz").addEventListener("click", () => generateBeginnerQuiz($("searchInput").value.trim()));
  $("quizAll").addEventListener("click", () => generateStandardQuiz(""));
  $("beginnerAll").addEventListener("click", () => generateBeginnerQuiz(""));
  $("quizExportQuestions").addEventListener("click", () => exportQuizMarkdown("standard", "questions"));
  $("quizExportAnswers").addEventListener("click", () => exportQuizMarkdown("standard", "answers"));
  $("beginnerExportQuestions").addEventListener("click", () => exportQuizMarkdown("beginner", "questions"));
  $("beginnerExportAnswers").addEventListener("click", () => exportQuizMarkdown("beginner", "answers"));
  renderStats();
  renderTermCloud();
  renderOutline();
  renderReview();
  runSearch("");
  generateBeginnerQuiz("", false);
}

init();
"""


HTML = """<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>深度学习期末复习资料库</title>
  <link rel="stylesheet" href="assets/course_site.css">
  <script>
    window.MathJax = {
      tex: {
        inlineMath: [['$', '$']],
        displayMath: [['$$', '$$']],
        processEscapes: true
      },
      svg: { fontCache: 'global' },
      startup: { typeset: false },
      options: { skipHtmlTags: ['script', 'noscript', 'style', 'textarea', 'pre', 'code'] }
    };
  </script>
  <script defer src="https://cdn.jsdelivr.net/npm/mathjax@3/es5/tex-svg.js"></script>
</head>
<body>
  <header>
    <div class="topbar">
      <div class="brand">
        <h1>深度学习期末复习资料库</h1>
        <p>检索 · 主题书库 · 基础练习 · 标准组卷</p>
      </div>
      <nav aria-label="主导航">
        <button class="active" data-view="search">搜索</button>
        <button data-view="outline">教材目录</button>
        <button data-view="beginner">基础练习</button>
        <button data-view="quiz">标准组卷</button>
        <button data-view="review">主题书库</button>
      </nav>
    </div>
  </header>
  <main>
    <section id="view-search" class="view active">
      <div class="hero-grid">
        <div class="panel">
          <h2>从一个词条开始复习</h2>
          <p class="muted">搜索会回到课件原文位置，也能继续跳到基础题、标准题、目录和主题书库。</p>
          <div class="toolbar">
            <input id="searchInput" type="search" placeholder="输入词条：如 感受野、BatchNorm、LSTM、CrossEntropyLoss">
            <button id="makeBeginnerQuiz" class="btn secondary">基础练习</button>
            <button id="makeStandardQuiz" class="btn">标准组卷</button>
          </div>
        </div>
        <div id="stats" class="stats"></div>
      </div>
      <div class="layout">
        <div>
          <div class="panel" style="margin-bottom:14px">
            <div class="result-title"><h2 style="margin:0">搜索结果</h2><span id="resultCount" class="meta"></span></div>
          </div>
          <div id="results" class="panel"></div>
        </div>
        <aside>
          <div id="explainBox" class="panel"></div>
          <div class="panel" style="margin-top:14px">
            <h3>高频词条</h3>
            <div id="termCloud" class="term-list"></div>
          </div>
        </aside>
      </div>
    </section>

    <section id="view-outline" class="view">
      <div class="course-layout">
        <aside class="panel side-index">
          <h3>课件目录</h3>
          <div id="courseIndex"></div>
        </aside>
        <div id="outlineBody"></div>
      </div>
    </section>

    <section id="view-beginner" class="view">
      <div class="panel">
        <div class="quiz-head">
          <div>
            <h2 id="beginnerTitle">基础练习</h2>
            <p id="beginnerMeta" class="muted">适合先熟悉课程脉络。</p>
          </div>
          <div class="quiz-actions">
            <button id="beginnerExportQuestions" class="btn secondary">导出题目.md</button>
            <button id="beginnerExportAnswers" class="btn secondary">导出答案.md</button>
            <button id="beginnerAll" class="btn">重新生成基础题</button>
          </div>
        </div>
        <div id="beginnerBody"></div>
      </div>
    </section>

    <section id="view-quiz" class="view">
      <div class="panel">
        <div class="quiz-head">
          <div>
            <h2 id="quizTitle">标准综合随机题</h2>
            <p id="quizMeta" class="muted">保留原标准题库，适合考前检测。</p>
          </div>
          <div class="quiz-actions">
            <button id="quizExportQuestions" class="btn secondary">导出题目.md</button>
            <button id="quizExportAnswers" class="btn secondary">导出答案.md</button>
            <button id="quizAll" class="btn">生成标准题</button>
          </div>
        </div>
        <div id="quizBody"></div>
      </div>
    </section>

    <section id="view-review" class="view">
      <div class="panel">
        <h2>主题学习讲义书库</h2>
        <p class="muted">这里按 MLP、CNN、RNN、Transformer、图学习、NLP、LLM 重新组织，像电子书一样阅读，也可以下载 Markdown 打印。</p>
        <div id="reviewBody"></div>
      </div>
    </section>
  </main>
  <script src="assets/course_data.js"></script>
  <script src="assets/course_site.js"></script>
</body>
</html>
"""


def write_json(path: Path, obj):
    path.write_text(json.dumps(obj, ensure_ascii=False, indent=2), encoding="utf-8")


TOPIC_GUIDE_SPECS = [
    {
        "id": "mlp",
        "title": "MLP学习资料-通俗版",
        "label": "MLP",
        "href": "topic_guides/MLP学习资料-通俗版.md",
        "tocImages": ["guide_toc_images/mlp-toc.png"],
    },
    {
        "id": "cnn",
        "title": "CNN学习资料-通俗版",
        "label": "CNN",
        "href": "topic_guides/CNN学习资料-通俗版.md",
        "tocImages": ["guide_toc_images/cnn-toc.png"],
    },
    {
        "id": "rnn",
        "title": "RNN学习资料-通俗版",
        "label": "RNN",
        "href": "topic_guides/RNN学习资料-通俗版.md",
        "tocImages": ["guide_toc_images/rnn-toc.png"],
    },
    {
        "id": "transformer",
        "title": "Transformer学习资料-通俗版",
        "label": "Transformer",
        "href": "topic_guides/Transformer学习资料-通俗版.md",
        "tocImages": ["guide_toc_images/transformer-toc.png"],
    },
    {
        "id": "graph",
        "title": "图学习资料-通俗版",
        "label": "图学习",
        "href": "topic_guides/图学习资料-通俗版.md",
        "tocImages": ["guide_toc_images/graph-toc.png"],
    },
    {
        "id": "nlp",
        "title": "NLP学习资料-通俗版",
        "label": "NLP 扩展",
        "href": "topic_guides/NLP学习资料-通俗版.md",
        "tocImages": ["guide_toc_images/nlp-toc.png"],
    },
    {
        "id": "llm",
        "title": "LLM学习资料-通俗版",
        "label": "LLM 扩展",
        "href": "topic_guides/LLM学习资料-通俗版.md",
        "tocImages": ["guide_toc_images/llm-toc.png"],
    },
    {
        "id": "practice_questions",
        "title": "课堂练习册-题目",
        "label": "练习题册",
        "href": "topic_guides/课堂练习册-题目.md",
        "tocImages": [],
    },
    {
        "id": "practice_answers",
        "title": "课堂练习册-答案解析",
        "label": "答案解析册",
        "href": "topic_guides/课堂练习册-答案解析.md",
        "tocImages": [],
    },
    {
        "id": "final_review",
        "title": "深度学习期末复习-总复习资料",
        "label": "总复习资料",
        "href": "topic_guides/深度学习期末复习-总复习资料.md",
        "tocImages": [],
    },
]


def load_topic_guides():
    items = []
    for spec in TOPIC_GUIDE_SPECS:
        path = SITE / spec["href"]
        content = path.read_text(encoding="utf-8") if path.exists() else ""
        headings = re.findall(r"^(#{1,3})\s+(.+)$", content, flags=re.M)
        items.append({
            **spec,
            "content": content,
            "wordCount": len(re.sub(r"\s+", "", content)),
            "headings": [{"level": len(mark), "title": title.strip()} for mark, title in headings[:80]],
            "exists": path.exists(),
            "tocImages": [img for img in spec["tocImages"] if (SITE / img).exists()],
        })
    return {"title": "主题学习讲义书库", "items": items}


def main():
    DATABASE.mkdir(parents=True, exist_ok=True)
    ASSETS.mkdir(parents=True, exist_ok=True)
    IMAGE_OUT.mkdir(parents=True, exist_ok=True)
    GUIDE_OUT.mkdir(parents=True, exist_ok=True)
    TOPIC_GUIDE_OUT.mkdir(parents=True, exist_ok=True)
    TOC_IMAGE_OUT.mkdir(parents=True, exist_ok=True)
    if MINDMAP_OUT.exists():
        shutil.rmtree(MINDMAP_OUT)
    mindmap_db = DATABASE / "mindmaps.json"
    if mindmap_db.exists():
        mindmap_db.unlink()

    standard_quiz, beginner_quiz = make_exam_banks()
    manifest, chunks, source_cache = extract_sources()
    terms = build_term_index(chunks)
    outline = build_course_outline(manifest, chunks, concept_lookup_from_quiz(standard_quiz))
    standard_quiz, beginner_quiz = extend_quiz_from_outline(standard_quiz, beginner_quiz, outline)
    validate_quiz_bank(standard_quiz, "标准题库")
    validate_quiz_bank(beginner_quiz, "基础题库")
    review = load_review_or_fallback(outline)
    guides = build_study_guides(outline)
    write_final_review_markdown(outline, manifest)
    topic_guides = load_topic_guides()

    for img in (ROOT / "image").glob("*.png"):
        if img.name.startswith("截屏"):
            continue
        shutil.copy2(img, IMAGE_OUT / img.name)

    write_json(DATABASE / "source_manifest.json", manifest)
    write_json(DATABASE / "source_cache.json", source_cache)
    write_json(DATABASE / "course_chunks.json", chunks)
    write_json(DATABASE / "course_terms.json", terms)
    write_json(DATABASE / "course_outline.json", outline)
    write_json(DATABASE / "review_material.json", review)
    write_json(DATABASE / "quiz_bank.json", standard_quiz)
    write_json(DATABASE / "beginner_quiz_bank.json", beginner_quiz)
    write_json(DATABASE / "study_guides.json", guides)
    write_json(DATABASE / "topic_guides.json", topic_guides)

    bundle = {
        "manifest": manifest,
        "chunks": chunks,
        "terms": terms,
        "outline": outline,
        "review": review,
        "guides": guides,
        "topicGuides": topic_guides,
        "quizBank": standard_quiz,
        "beginnerQuizBank": beginner_quiz,
    }
    (ASSETS / "course_data.js").write_text("window.COURSE_DB = " + json.dumps(bundle, ensure_ascii=False) + ";\n", encoding="utf-8")
    (ASSETS / "course_site.css").write_text(CSS.strip() + "\n", encoding="utf-8")
    (ASSETS / "course_site.js").write_text(JS.strip() + "\n", encoding="utf-8")
    (SITE / "index.html").write_text(HTML, encoding="utf-8")

    print(
        f"sources={len(manifest)} chunks={len(chunks)} terms={len(terms)} "
        f"standard_quiz={len(standard_quiz)} beginner_quiz={len(beginner_quiz)} "
        f"guides={len(guides['items'])} "
        f"topic_guides={sum(1 for g in topic_guides['items'] if g['exists'])}"
    )
    print(SITE / "index.html")


if __name__ == "__main__":
    main()
